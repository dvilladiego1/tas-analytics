#!/usr/bin/env python3
"""
Kia Coapa / Continental Coapa — Monthly Business Review — Febrero 2026
Generates a 10-slide widescreen PPTX using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# ── Colours ──────────────────────────────────────────────────────────────
KAVAK_BLUE  = RGBColor(0x00, 0x4E, 0x98)
KAVAK_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
RED         = RGBColor(0xEF, 0x44, 0x44)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
CARD_BG     = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE_LIGHT= RGBColor(0xED, 0xE9, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_blue_bar(slide):
    """Top accent bar"""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=KAVAK_BLUE)

def add_slide_title(slide, title_text, subtitle_text=None):
    """Standard slide title below blue bar"""
    add_blue_bar(slide)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.5),
                 title_text, font_size=24, bold=True, color=KAVAK_DARK)
    if subtitle_text:
        add_text_box(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
                     subtitle_text, font_size=12, color=MED_GRAY)

def add_kpi_card(slide, left, top, width, height, label, value, delta, delta_color, accent_color=KAVAK_BLUE):
    """Rounded rectangle KPI card with accent top"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, line_color=LIGHT_GRAY)
    # accent top bar inside card
    add_shape(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.06), fill_color=accent_color)
    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), Inches(0.3),
                 label, font_size=10, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), Inches(0.5),
                 value, font_size=28, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
    # Delta
    add_text_box(slide, left + Inches(0.15), top + Inches(1.0), width - Inches(0.3), Inches(0.3),
                 delta, font_size=11, bold=True, color=delta_color, alignment=PP_ALIGN.CENTER)

def set_cell_text(cell, text, font_size=9, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def style_table_header(table, cols, fill_color=KAVAK_BLUE):
    for i in range(cols):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        set_cell_text(cell, cell.text, font_size=9, bold=True, color=WHITE)

def style_table_rows(table, rows, cols):
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


def add_insight_box(slide, left, top, width, height, text, accent_color=KAVAK_BLUE, bg_color=BLUE_LIGHT):
    """Rounded rectangle insight box with left accent"""
    box = add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    # left accent
    add_shape(slide, left, top + Inches(0.1), Inches(0.06), height - Inches(0.2), fill_color=accent_color)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), width - Inches(0.35), height - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = KAVAK_DARK
        p.font.name = "Calibri"
        p.space_after = Pt(2)
        if line.startswith(">>"):
            p.text = line[2:].strip()
            p.font.bold = True
            p.font.size = Pt(10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
# Full dark bg
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=KAVAK_DARK)
# Blue accent bar top
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.12), fill_color=KAVAK_BLUE)
# Blue accent line in middle
add_shape(slide1, Inches(1.2), Inches(3.0), Inches(1.5), Inches(0.06), fill_color=KAVAK_BLUE)
# Title text
add_text_box(slide1, Inches(1.2), Inches(1.8), Inches(10), Inches(0.6),
             "Continental — Kia Coapa", font_size=14, bold=True, color=KAVAK_BLUE)
add_text_box(slide1, Inches(1.2), Inches(2.3), Inches(10), Inches(0.7),
             "Monthly Business Review", font_size=36, bold=True, color=WHITE)
add_text_box(slide1, Inches(1.2), Inches(3.3), Inches(10), Inches(0.5),
             "Febrero 2026", font_size=22, bold=False, color=WHITE)
add_text_box(slide1, Inches(1.2), Inches(4.2), Inches(10), Inches(0.4),
             "Kavak Marketplace B2B — Programa TAS", font_size=14, color=MED_GRAY)
# Bottom right date
add_text_box(slide1, Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.4),
             "1 de Marzo 2026", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_slide_title(slide2, "Executive Summary — Febrero 2026", "Kia Coapa | Continental Coapa")

# 4 KPI Cards
card_w = Inches(2.8)
card_h = Inches(1.4)
card_y = Inches(1.1)
gap = Inches(0.3)
start_x = Inches(0.6)

add_kpi_card(slide2, start_x, card_y, card_w, card_h,
             "COTIZACIONES", "908", "▼ -8.2% vs Ene (989)", RED, KAVAK_BLUE)
add_kpi_card(slide2, start_x + card_w + gap, card_y, card_w, card_h,
             "INSPECCIONES", "220", "▲ +22.9% vs Ene (179)", GREEN, KAVAK_BLUE)
add_kpi_card(slide2, start_x + 2*(card_w + gap), card_y, card_w, card_h,
             "COMPRAS NETAS", "69", "▲ +19.0% vs Ene (58)", GREEN, GREEN)
add_kpi_card(slide2, start_x + 3*(card_w + gap), card_y, card_w, card_h,
             "Q→P CONVERSION", "7.6%", "▲ +1.7pp vs Ene (5.9%)", GREEN, PURPLE)

# Highlights box
hl_text = (">> HIGHLIGHTS\n"
           "• 69 compras netas — record historico de la celula (+19% MoM)\n"
           "• 8vo mes consecutivo de crecimiento en compras\n"
           "• Q→M mejoro a 24.2% (vs 18.1% en Ene, +6.1pp)\n"
           "• Aliados consolidados: 51 compras (73.9% del total)\n"
           "• Semana 9 fue la mejor: 19 compras con Q→P 11.3%")
add_insight_box(slide2, Inches(0.6), Inches(2.75), Inches(5.8), Inches(2.05),
                hl_text, accent_color=GREEN, bg_color=GREEN_LIGHT)

# Lowlights box
ll_text = (">> LOWLIGHTS\n"
           "• A→P cayo a 36.9% (vs 40% en Ene). Cuesta mas cerrar aprobados\n"
           "• M→A bajo a 85% — 33 inspecciones rechazadas en el mes\n"
           "• Quotes bajaron -8.2% (908 vs 989). Primera caida en meses\n"
           "• Q→P de Aliado (6.9%) vs TAS (10.8%) — gap de eficiencia")
add_insight_box(slide2, Inches(6.9), Inches(2.75), Inches(5.8), Inches(2.05),
                ll_text, accent_color=RED, bg_color=RED_LIGHT)

# Bottom insight
add_insight_box(slide2, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.95),
                ">> Resumen: Kia Coapa alcanzo su record de 69 compras en Feb, consolidando 8 meses de crecimiento ininterrumpido desde su lanzamiento en ~May 2025. El principal reto es mejorar A→P (36.9%) donde 118 vehiculos aprobados no se compraron.",
                accent_color=KAVAK_BLUE, bg_color=BLUE_LIGHT)

# Footer
add_text_box(slide2, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ENE vs FEB COMPARISON
# ═══════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_slide_title(slide3, "Enero vs Febrero 2026 — Comparativo", "Cambios mes a mes en metricas clave del funnel")

# Table: Metrica | Ene 26 | Feb 26 | Cambio | Var%
rows_data = [
    ["Quotes",          "989",  "908",  "-81",   "-8.2%",  RED],
    ["Made",            "179",  "220",  "+41",   "+22.9%", GREEN],
    ["Approved",        "145",  "187",  "+42",   "+29.0%", GREEN],
    ["Purchased",       "59",   "69",   "+10",   "+16.9%", GREEN],
    ["Devoluciones",    "1",    "0",    "-1",    "-100%",  GREEN],
    ["Compras Netas",   "58",   "69",   "+11",   "+19.0%", GREEN],
    ["Q→M %",          "18.1%","24.2%","+6.1pp", "—",     GREEN],
    ["M→A %",          "81.0%","85.0%","+4.0pp", "—",     GREEN],
    ["A→P %",          "40.0%","36.9%","-3.1pp", "—",     RED],
    ["Q→P %",          "5.9%", "7.6%", "+1.7pp", "—",     GREEN],
]

tbl_rows = len(rows_data) + 1
tbl_cols = 6
table_shape = slide3.shapes.add_table(tbl_rows, tbl_cols, Inches(0.6), Inches(1.1), Inches(9.5), Inches(3.8))
table = table_shape.table

headers = ["Metrica", "Ene 26", "Feb 26", "Cambio", "Var %", "Tendencia"]
for i, h in enumerate(headers):
    set_cell_text(table.cell(0, i), h, font_size=10, bold=True, color=WHITE)
style_table_header(table, tbl_cols)

for r, row in enumerate(rows_data, start=1):
    set_cell_text(table.cell(r, 0), row[0], font_size=10, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.LEFT)
    set_cell_text(table.cell(r, 1), row[1], font_size=10)
    set_cell_text(table.cell(r, 2), row[2], font_size=10, bold=True)
    set_cell_text(table.cell(r, 3), row[3], font_size=10, bold=True, color=row[5])
    set_cell_text(table.cell(r, 4), row[4], font_size=10, color=row[5])
    arrow = "▲" if row[5] == GREEN else "▼"
    set_cell_text(table.cell(r, 5), arrow, font_size=14, bold=True, color=row[5])

style_table_rows(table, tbl_rows, tbl_cols)

# Column widths
col_widths = [Inches(2.0), Inches(1.3), Inches(1.3), Inches(1.5), Inches(1.4), Inches(1.2)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Insight box
add_insight_box(slide3, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.2),
                ">> Lectura clave\nFeb consolida el crecimiento: 8 meses seguidos de aumento en compras. El Q→M mejoro +6pp, lo que mas que compensa la caida en quotes. El reto principal es A→P que baja 3.1pp — hay 118 aprobados que no se compraron. La mejora en volumen de inspecciones (+22.9%) fue el gran motor del mes.",
                accent_color=KAVAK_BLUE, bg_color=BLUE_LIGHT)

add_text_box(slide3, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HISTORICAL EVOLUTION
# ═══════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_slide_title(slide4, "Evolucion Historica — Compras Netas", "Mayo 2025 – Febrero 2026 | 8 meses consecutivos de crecimiento")

months = ["May 25", "Jun 25", "Jul 25", "Ago 25", "Sep 25", "Oct 25", "Nov 25", "Dic 25", "Ene 26", "Feb 26"]
compras = [3, 0, 11, 20, 21, 32, 42, 45, 58, 69]
qp_pct  = [25.0, 0, 14.5, 9.8, 10.6, 6.5, 7.9, 9.6, 5.9, 7.6]

chart_data = CategoryChartData()
chart_data.categories = months
chart_data.add_series("Compras Netas", compras)

chart_frame = slide4.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(1.1), Inches(8.5), Inches(4.8),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = False

# Style the chart
plot = chart.plots[0]
plot.gap_width = 80
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = KAVAK_BLUE

# Data labels
series.has_data_labels = True
data_labels = series.data_labels
data_labels.font.size = Pt(11)
data_labels.font.bold = True
data_labels.font.color.rgb = KAVAK_DARK
data_labels.number_format = '0'
data_labels.show_value = True
data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

# Value axis
value_axis = chart.value_axis
value_axis.visible = True
value_axis.has_title = False
value_axis.major_gridlines.format.line.color.rgb = LIGHT_GRAY
value_axis.format.line.color.rgb = LIGHT_GRAY
value_axis.tick_labels.font.size = Pt(9)
value_axis.tick_labels.font.color.rgb = MED_GRAY

# Category axis
cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(9)
cat_axis.tick_labels.font.color.rgb = DARK_GRAY
cat_axis.format.line.color.rgb = LIGHT_GRAY

# Q→P table on the right
qp_tbl_shape = slide4.shapes.add_table(11, 2, Inches(9.5), Inches(1.1), Inches(3.2), Inches(4.0))
qp_tbl = qp_tbl_shape.table
qp_tbl.columns[0].width = Inches(1.4)
qp_tbl.columns[1].width = Inches(1.4)

set_cell_text(qp_tbl.cell(0, 0), "Mes", font_size=9, bold=True, color=WHITE)
set_cell_text(qp_tbl.cell(0, 1), "Q→P %", font_size=9, bold=True, color=WHITE)
style_table_header(qp_tbl, 2)

for i, (m, q) in enumerate(zip(months, qp_pct)):
    set_cell_text(qp_tbl.cell(i+1, 0), m, font_size=9)
    set_cell_text(qp_tbl.cell(i+1, 1), f"{q}%", font_size=9, bold=(i==len(months)-1),
                  color=KAVAK_BLUE if i==len(months)-1 else DARK_GRAY)
style_table_rows(qp_tbl, 11, 2)

# Insight
add_insight_box(slide4, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7),
                ">> Trayectoria: Celula lanzo ~May 2025. Desde Jul 25, 8 meses consecutivos de crecimiento en compras (11→20→21→32→42→45→58→69). Ritmo de crecimiento MoM promedio ~+27% en ultimos 4 meses.",
                accent_color=GREEN, bg_color=GREEN_LIGHT)

add_text_box(slide4, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WEEKLY & DAILY
# ═══════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_slide_title(slide5, "Estacionalidad — Semanal y Diario", "Febrero 2026 | Distribucion de compras en el mes")

# Weekly table
wk_headers = ["Metrica", "S5", "S6", "S7", "S8", "S9"]
wk_data = [
    ["Quotes",       "259", "253", "242", "219", "168"],
    ["Made",         "44",  "45",  "57",  "66",  "52"],
    ["Purchased",    "13",  "14",  "19",  "17",  "19"],
    ["Compras Netas","12",  "14",  "19",  "17",  "19"],
    ["Q→M %",       "17.0%","17.8%","23.6%","30.1%","31.0%"],
    ["Q→P %",       "4.6%","5.5%","7.9%","7.8%","11.3%"],
]

wk_rows = len(wk_data) + 1
wk_cols = len(wk_headers)
wk_shape = slide5.shapes.add_table(wk_rows, wk_cols, Inches(0.6), Inches(1.1), Inches(7.5), Inches(2.5))
wk_tbl = wk_shape.table
wk_tbl.columns[0].width = Inches(2.0)
for c in range(1, wk_cols):
    wk_tbl.columns[c].width = Inches(1.1)

for i, h in enumerate(wk_headers):
    set_cell_text(wk_tbl.cell(0, i), h, font_size=9, bold=True, color=WHITE)
style_table_header(wk_tbl, wk_cols)

for r, row in enumerate(wk_data, start=1):
    set_cell_text(wk_tbl.cell(r, 0), row[0], font_size=9, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.LEFT)
    for c in range(1, wk_cols):
        set_cell_text(wk_tbl.cell(r, c), row[c], font_size=9)
style_table_rows(wk_tbl, wk_rows, wk_cols)

# Note about weekly overlap
add_text_box(slide5, Inches(0.6), Inches(3.7), Inches(7.5), Inches(0.3),
             "Nota: Sumas semanales (81) exceden el total mensual (69) porque las semanas se solapan con meses adyacentes. El total mensual de 69 es autoritativo.",
             font_size=8, color=MED_GRAY)

# Quincena comparison cards
q1_card = add_rounded_rect(slide5, Inches(8.6), Inches(1.1), Inches(2.1), Inches(1.6), fill_color=WHITE, line_color=LIGHT_GRAY)
add_shape(slide5, Inches(8.65), Inches(1.15), Inches(2.0), Inches(0.06), fill_color=KAVAK_BLUE)
add_text_box(slide5, Inches(8.7), Inches(1.3), Inches(1.9), Inches(0.25),
             "1ra QUINCENA (1-14 Feb)", font_size=9, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(8.7), Inches(1.6), Inches(1.9), Inches(0.4),
             "33 compras", font_size=22, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(8.7), Inches(2.1), Inches(1.9), Inches(0.3),
             "2.4 / dia", font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

q2_card = add_rounded_rect(slide5, Inches(11.0), Inches(1.1), Inches(2.1), Inches(1.6), fill_color=WHITE, line_color=LIGHT_GRAY)
add_shape(slide5, Inches(11.05), Inches(1.15), Inches(2.0), Inches(0.06), fill_color=GREEN)
add_text_box(slide5, Inches(11.1), Inches(1.3), Inches(1.9), Inches(0.25),
             "2da QUINCENA (15-28 Feb)", font_size=9, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(11.1), Inches(1.6), Inches(1.9), Inches(0.4),
             "36 compras", font_size=22, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(11.1), Inches(2.1), Inches(1.9), Inches(0.3),
             "2.6 / dia", font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Key stats cards
stats_card = add_rounded_rect(slide5, Inches(8.6), Inches(2.9), Inches(4.5), Inches(1.1), fill_color=AMBER_LIGHT, line_color=AMBER)
add_text_box(slide5, Inches(8.8), Inches(3.0), Inches(4.1), Inches(0.3),
             "DATOS CLAVE DEL MES", font_size=10, bold=True, color=AMBER)
add_text_box(slide5, Inches(8.8), Inches(3.3), Inches(4.1), Inches(0.6),
             "• Mejor dia: 13 Feb (7 compras)\n• Dias con 0 compras: 6 de 28 (1, 2, 8, 15, 20, 22 Feb)\n• Operacion bien distribuida vs otras celulas",
             font_size=9, color=DARK_GRAY)

# Daily compras mini-table (compact)
daily_data_1 = [("1","0"),("2","0"),("3","3"),("4","2"),("5","1"),("6","6"),("7","2"),
                ("8","0"),("9","2"),("10","4"),("11","2"),("12","3"),("13","7"),("14","1")]
daily_data_2 = [("15","0"),("16","5"),("17","2"),("18","2"),("19","3"),("20","0"),("21","5"),
                ("22","0"),("23","1"),("24","3"),("25","4"),("26","3"),("27","3"),("28","5")]

# Row 1: days 1-14
d_shape1 = slide5.shapes.add_table(2, 14, Inches(0.6), Inches(4.3), Inches(12.1), Inches(0.7))
d_tbl1 = d_shape1.table
for c in range(14):
    d_tbl1.columns[c].width = Inches(0.864)
    set_cell_text(d_tbl1.cell(0, c), daily_data_1[c][0], font_size=8, bold=True, color=WHITE)
    val = daily_data_1[c][1]
    vc = RED if val == "0" else (GREEN if int(val) >= 5 else DARK_GRAY)
    set_cell_text(d_tbl1.cell(1, c), val, font_size=10, bold=True, color=vc)
style_table_header(d_tbl1, 14)

# Row 2: days 15-28
d_shape2 = slide5.shapes.add_table(2, 14, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.7))
d_tbl2 = d_shape2.table
for c in range(14):
    d_tbl2.columns[c].width = Inches(0.864)
    set_cell_text(d_tbl2.cell(0, c), daily_data_2[c][0], font_size=8, bold=True, color=WHITE)
    val = daily_data_2[c][1]
    vc = RED if val == "0" else (GREEN if int(val) >= 5 else DARK_GRAY)
    set_cell_text(d_tbl2.cell(1, c), val, font_size=10, bold=True, color=vc)
style_table_header(d_tbl2, 14)

add_text_box(slide5, Inches(0.6), Inches(4.1), Inches(5), Inches(0.25),
             "Compras Netas diarias — Febrero 2026 (total: 69)", font_size=9, bold=True, color=KAVAK_DARK)

# Insight
add_insight_box(slide5, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.65),
                ">> Operacion bien distribuida — solo 6 dias con 0 compras de 28 dias en el mes. Ambas quincenas casi iguales (33 vs 36). Tendencia semanal positiva: Q→M subio de 17% (S5) a 31% (S9).",
                accent_color=GREEN, bg_color=GREEN_LIGHT)

add_text_box(slide5, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TAS vs ALIADO
# ═══════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_slide_title(slide6, "TAS vs Aliado — Febrero 2026", "Desglose por canal de originacion")

# TAS Card
tas_card = add_rounded_rect(slide6, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.5), fill_color=WHITE, line_color=KAVAK_BLUE)
add_shape(slide6, Inches(0.65), Inches(1.25), Inches(5.7), Inches(0.08), fill_color=KAVAK_BLUE)
add_text_box(slide6, Inches(0.8), Inches(1.45), Inches(5.4), Inches(0.35),
             "TAS (Directo)", font_size=18, bold=True, color=KAVAK_BLUE)

tas_metrics = [("Quotes", "167", "18.4%"), ("Made", "50", "22.7%"),
               ("Approved", "46", "24.6%"), ("Purchased", "18", "26.1%")]
for i, (label, val, pct) in enumerate(tas_metrics):
    x = Inches(0.8) + Inches(1.35) * i
    add_text_box(slide6, x, Inches(1.9), Inches(1.2), Inches(0.2),
                 label, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, Inches(2.1), Inches(1.2), Inches(0.35),
                 val, font_size=22, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, Inches(2.5), Inches(1.2), Inches(0.2),
                 f"({pct} del total)", font_size=8, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide6, Inches(0.8), Inches(2.9), Inches(5.4), Inches(0.35),
             "Q→P: 10.8%", font_size=16, bold=True, color=KAVAK_BLUE, alignment=PP_ALIGN.CENTER)

# Aliado Card
ali_card = add_rounded_rect(slide6, Inches(6.9), Inches(1.2), Inches(5.8), Inches(2.5), fill_color=WHITE, line_color=PURPLE)
add_shape(slide6, Inches(6.95), Inches(1.25), Inches(5.7), Inches(0.08), fill_color=PURPLE)
add_text_box(slide6, Inches(7.1), Inches(1.45), Inches(5.4), Inches(0.35),
             "ALIADO (Partners)", font_size=18, bold=True, color=PURPLE)

ali_metrics = [("Quotes", "741", "81.6%"), ("Made", "170", "77.3%"),
               ("Approved", "141", "75.4%"), ("Purchased", "51", "73.9%")]
for i, (label, val, pct) in enumerate(ali_metrics):
    x = Inches(7.1) + Inches(1.35) * i
    add_text_box(slide6, x, Inches(1.9), Inches(1.2), Inches(0.2),
                 label, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, Inches(2.1), Inches(1.2), Inches(0.35),
                 val, font_size=22, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, Inches(2.5), Inches(1.2), Inches(0.2),
                 f"({pct} del total)", font_size=8, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide6, Inches(7.1), Inches(2.9), Inches(5.4), Inches(0.35),
             "Q→P: 6.9%", font_size=16, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)

# Historical Aliado growth table
ali_hist_headers = ["Canal", "Sep 25", "Oct 25", "Nov 25", "Dic 25", "Ene 26", "Feb 26"]
# For historical TAS vs Aliado, we derive from total - aliado.
# Sep: total 21 purchased, TAS only (no aliado context before). Let's show total + aliado where available.
# Using approximate data: Aliado started growing significantly from Oct.
ali_hist = [
    ["Aliado Compras", "—", "—", "—", "—", "—", "51"],
    ["TAS Compras", "—", "—", "—", "—", "—", "18"],
    ["Aliado % Total", "—", "—", "—", "—", "—", "73.9%"],
]

# Simpler: show a summary table of the split for Feb only since we only have that breakdown
split_shape = slide6.shapes.add_table(4, 6, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.4))
split_tbl = split_shape.table

split_headers = ["Canal", "Quotes", "Made", "Approved", "Purchased", "Q→P"]
for i, h in enumerate(split_headers):
    set_cell_text(split_tbl.cell(0, i), h, font_size=9, bold=True, color=WHITE)
style_table_header(split_tbl, 6)

split_rows = [
    ["TAS", "167", "50", "46", "18", "10.8%"],
    ["Aliado", "741", "170", "141", "51", "6.9%"],
    ["TOTAL", "908", "220", "187", "69", "7.6%"],
]
for r, row in enumerate(split_rows, start=1):
    bold = (r == 3)
    for c, val in enumerate(row):
        clr = KAVAK_BLUE if r==1 and c>0 else (PURPLE if r==2 and c>0 else KAVAK_DARK)
        set_cell_text(split_tbl.cell(r, c), val, font_size=10, bold=bold, color=clr,
                      alignment=PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER)
style_table_rows(split_tbl, 4, 6)

split_tbl.columns[0].width = Inches(2.0)
for c in range(1, 6):
    split_tbl.columns[c].width = Inches(2.0)

# Insight
add_insight_box(slide6, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.8),
                ">> Aliado es el motor de crecimiento: 81.6% de quotes y 73.9% de compras. Sin embargo, su Q→P (6.9%) es 3.9pp menor que TAS (10.8%). Cada punto porcentual de mejora en Q→P de Aliado = ~7 compras adicionales. Oportunidad clara de coaching y soporte a partners.",
                accent_color=PURPLE, bg_color=PURPLE_LIGHT)

add_text_box(slide6, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FUNNEL DIAGNOSTIC
# ═══════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
add_slide_title(slide7, "Diagnostico del Funnel — Febrero 2026", "Donde se pierden los autos? Analisis de conversion por etapa")

# Funnel stages as descending cards
funnel_stages = [
    ("QUOTES", "908", "100%", KAVAK_BLUE, Inches(9.0)),
    ("MADE", "220", "24.2%", KAVAK_BLUE, Inches(7.5)),
    ("APPROVED", "187", "85.0%", GREEN, Inches(6.0)),
    ("PURCHASED", "69", "36.9%", GREEN, Inches(4.5)),
]

funnel_x = Inches(0.8)
for i, (stage, count, conv, color, width) in enumerate(funnel_stages):
    y = Inches(1.3) + Inches(1.2) * i
    # Centered funnel bar
    offset = (Inches(9.0) - width) / 2
    bar = add_rounded_rect(slide7, funnel_x + offset, y, width, Inches(0.7), fill_color=color)
    bar.adjustments[0] = 0.15
    # Stage name and count
    add_text_box(slide7, funnel_x + offset + Inches(0.2), y + Inches(0.05), Inches(3), Inches(0.3),
                 f"{stage}: {count}", font_size=16, bold=True, color=WHITE)
    # Conversion label on right
    if i > 0:
        prev_stage = funnel_stages[i-1][0]
        label_map = {1: "Q→M", 2: "M→A", 3: "A→P"}
        add_text_box(slide7, funnel_x + offset + width + Inches(0.15), y + Inches(0.1), Inches(2), Inches(0.25),
                     f"{label_map[i]}: {conv}", font_size=13, bold=True, color=color)
    # Drop count
    if i > 0:
        prev_count = int(funnel_stages[i-1][1])
        curr_count = int(count)
        drop = prev_count - curr_count
        add_text_box(slide7, funnel_x + offset + width + Inches(0.15), y + Inches(0.4), Inches(2.5), Inches(0.2),
                     f"-{drop} perdidos", font_size=10, color=RED)

# A→P Historical comparison
ap_title = add_text_box(slide7, Inches(8.5), Inches(1.3), Inches(4), Inches(0.3),
                        "A→P Historico (tendencia descendente):", font_size=11, bold=True, color=RED)

ap_data = [("Sep 25", "61.8%"), ("Oct 25", "34.8%"), ("Nov 25", "32.1%"),
           ("Dic 25", "45.5%"), ("Ene 26", "40.0%"), ("Feb 26", "36.9%")]
ap_shape = slide7.shapes.add_table(7, 2, Inches(8.5), Inches(1.7), Inches(3.5), Inches(2.5))
ap_tbl = ap_shape.table
ap_tbl.columns[0].width = Inches(1.5)
ap_tbl.columns[1].width = Inches(2.0)
set_cell_text(ap_tbl.cell(0, 0), "Mes", font_size=9, bold=True, color=WHITE)
set_cell_text(ap_tbl.cell(0, 1), "A→P %", font_size=9, bold=True, color=WHITE)
style_table_header(ap_tbl, 2)
for r, (m, v) in enumerate(ap_data, start=1):
    set_cell_text(ap_tbl.cell(r, 0), m, font_size=9)
    is_last = (r == len(ap_data))
    set_cell_text(ap_tbl.cell(r, 1), v, font_size=9, bold=is_last, color=RED if is_last else DARK_GRAY)
style_table_rows(ap_tbl, 7, 2)

# Impact box
add_insight_box(slide7, Inches(8.5), Inches(4.4), Inches(4.2), Inches(0.9),
                ">> Impacto simulado\nSi A→P fuera 45% (nivel Dic):\n187 × 45% = 84 compras (+15 vs actual)",
                accent_color=AMBER, bg_color=AMBER_LIGHT)

# Main insight
add_insight_box(slide7, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.8),
                ">> Cuello de botella: A→P baja a 36.9%. De 187 aprobados, solo 69 se compraron — 118 vehiculos aprobados se perdieron. La calidad de oferta/precio en la etapa final es el factor limitante. Q→M y M→A estan saludables.",
                accent_color=RED, bg_color=RED_LIGHT)

add_text_box(slide7, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ACTION PLAN MARCH
# ═══════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_slide_title(slide8, "Plan de Accion — Marzo 2026", "Meta: 90 compras netas (+30.4% vs Feb)")

# Target card
target_card = add_rounded_rect(slide8, Inches(0.6), Inches(1.1), Inches(3.5), Inches(1.2), fill_color=GREEN, line_color=GREEN)
add_text_box(slide8, Inches(0.8), Inches(1.2), Inches(3.1), Inches(0.25),
             "META MARZO 2026", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide8, Inches(0.8), Inches(1.5), Inches(3.1), Inches(0.5),
             "90 compras netas", font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide8, Inches(0.8), Inches(2.0), Inches(3.1), Inches(0.2),
             "+30.4% vs Feb (69)", font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

# Feb reference
feb_card = add_rounded_rect(slide8, Inches(4.5), Inches(1.1), Inches(2.2), Inches(1.2), fill_color=LIGHT_GRAY)
add_text_box(slide8, Inches(4.6), Inches(1.2), Inches(2.0), Inches(0.25),
             "FEB 2026 (ACTUAL)", font_size=9, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide8, Inches(4.6), Inches(1.5), Inches(2.0), Inches(0.5),
             "69", font_size=28, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)

# Palancas
palanca_data = [
    ("1", "Mantener volumen de quotes >900", "+0 riesgo — fue 989 en Ene, 908 en Feb. Sostener flujo de entrada.", KAVAK_BLUE, "Base"),
    ("2", "Mejorar A→P de 36.9% a 42%", "Focus en pricing/calidad de oferta para vehiculos aprobados. 187 × 42% = 79 → +10 compras", RED, "+10"),
    ("3", "Mejorar Q→P Aliado de 6.9% a 8.5%", "Training y soporte a partners aliados. 741 × 8.5% = 63 vs 51 actual → +12 compras", PURPLE, "+12"),
    ("4", "Sostener Q→M >24%", "Mas inspecciones del volumen entrante. Cada 5pp en Q→M = ~45 made adicionales → +3 compras", GREEN, "+3"),
]

for i, (num, title, desc, color, impact) in enumerate(palanca_data):
    y = Inches(2.6) + Inches(0.85) * i
    # Number badge
    badge = add_rounded_rect(slide8, Inches(0.6), y, Inches(0.45), Inches(0.45), fill_color=color)
    add_text_box(slide8, Inches(0.6), y + Inches(0.05), Inches(0.45), Inches(0.35),
                 num, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide8, Inches(1.2), y, Inches(6), Inches(0.3),
                 title, font_size=12, bold=True, color=KAVAK_DARK)
    # Description
    add_text_box(slide8, Inches(1.2), y + Inches(0.3), Inches(7), Inches(0.4),
                 desc, font_size=9, color=MED_GRAY)
    # Impact badge
    imp_card = add_rounded_rect(slide8, Inches(8.5), y + Inches(0.05), Inches(1.2), Inches(0.5), fill_color=color)
    add_text_box(slide8, Inches(8.5), y + Inches(0.1), Inches(1.2), Inches(0.4),
                 impact, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Waterfall summary
wf_card = add_rounded_rect(slide8, Inches(10.2), Inches(1.1), Inches(2.8), Inches(4.8), fill_color=BLUE_LIGHT, line_color=KAVAK_BLUE)
add_text_box(slide8, Inches(10.4), Inches(1.2), Inches(2.4), Inches(0.3),
             "WATERFALL", font_size=12, bold=True, color=KAVAK_BLUE, alignment=PP_ALIGN.CENTER)

wf_items = [("Base Feb", "69", KAVAK_DARK), ("+ A→P mejora", "+10", RED),
            ("+ Aliado eficiencia", "+12", PURPLE), ("+ Q→M sostenido", "+3", GREEN),
            ("= TOTAL", "94", KAVAK_BLUE)]
for i, (label, val, color) in enumerate(wf_items):
    y = Inches(1.7) + Inches(0.55) * i
    is_total = (i == len(wf_items) - 1)
    if is_total:
        add_shape(slide8, Inches(10.4), y - Inches(0.05), Inches(2.4), Inches(0.02), fill_color=KAVAK_BLUE)
    add_text_box(slide8, Inches(10.4), y, Inches(1.6), Inches(0.3),
                 label, font_size=10, bold=is_total, color=color)
    add_text_box(slide8, Inches(12.0), y, Inches(0.8), Inches(0.3),
                 val, font_size=14 if is_total else 12, bold=True, color=color, alignment=PP_ALIGN.RIGHT)

# Conservative range
add_text_box(slide8, Inches(10.4), Inches(4.7), Inches(2.4), Inches(0.4),
             "Rango conservador:\n85 – 94 compras", font_size=10, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_insight_box(slide8, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
                ">> Meta ambiciosa pero la trayectoria la soporta: 8 meses de crecimiento continuo. Las palancas son especificas y medibles. El rango conservador (85-94) ya representaria un nuevo record.",
                accent_color=GREEN, bg_color=GREEN_LIGHT)

add_text_box(slide8, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KPIs TO MONITOR
# ═══════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_slide_title(slide9, "KPIs a Monitorear — Marzo 2026", "Indicadores clave y acciones inmediatas para semana 1")

# KPI Table
kpi_headers = ["Indicador", "Feb Actual", "Meta Mar", "Prioridad", "Owner"]
kpi_rows_data = [
    ["Compras Netas",     "69",     "90",     "CRITICA",   "Celula"],
    ["Quotes",            "908",    ">900",   "ALTA",      "Aliado + TAS"],
    ["Q→M %",            "24.2%",  ">24%",   "ALTA",      "TAS Manager"],
    ["M→A %",            "85.0%",  ">87%",   "MEDIA",     "Inspecciones"],
    ["A→P %",            "36.9%",  "42%",    "CRITICA",   "Pricing/Ops"],
    ["Q→P %",            "7.6%",   ">9%",    "ALTA",      "Celula"],
    ["Q→P Aliado",       "6.9%",   "8.5%",   "CRITICA",   "Partner Mgmt"],
    ["Q→P TAS",          "10.8%",  ">11%",   "MEDIA",     "TAS Manager"],
    ["Dias con 0 compras","6/28",  "<5/31",  "MEDIA",     "Ops"],
]

kpi_shape = slide9.shapes.add_table(len(kpi_rows_data)+1, 5, Inches(0.6), Inches(1.1), Inches(12.1), Inches(3.6))
kpi_tbl = kpi_shape.table

kpi_tbl.columns[0].width = Inches(2.8)
kpi_tbl.columns[1].width = Inches(2.0)
kpi_tbl.columns[2].width = Inches(2.0)
kpi_tbl.columns[3].width = Inches(2.3)
kpi_tbl.columns[4].width = Inches(3.0)

for i, h in enumerate(kpi_headers):
    set_cell_text(kpi_tbl.cell(0, i), h, font_size=10, bold=True, color=WHITE)
style_table_header(kpi_tbl, 5)

for r, row in enumerate(kpi_rows_data, start=1):
    set_cell_text(kpi_tbl.cell(r, 0), row[0], font_size=10, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.LEFT)
    set_cell_text(kpi_tbl.cell(r, 1), row[1], font_size=10)
    set_cell_text(kpi_tbl.cell(r, 2), row[2], font_size=10, bold=True, color=KAVAK_BLUE)
    # Priority coloring
    prio = row[3]
    prio_color = RED if prio == "CRITICA" else (AMBER if prio == "ALTA" else MED_GRAY)
    set_cell_text(kpi_tbl.cell(r, 3), prio, font_size=10, bold=True, color=prio_color)
    set_cell_text(kpi_tbl.cell(r, 4), row[4], font_size=10, color=MED_GRAY)
style_table_rows(kpi_tbl, len(kpi_rows_data)+1, 5)

# Immediate actions for Week 1
add_text_box(slide9, Inches(0.6), Inches(4.9), Inches(5), Inches(0.3),
             "ACCIONES INMEDIATAS — SEMANA 1 DE MARZO", font_size=12, bold=True, color=KAVAK_DARK)

actions = [
    ("1", "Revisar los 118 aprobados no comprados de Feb — identificar razones de caida (precio, condicion, timing)", RED),
    ("2", "Sesion de coaching con top 5 aliados por volumen — compartir mejores practicas de conversion", PURPLE),
    ("3", "Establecer daily standup de compras con meta diaria de 3+ compras/dia", GREEN),
    ("4", "Auditar rechazos de M→A (33 en Feb) — identificar patrones y reducir tasa de rechazo", AMBER),
]

for i, (num, text, color) in enumerate(actions):
    y = Inches(5.25) + Inches(0.35) * i
    badge = add_rounded_rect(slide9, Inches(0.6), y, Inches(0.3), Inches(0.3), fill_color=color)
    add_text_box(slide9, Inches(0.6), y + Inches(0.02), Inches(0.3), Inches(0.25),
                 num, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, Inches(1.05), y + Inches(0.02), Inches(11.5), Inches(0.3),
                 text, font_size=10, color=DARK_GRAY)

add_text_box(slide9, Inches(0.6), Inches(6.9), Inches(5), Inches(0.3),
             "Continental Coapa — Febrero 2026", font_size=8, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
# Dark background
add_shape(slide10, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=KAVAK_DARK)
add_shape(slide10, Inches(0), Inches(0), SLIDE_W, Inches(0.12), fill_color=KAVAK_BLUE)

# Blue accent line
add_shape(slide10, Inches(1.2), Inches(2.8), Inches(1.5), Inches(0.06), fill_color=KAVAK_BLUE)

# Summary title
add_text_box(slide10, Inches(1.2), Inches(1.5), Inches(10), Inches(0.5),
             "Continental — Kia Coapa", font_size=14, bold=True, color=KAVAK_BLUE)
add_text_box(slide10, Inches(1.2), Inches(2.0), Inches(10), Inches(0.7),
             "Febrero 2026 — Resumen", font_size=32, bold=True, color=WHITE)

# Three key summary points
summary_points = [
    ("69", "compras netas", "Record historico de la celula", GREEN),
    ("8", "meses de crecimiento", "Consecutivos desde julio 2025", KAVAK_BLUE),
    ("90", "meta marzo", "Rango conservador: 85–94", AMBER),
]

for i, (big_num, label, desc, color) in enumerate(summary_points):
    x = Inches(1.2) + Inches(3.8) * i
    y = Inches(3.3)
    # Big number
    add_text_box(slide10, x, y, Inches(3.0), Inches(0.7),
                 big_num, font_size=48, bold=True, color=color)
    add_text_box(slide10, x, y + Inches(0.7), Inches(3.0), Inches(0.3),
                 label, font_size=14, bold=True, color=WHITE)
    add_text_box(slide10, x, y + Inches(1.0), Inches(3.0), Inches(0.3),
                 desc, font_size=11, color=MED_GRAY)

# Bottom
add_text_box(slide10, Inches(1.2), Inches(5.8), Inches(10), Inches(0.3),
             "El camino a 90 compras: mejorar A→P, potenciar aliados, sostener volumen.", font_size=13, color=WHITE)

add_shape(slide10, Inches(1.2), Inches(6.3), Inches(11), Inches(0.01), fill_color=MED_GRAY)

add_text_box(slide10, Inches(1.2), Inches(6.5), Inches(5), Inches(0.3),
             "TAS Manager Analytics", font_size=11, color=MED_GRAY)
add_text_box(slide10, Inches(8.5), Inches(6.5), Inches(4.5), Inches(0.3),
             "1 de Marzo 2026", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

# ── Save ─────────────────────────────────────────────────────────────────
output_path = "/Users/danielvilladiego/Downloads/KiaCoapa_MBR_Febrero_2026.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
