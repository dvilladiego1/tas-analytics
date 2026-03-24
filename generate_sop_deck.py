#!/usr/bin/env python3
"""
S&Op Resource Request Deck Generator
Cells: ISMO Leon & KIA Coapa
Purpose: Justify 2 inspectors + 2 lead managers per cell
"""

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from datetime import datetime
import os

# ── Colors ──────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
BRAND_BLUE = RGBColor(0x00, 0x4E, 0x98)
TBL_HDR    = RGBColor(0x2B, 0x47, 0x8B)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED        = RGBColor(0xEF, 0x44, 0x44)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
BLUE_MED   = RGBColor(0x3B, 0x7D, 0xDD)
TEAL       = RGBColor(0x00, 0xB4, 0x8A)

# Insight box colors
HIGH_BG    = RGBColor(0xDC, 0xFC, 0xE7)
HIGH_BAR   = RGBColor(0x22, 0xC5, 0x5E)
LOW_BG     = RGBColor(0xFE, 0xE2, 0xE2)
LOW_BAR    = RGBColor(0xEF, 0x44, 0x44)
SUM_BG     = RGBColor(0xDB, 0xEA, 0xFE)
SUM_BAR    = RGBColor(0x00, 0x4E, 0x98)
LECT_BG    = RGBColor(0xE3, 0xEE, 0xFB)

# ── Load Data ───────────────────────────────────────────────────────────
CSV_PATH = os.path.expanduser("~/Downloads/Summary MKP _ TAS & BULK - Snapshot_Diario (1).csv")
df = pd.read_csv(CSV_PATH, low_memory=False)
df['fecha_creacion'] = pd.to_datetime(df['Fecha de creación'], format='%d/%m/%Y', errors='coerce')
df['purchase_dt'] = pd.to_datetime(df['purchase_date'], format='%d/%m/%Y', errors='coerce')
df['fecha_inspeccion'] = pd.to_datetime(df['Fecha Inspección'], format='%d/%m/%Y', errors='coerce')
df['made_dt'] = pd.to_datetime(df['made'], format='%d/%m/%Y', errors='coerce')
df['scheduled_dt'] = pd.to_datetime(df['scheduled'], format='%d/%m/%Y', errors='coerce')

CH = 'BULK\nTAS\nALIADO'
SLA_COL = 'SLA Cotizado a Inspeccionando'

# ── Helper Functions ────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_top_bar(slide, color=BRAND_BLUE, h=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_footer(slide, text="S&Op Resource Request — TAS Marketplace", date_str="Marzo 2026"):
    add_textbox(slide, 0.5, 7.05, 8, 0.3, text, 8, False, LIGHT_GRAY)
    add_textbox(slide, 10.5, 7.05, 2.5, 0.3, date_str, 8, False, LIGHT_GRAY, PP_ALIGN.RIGHT)

def add_kpi_card(slide, left, top, width, height, label, value, delta=None, accent_color=BRAND_BLUE):
    """Add a KPI card with accent bar on top"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Label
    add_textbox(slide, left + 0.15, top + 0.12, width - 0.3, 0.25, label, 9, True, LIGHT_GRAY)
    # Value
    add_textbox(slide, left + 0.15, top + 0.35, width - 0.3, 0.45, str(value), 26, True, DARK_TEXT)
    # Delta
    if delta:
        d_color = GREEN if delta.startswith('+') or delta.startswith('↑') else RED
        add_textbox(slide, left + 0.15, top + 0.75, width - 0.3, 0.25, delta, 10, True, d_color)

def add_insight_box(slide, left, top, width, height, title, text, bg_color, bar_color):
    """Add an insight box with left accent bar"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    add_textbox(slide, left + 0.15, top + 0.05, width - 0.25, 0.2, title, 9, True, bar_color)
    add_textbox(slide, left + 0.15, top + 0.22, width - 0.25, height - 0.27, text, 9, False, DARK_TEXT)

def make_table(slide, left, top, width, height, data, col_widths=None):
    """data = list of lists, first row = headers"""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
                    paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_HDR
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF1, 0xF5, 0xF9)

    return table_shape

# ── Compute Data ────────────────────────────────────────────────────────
cells_config = {
    'ISMO LEON': {
        'celulas': ['ISMO LEON', 'ALIADO LEON'],
        'tas_celula': 'ISMO LEON',
        'ali_celula': 'ALIADO LEON',
        'grupo': 'GRUPO ISMO',
        'region': 'León, Guanajuato',
        'current_inspectors': 1,
        'current_lms': 1,
        'proposed_inspectors': 2,
        'proposed_lms': 2,
    },
    'KIA COAPA': {
        'celulas': ['CONTINENTAL COAPA', 'ALIADO COAPA'],
        'tas_celula': 'CONTINENTAL COAPA',
        'ali_celula': 'ALIADO COAPA',
        'grupo': 'GRUPO CONTINENTAL',
        'region': 'CDMX Sur',
        'current_inspectors': 1,  # main inspector
        'current_lms': 1,
        'proposed_inspectors': 2,
        'proposed_lms': 2,
    },
}

months = [('Oct 2025', 2025, 10), ('Nov 2025', 2025, 11), ('Dec 2025', 2025, 12),
          ('Jan 2026', 2026, 1), ('Feb 2026', 2026, 2)]

def get_funnel(cell_data, yr, mn, channel=None):
    if channel:
        month_data = cell_data[(cell_data['fecha_creacion'].dt.month == mn) &
                               (cell_data['fecha_creacion'].dt.year == yr) &
                               (cell_data[CH] == channel)]
    else:
        month_data = cell_data[(cell_data['fecha_creacion'].dt.month == mn) &
                               (cell_data['fecha_creacion'].dt.year == yr)]

    quotes = month_data[month_data['scheduled.1'] == 1].shape[0]
    made = month_data[month_data['made.1'] == 1].shape[0]
    approved = month_data[month_data['approved'] == 1].shape[0]

    # Purchases by purchase_date
    if channel:
        purch_data = cell_data[(cell_data['purchase_dt'].dt.month == mn) &
                               (cell_data['purchase_dt'].dt.year == yr) &
                               (cell_data['purchased'] == 1) &
                               (cell_data[CH] == channel)]
    else:
        purch_data = cell_data[(cell_data['purchase_dt'].dt.month == mn) &
                               (cell_data['purchase_dt'].dt.year == yr) &
                               (cell_data['purchased'] == 1)]
    purchased = purch_data.shape[0]

    return {'Q': quotes, 'M': made, 'A': approved, 'P': purchased,
            'Q2M': made/quotes*100 if quotes else 0,
            'M2A': approved/made*100 if made else 0,
            'A2P': purchased/approved*100 if approved else 0}

def get_daily_inspections(cell_data, yr, mn):
    feb_data = cell_data[(cell_data['fecha_inspeccion'].dt.month == mn) &
                         (cell_data['fecha_inspeccion'].dt.year == yr) &
                         (cell_data['made.1'] == 1)]
    daily = feb_data.groupby(feb_data['fecha_inspeccion'].dt.date).size()
    return daily

# ── Build Presentation ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)
add_top_bar(slide, BRAND_BLUE, 0.12)

# Separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06))
sep.fill.solid()
sep.fill.fore_color.rgb = BRAND_BLUE
sep.line.fill.background()

add_textbox(slide, 1.0, 1.5, 11, 0.5, "TAS Marketplace — S&Op Resource Request", 14, True, BRAND_BLUE)
add_textbox(slide, 1.0, 2.1, 11, 0.8, "Plan de Contratación: ISMO León & Kia Coapa", 36, True, WHITE)
add_textbox(slide, 1.0, 3.3, 11, 0.5, "Marzo 2026", 22, False, WHITE)
add_textbox(slide, 1.0, 4.2, 11, 0.5, "Justificación basada en datos para expansión de capacidad operativa", 14, False, LIGHT_GRAY)
add_textbox(slide, 1.0, 5.0, 11, 0.4, "Análisis de capacidad instalada • Demanda perdida en Made • Plan de headcount", 13, False, LIGHT_GRAY)
add_textbox(slide, 9.5, 6.8, 3.5, 0.3, "Kavak Marketplace B2B — Programa TAS", 11, False, LIGHT_GRAY, PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY - THE ASK
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_bar(slide, BRAND_BLUE, 0.06)

add_textbox(slide, 0.5, 0.25, 12, 0.5, "El Ask: Expansión de Capacidad para Capturar Demanda", 24, True, DARK_TEXT)
add_textbox(slide, 0.5, 0.7, 12, 0.4, "Ambas células están perdiendo compras por restricciones de capacidad en inspectores y lead managers", 12, False, LIGHT_GRAY)

# The ASK cards
for i, (cell, config) in enumerate(cells_config.items()):
    x = 0.5 + i * 6.3
    # Cell header
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(5.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    add_textbox(slide, x + 0.2, 1.4, 5.4, 0.4, f"📍 {cell} — {config['region']}", 16, True, BRAND_BLUE)

    # Current vs Proposed
    add_textbox(slide, x + 0.2, 1.85, 2.5, 0.25, "ACTUAL", 10, True, LIGHT_GRAY)
    add_textbox(slide, x + 3.0, 1.85, 2.5, 0.25, "PROPUESTO", 10, True, GREEN)

    add_textbox(slide, x + 0.2, 2.1, 2.5, 0.3, f"{config['current_inspectors']} Inspector + {config['current_lms']} LM", 14, True, RED)
    add_textbox(slide, x + 3.0, 2.1, 2.5, 0.3, f"{config['proposed_inspectors']} Inspectores + {config['proposed_lms']} LMs", 14, True, GREEN)

    # Key metric
    cell_data = df[df['Célula'].isin(config['celulas'])]
    feb_total = get_funnel(cell_data, 2026, 2)

    add_textbox(slide, x + 0.2, 2.6, 5.4, 0.25,
                f"Feb: {feb_total['Q']} quotes → {feb_total['M']} inspecciones → {feb_total['P']} compras",
                11, False, DARK_TEXT)
    add_textbox(slide, x + 0.2, 2.85, 5.4, 0.25,
                f"Q→M: {feb_total['Q2M']:.0f}% | Demanda perdida: {feb_total['Q'] - feb_total['M']} leads sin inspeccionar",
                11, True, RED)
    add_textbox(slide, x + 0.2, 3.1, 5.4, 0.25,
                f"Seguimiento de pipe aliado se haría desde central",
                10, False, LIGHT_GRAY)

# Bottom section - Operating Model
add_textbox(slide, 0.5, 4.1, 12, 0.35, "Modelo Operativo Propuesto", 16, True, DARK_TEXT)

# 3 pillars
pillars = [
    ("Inspector Dedicado TAS", "Enfocado en capturar demanda TAS\n(marketing + walk-ins)\nMeta: 6-8 insp/día TAS", BRAND_BLUE),
    ("Inspector Dedicado Aliado", "Atiende volumen de aliados\nen agencia y domicilio\nMeta: 6-8 insp/día Aliado", BLUE_MED),
    ("2 Lead Managers en Piso", "LM1: Seguimiento pipe TAS\nLM2: Seguimiento pipe Aliado\nPipe aliado migra a central", GREEN),
]

for i, (title, desc, color) in enumerate(pillars):
    x = 0.5 + i * 4.2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.55), Inches(3.9), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = color
    box.line.width = Pt(2)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.55), Inches(3.9), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_textbox(slide, x + 0.15, 4.65, 3.6, 0.3, title, 12, True, color)
    add_textbox(slide, x + 0.15, 4.95, 3.6, 1.0, desc, 10, False, DARK_TEXT)

# Investment summary
add_textbox(slide, 0.5, 6.35, 12, 0.3,
            "Inversión total: +2 inspectores + 2 LMs (piso) | Ahorro: seguimiento aliado migra a equipo central",
            11, True, BRAND_BLUE)

add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: ISMO LEON - DIAGNÓSTICO DE CAPACIDAD
# ════════════════════════════════════════════════════════════════════════
for cell_name, config in cells_config.items():
    cell_data = df[df['Célula'].isin(config['celulas'])].copy()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, BRAND_BLUE, 0.06)

    add_textbox(slide, 0.5, 0.2, 12, 0.5, f"{cell_name} — Diagnóstico de Capacidad", 24, True, DARK_TEXT)
    add_textbox(slide, 0.5, 0.65, 12, 0.3, f"{config['region']} | {config['current_inspectors']} inspector actual | Capacidad máx: 8 insp/día", 11, False, LIGHT_GRAY)

    # Funnel Evolution Table
    funnel_data = [['Mes', 'Quotes', 'Made', 'Approved', 'Purchased', 'Q→M%', 'M→A%', 'Demanda Perdida']]
    for mname, yr, mn in months:
        for ch_label, ch_val in [('TAS', 'TAS'), ('ALIADO', 'ALIADO')]:
            f = get_funnel(cell_data, yr, mn, ch_val)
            lost = f['Q'] - f['M']
            funnel_data.append([f"{mname} {ch_label}", f['Q'], f['M'], f['A'], f['P'],
                               f"{f['Q2M']:.0f}%", f"{f['M2A']:.0f}%", lost])

    # Simplified table: total per month
    total_data = [['Mes', 'Q Total', 'M Total', 'Q TAS', 'M TAS', 'Q→M TAS', 'Q Ali', 'M Ali', 'Q→M Ali', 'P Total', 'Lost']]
    for mname, yr, mn in months:
        ft = get_funnel(cell_data, yr, mn)
        ftas = get_funnel(cell_data, yr, mn, 'TAS')
        fali = get_funnel(cell_data, yr, mn, 'ALIADO')
        lost = ft['Q'] - ft['M']
        total_data.append([mname, ft['Q'], ft['M'],
                          ftas['Q'], ftas['M'], f"{ftas['Q2M']:.0f}%",
                          fali['Q'], fali['M'], f"{fali['Q2M']:.0f}%",
                          ft['P'], lost])

    make_table(slide, 0.5, 1.1, 12.3, 2.5, total_data)

    # Daily inspection heatmap for Feb 2026
    add_textbox(slide, 0.5, 3.7, 6, 0.3, "Inspecciones Diarias — Feb 2026 (1 inspector)", 14, True, DARK_TEXT)

    daily = get_daily_inspections(cell_data, 2026, 2)
    daily_data = [['Día', 'Total', 'TAS', 'Aliado', 'Status']]

    feb_made = cell_data[(cell_data['fecha_inspeccion'].dt.month == 2) &
                         (cell_data['fecha_inspeccion'].dt.year == 2026) &
                         (cell_data['made.1'] == 1)]

    for date in sorted(daily.index):
        day_data = feb_made[feb_made['fecha_inspeccion'].dt.date == date]
        tas_n = (day_data[CH] == 'TAS').sum()
        ali_n = (day_data[CH] == 'ALIADO').sum()
        total = tas_n + ali_n
        status = "🔴 FULL" if total >= 8 else ("🟡 HIGH" if total >= 6 else "🟢 OK")
        daily_data.append([str(date), total, tas_n, ali_n, status])

    # Split into 2 columns of ~12 days each
    mid = len(daily_data) // 2 + 1
    if len(daily_data) > 13:
        left_data = daily_data[:mid]
        right_data = [daily_data[0]] + daily_data[mid:]
        make_table(slide, 0.5, 4.05, 5.8, 2.8, left_data)
        make_table(slide, 6.5, 4.05, 5.8, 2.8, right_data)
    else:
        make_table(slide, 0.5, 4.05, 6, 2.8, daily_data)

    # Capacity summary box
    days_full = (daily >= 8).sum()
    days_high = (daily >= 6).sum()
    avg_daily = daily.mean()

    add_insight_box(slide, 0.5, 6.9, 4, 0.45, "CAPACIDAD",
                   f"Avg: {avg_daily:.1f}/día | {days_high}/{len(daily)} días ≥6 | {days_full}/{len(daily)} días ≥8 (FULL)",
                   LOW_BG, LOW_BAR)

    add_insight_box(slide, 4.7, 6.9, 4, 0.45, "SLA",
                   f"TAS: ~2 días Q→Insp | ALIADO: ~4 días Q→Insp",
                   SUM_BG, SUM_BAR)

    tas_feb = get_funnel(cell_data, 2026, 2, 'TAS')
    ali_feb = get_funnel(cell_data, 2026, 2, 'ALIADO')
    add_insight_box(slide, 8.9, 6.9, 4, 0.45, "DEMANDA PERDIDA FEB",
                   f"TAS: {tas_feb['Q']-tas_feb['M']} sin inspeccionar | ALI: {ali_feb['Q']-ali_feb['M']} sin inspeccionar",
                   LOW_BG, LOW_BAR)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5-6: DEMAND LOSS & OPPORTUNITY (one per cell)
# ════════════════════════════════════════════════════════════════════════
for cell_name, config in cells_config.items():
    cell_data = df[df['Célula'].isin(config['celulas'])].copy()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, BRAND_BLUE, 0.06)

    add_textbox(slide, 0.5, 0.2, 12, 0.5, f"{cell_name} — Demanda Perdida & Oportunidad", 24, True, DARK_TEXT)

    # Problem statement
    jan = get_funnel(cell_data, 2026, 1)
    feb = get_funnel(cell_data, 2026, 2)

    jan_tas = get_funnel(cell_data, 2026, 1, 'TAS')
    feb_tas = get_funnel(cell_data, 2026, 2, 'TAS')
    jan_ali = get_funnel(cell_data, 2026, 1, 'ALIADO')
    feb_ali = get_funnel(cell_data, 2026, 2, 'ALIADO')

    # 3 Problem Cards
    problems = [
        ("1. Capacidad Instalada Limitada",
         f"1 inspector hace max 8 insp/día.\n"
         f"Feb: promedio {get_daily_inspections(cell_data, 2026, 2).mean():.1f} insp/día (TAS+Aliado combinados).\n"
         f"En días pico llega a {get_daily_inspections(cell_data, 2026, 2).max()}, bloqueando nuevas citas.\n"
         f"Resultado: {feb['Q']-feb['M']} quotes sin inspeccionar en Feb.",
         LOW_BG, LOW_BAR),
        ("2. Lead Manager Saturado",
         f"1 LM atiende TAS + Aliado simultáneamente.\n"
         f"TAS Feb: solo {feb_tas['Q']} quotes (vs {jan_tas['Q']} en Ene).\n"
         f"El foco en aliados limita el crecimiento de cotizaciones TAS.\n"
         f"Con 2 LMs: uno dedicado a TAS, otro a seguimiento en piso.",
         LOW_BG, LOW_BAR),
        ("3. Marketing Limitado por CX",
         f"No se puede invertir más en mktg porque la experiencia\n"
         f"se degrada: SLA de 4+ días para inspección en aliados.\n"
         f"Generar más demanda sin capacidad = peor NPS.\n"
         f"Con 2 inspectores: se desbloquea inversión en mktg.",
         LOW_BG, LOW_BAR),
    ]

    for i, (title, desc, bg, bar) in enumerate(problems):
        x = 0.5 + i * 4.2
        add_insight_box(slide, x, 0.8, 3.9, 2.2, title, desc, bg, bar)

    # Opportunity Table
    add_textbox(slide, 0.5, 3.2, 12, 0.35, "Escenario de Oportunidad — Con 2 Inspectores + 2 LMs", 16, True, GREEN)

    # Current vs Target
    opp_data = [
        ['Métrica', 'Actual (Feb)', 'Target con 2I+2LM', 'Delta', 'Impacto'],
        ['Inspecciones/día', f"{get_daily_inspections(cell_data, 2026, 2).mean():.1f}", '12-16',
         f"+{16 - get_daily_inspections(cell_data, 2026, 2).mean():.0f}", '2x capacidad'],
        ['Quotes TAS/mes', str(feb_tas['Q']), str(int(feb_tas['Q'] * 2.5)),
         f"+{int(feb_tas['Q'] * 1.5)}", 'LM dedicado + mktg'],
        ['Made TAS/mes', str(feb_tas['M']), str(int(feb_tas['Q'] * 2.5 * 0.35)),
         f"+{int(feb_tas['Q'] * 2.5 * 0.35) - feb_tas['M']}", 'Inspector dedicado'],
        ['Made Aliado/mes', str(feb_ali['M']), str(int(feb_ali['Q'] * 0.30)),
         f"+{int(feb_ali['Q'] * 0.30) - feb_ali['M']}", 'Inspector dedicado'],
        ['Compras TAS/mes', str(feb_tas['P']), str(int(feb_tas['Q'] * 2.5 * 0.35 * 0.85 * 0.45)),
         f"+{int(feb_tas['Q'] * 2.5 * 0.35 * 0.85 * 0.45) - feb_tas['P']}", 'Más made → más compras'],
        ['Compras Aliado/mes', str(feb_ali['P']), str(int(feb_ali['Q'] * 0.30 * 0.85 * 0.40)),
         f"+{int(feb_ali['Q'] * 0.30 * 0.85 * 0.40) - feb_ali['P']}", 'Más made → más compras'],
        ['SLA Q→Insp (Ali)', '4+ días', '2-3 días', '-50%', 'Mejor CX, más cierre'],
    ]

    make_table(slide, 0.5, 3.6, 12.3, 2.6, opp_data)

    # Bottom summary
    target_p_tas = int(feb_tas['Q'] * 2.5 * 0.35 * 0.85 * 0.45)
    target_p_ali = int(feb_ali['Q'] * 0.30 * 0.85 * 0.40)
    total_target = target_p_tas + target_p_ali
    current_total = feb_tas['P'] + feb_ali['P']

    add_insight_box(slide, 0.5, 6.4, 6, 0.65, "UPSIDE ESTIMADO",
                   f"De {current_total} compras/mes → {total_target} compras/mes (+{total_target - current_total} incremental)\n"
                   f"Crecimiento: +{((total_target/current_total)-1)*100:.0f}% en compras netas",
                   HIGH_BG, HIGH_BAR)

    add_insight_box(slide, 6.8, 6.4, 5.7, 0.65, "MODELO OPERATIVO",
                   f"Inspector 1: Dedicado TAS (6-8 insp/día)\n"
                   f"Inspector 2: Dedicado Aliado (6-8 insp/día) | Pipe aliado: seguimiento central",
                   SUM_BG, SUM_BAR)

    add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: ISMO LEON - SPECIFIC ISSUE: TAS VOLUME COLLAPSE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_bar(slide, BRAND_BLUE, 0.06)

add_textbox(slide, 0.5, 0.2, 12, 0.5, "ISMO León — TAS Quote Volume: Señal de Alerta", 24, True, DARK_TEXT)
add_textbox(slide, 0.5, 0.65, 12, 0.3, "El canal TAS de León tiene un volumen de cotizaciones críticamente bajo", 11, False, RED)

# TAS Volume evolution
ismo_data = df[df['Célula'].isin(['ISMO LEON', 'ALIADO LEON'])].copy()

ismo_vol = [['Mes', 'Quotes TAS', 'Made TAS', 'Q→M%', 'Quotes ALI', 'Made ALI', 'Q→M%', 'Ratio ALI/TAS']]
for mname, yr, mn in months:
    ftas = get_funnel(ismo_data, yr, mn, 'TAS')
    fali = get_funnel(ismo_data, yr, mn, 'ALIADO')
    ratio = f"{fali['Q']/ftas['Q']:.0f}x" if ftas['Q'] > 0 else "-"
    ismo_vol.append([mname, ftas['Q'], ftas['M'], f"{ftas['Q2M']:.0f}%",
                     fali['Q'], fali['M'], f"{fali['Q2M']:.0f}%", ratio])

make_table(slide, 0.5, 1.1, 12.3, 2.2, ismo_vol)

add_insight_box(slide, 0.5, 3.5, 6, 1.5, "DIAGNÓSTICO",
    "1. TAS solo genera 17-23 quotes/mes vs 400-485 de Aliado\n"
    "2. Ratio Aliado/TAS es 20-28x — el TAS es marginal en volumen\n"
    "3. El LM dedica ~90% de su tiempo al pipe de aliados\n"
    "4. Sin marketing activo para TAS (se pausó por capacity)\n"
    "5. El inspector hace 91% inspecciones de aliado, 9% TAS",
    LOW_BG, LOW_BAR)

add_insight_box(slide, 6.8, 3.5, 5.7, 1.5, "PLAN DE ACCIÓN",
    "1. LM dedicado TAS: generar demanda propia + walk-ins\n"
    "2. Reactivar inversión marketing para León (TAS)\n"
    "3. Inspector dedicado TAS: reducir SLA de 2 a 1 día\n"
    "4. Inspector dedicado Aliado: absorber volumen actual\n"
    "5. Seguimiento pipe Aliado → equipo central (libera LM)",
    HIGH_BG, HIGH_BAR)

# Target for TAS Growth in Leon
add_textbox(slide, 0.5, 5.2, 12, 0.35, "Ramp-Up Proyectado — TAS León (con LM + Inspector dedicados + Mktg)", 14, True, BRAND_BLUE)

ramp_data = [
    ['', 'Actual (Feb)', 'Mes 1', 'Mes 2', 'Mes 3', 'Mes 4', 'Mes 6'],
    ['Quotes TAS', '17', '40', '60', '80', '100', '120+'],
    ['Made TAS', '10', '20', '30', '35', '40', '50+'],
    ['Q→M% TAS', '59%', '50%', '50%', '44%', '40%', '42%'],
    ['Compras TAS', '5', '8', '12', '15', '18', '22+'],
    ['Quotes ALI (ref)', '485', '450', '450', '450', '450', '450'],
    ['Made ALI', '105', '110', '115', '120', '125', '130'],
    ['Total Compras', '25', '30', '38', '43', '50', '58+'],
]
make_table(slide, 0.5, 5.55, 12.3, 1.6, ramp_data)

add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: KIA COAPA - SPECIFIC: BOTH CHANNELS DEMANDING
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_bar(slide, BRAND_BLUE, 0.06)

add_textbox(slide, 0.5, 0.2, 12, 0.5, "Kia Coapa — Presión en Ambos Canales", 24, True, DARK_TEXT)
add_textbox(slide, 0.5, 0.65, 12, 0.3, "Volumen competitivo TAS + Aliado con inspector compartido genera cuellos de botella", 11, False, RED)

coapa_data = df[df['Célula'].isin(['CONTINENTAL COAPA', 'ALIADO COAPA'])].copy()

coapa_vol = [['Mes', 'Q TAS', 'M TAS', 'Q→M%', 'Q ALI', 'M ALI', 'Q→M%', 'M Total', 'Insp/día']]
for mname, yr, mn in months:
    ftas = get_funnel(coapa_data, yr, mn, 'TAS')
    fali = get_funnel(coapa_data, yr, mn, 'ALIADO')
    ft = get_funnel(coapa_data, yr, mn)
    daily = get_daily_inspections(coapa_data, yr, mn)
    avg_d = f"{daily.mean():.1f}" if len(daily) > 0 else "-"
    coapa_vol.append([mname, ftas['Q'], ftas['M'], f"{ftas['Q2M']:.0f}%",
                      fali['Q'], fali['M'], f"{fali['Q2M']:.0f}%", ft['M'], avg_d])

make_table(slide, 0.5, 1.1, 12.3, 2.2, coapa_vol)

# Key insight: Jan was problematic
jan_tas = get_funnel(coapa_data, 2026, 1, 'TAS')
feb_tas = get_funnel(coapa_data, 2026, 2, 'TAS')

add_insight_box(slide, 0.5, 3.5, 6, 1.3, "DIAGNÓSTICO",
    f"1. TAS Q→M cayó a {jan_tas['Q2M']:.0f}% en Ene (capacity constraint)\n"
    f"2. En Feb mejoró a {feb_tas['Q2M']:.0f}% pero con apoyo de inspectores extras\n"
    "3. Jonathan (inspector principal) hizo 86/132 insp en Feb\n"
    "4. SLA Aliado: 4.8 días promedio → mala experiencia\n"
    "5. LM atiende 500+ quotes combinados → seguimiento débil",
    LOW_BG, LOW_BAR)

add_insight_box(slide, 6.8, 3.5, 5.7, 1.3, "BENEFICIO DE 2I + 2LM",
    "1. Inspector TAS dedicado → Q→M sostenido en 35-40%\n"
    "2. Inspector Aliado dedicado → SLA de 4.8 a 2-3 días\n"
    "3. LM TAS: foco en conversión y crecimiento de pipe\n"
    "4. LM Aliado: coordinación en piso + seguimiento\n"
    "5. Pipe de aliado residual → seguimiento central",
    HIGH_BG, HIGH_BAR)

# Ramp-up Coapa
add_textbox(slide, 0.5, 5.0, 12, 0.35, "Ramp-Up Proyectado — Kia Coapa (con 2 Inspectores + 2 LMs)", 14, True, BRAND_BLUE)

feb_tas_c = get_funnel(coapa_data, 2026, 2, 'TAS')
feb_ali_c = get_funnel(coapa_data, 2026, 2, 'ALIADO')

ramp_coapa = [
    ['', 'Actual (Feb)', 'Mes 1', 'Mes 2', 'Mes 3', 'Mes 4', 'Mes 6'],
    ['Quotes TAS', str(feb_tas_c['Q']), '200', '230', '260', '280', '300+'],
    ['Made TAS', str(feb_tas_c['M']), '70', '85', '95', '105', '115+'],
    ['Q→M% TAS', f"{feb_tas_c['Q2M']:.0f}%", '35%', '37%', '37%', '38%', '38%'],
    ['Compras TAS', str(feb_tas_c['P']), '25', '30', '35', '38', '42+'],
    ['Made ALI', str(feb_ali_c['M']), '90', '95', '100', '105', '110+'],
    ['Compras ALI', str(feb_ali_c['P']), '35', '38', '42', '45', '48+'],
    ['Total Compras', str(feb_tas_c['P'] + feb_ali_c['P']), '60', '68', '77', '83', '90+'],
]
make_table(slide, 0.5, 5.35, 12.3, 1.65, ramp_coapa)

add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: COMPARATIVE ROI / BUSINESS CASE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_bar(slide, BRAND_BLUE, 0.06)

add_textbox(slide, 0.5, 0.2, 12, 0.5, "Business Case — ROI de la Inversión en Headcount", 24, True, DARK_TEXT)

# Cost assumptions
add_textbox(slide, 0.5, 0.75, 12, 0.3, "Supuestos de costo y revenue", 14, True, DARK_TEXT)

assumptions = [
    ['Concepto', 'ISMO León', 'Kia Coapa', 'Nota'],
    ['Costo Inspector/mes', '$25,000', '$25,000', 'Salario + comisiones promedio'],
    ['Costo Lead Manager/mes', '$20,000', '$20,000', 'Salario + comisiones promedio'],
    ['Inversión mensual (+1I +1LM)', '$45,000', '$45,000', 'Incremental al headcount actual'],
    ['Avg Oferta Final', '$293K', '$201K', 'Precio promedio de compra Feb'],
    ['Compras actuales/mes', '25', '49', 'Feb 2026'],
    ['Compras target Mes 3', '43', '77', 'Con 2I + 2LM'],
    ['Compras incrementales', '+18', '+28', 'Compras adicionales/mes'],
    ['Fee por compra (~3-5%)', '~$10K', '~$7K', 'Revenue incremental por unidad'],
    ['Revenue incremental/mes', '~$180K', '~$196K', '+18 x $10K / +28 x $7K'],
    ['ROI', '4.0x', '4.4x', 'Revenue incr / Costo incr'],
]

make_table(slide, 0.5, 1.1, 12.3, 3.2, assumptions)

add_insight_box(slide, 0.5, 4.5, 6, 1.1, "ISMO LEÓN — PAYBACK",
    "Inversión: $45K/mes incremental\n"
    "Revenue incremental: ~$180K/mes (Mes 3)\n"
    "ROI: 4.0x — se paga en primera semana\n"
    "Breakeven: +5 compras incrementales/mes",
    HIGH_BG, HIGH_BAR)

add_insight_box(slide, 6.8, 4.5, 5.7, 1.1, "KIA COAPA — PAYBACK",
    "Inversión: $45K/mes incremental\n"
    "Revenue incremental: ~$196K/mes (Mes 3)\n"
    "ROI: 4.4x — se paga en primera semana\n"
    "Breakeven: +7 compras incrementales/mes",
    HIGH_BG, HIGH_BAR)

add_insight_box(slide, 0.5, 5.8, 12, 0.7, "RESUMEN DE INVERSIÓN TOTAL",
    "Total: +2 inspectores + 2 LMs (2 células) = $90K/mes inversión incremental\n"
    "Revenue incremental combinado: ~$376K/mes (Mes 3) | ROI combinado: 4.2x\n"
    "Riesgo bajo: Si se logra solo 50% del target, ROI sigue siendo 2.1x — ampliamente positivo",
    SUM_BG, SUM_BAR)

add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10: TIMELINE & NEXT STEPS
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_bar(slide, BRAND_BLUE, 0.06)

add_textbox(slide, 0.5, 0.2, 12, 0.5, "Plan de Implementación & Next Steps", 24, True, DARK_TEXT)

timeline = [
    ['Fase', 'Timeline', 'Acción', 'Owner', 'KPI'],
    ['1. Aprobación', 'Semana 1 Mar', 'Aprobación de headcount por Central', 'Daniel + Central', 'Headcount aprobado'],
    ['2. Recruiting', 'Semana 2-3 Mar', 'Publicar vacantes + entrevistas', 'Ops + HR', '4 candidatos seleccionados'],
    ['3. Onboarding', 'Semana 4 Mar - Sem 1 Abr', 'Training de inspectores y LMs', 'Daniel + Cell Leads', 'Training completado'],
    ['4. Ramp-Up', 'Abril 2026', 'Operación con 2I+2LM en cada célula', 'Cell Leads', 'Made +30% vs Feb'],
    ['5. Estabilización', 'Mayo 2026', 'Optimización y ajuste de targets', 'Daniel', 'Compras +50% vs Feb'],
    ['6. Full Run Rate', 'Junio 2026', 'Operación plena con mktg activo', 'Daniel + Mktg', 'Compras +80% vs Feb'],
]

make_table(slide, 0.5, 0.85, 12.3, 2.5, timeline)

# Migration plan for Aliado pipe
add_textbox(slide, 0.5, 3.6, 12, 0.35, "Migración de Seguimiento Aliado a Central", 16, True, BRAND_BLUE)

migration = [
    ['Etapa', 'Qué migra', 'Qué queda en célula', 'Timeline'],
    ['Fase 1', 'Follow-up telefónico de pipe Aliado', 'Agendamiento en piso + inspección', 'Mes 1'],
    ['Fase 2', '+ Cotización remota de Aliado', 'Inspección en agencia + cierre', 'Mes 2'],
    ['Fase 3', '+ Seguimiento post-oferta Aliado', 'Solo inspección + relación agencia', 'Mes 3'],
]
make_table(slide, 0.5, 4.0, 12.3, 1.3, migration)

# KPIs to monitor
add_textbox(slide, 0.5, 5.6, 12, 0.35, "KPIs de Monitoreo", 16, True, BRAND_BLUE)

kpi_data = [
    ['KPI', 'ISMO León Target', 'Kia Coapa Target', 'Frecuencia'],
    ['Quotes TAS/mes', '80+ (Mes 3)', '260+ (Mes 3)', 'Semanal'],
    ['Q→M% TAS', '≥40%', '≥35%', 'Semanal'],
    ['Inspecciones/día', '12-16 (combinado)', '12-16 (combinado)', 'Diario'],
    ['SLA Q→Inspección', '≤2 días TAS / ≤3 días ALI', '≤2 días TAS / ≤3 días ALI', 'Semanal'],
    ['Compras netas/mes', '43+ (Mes 3)', '77+ (Mes 3)', 'Mensual'],
    ['Utilización inspector', '≥75% (6+ insp/día)', '≥75% (6+ insp/día)', 'Semanal'],
]
make_table(slide, 0.5, 5.95, 12.3, 1.35, kpi_data)

add_footer(slide)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11: CLOSING
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_top_bar(slide, BRAND_BLUE, 0.12)

sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06))
sep.fill.solid()
sep.fill.fore_color.rgb = BRAND_BLUE
sep.line.fill.background()

add_textbox(slide, 1.0, 1.5, 11, 0.5, "Resumen Ejecutivo", 14, True, BRAND_BLUE)
add_textbox(slide, 1.0, 2.0, 11, 0.8, "2 Inspectores + 2 Lead Managers", 36, True, WHITE)

add_textbox(slide, 1.0, 3.3, 11, 0.4, "para ISMO León y Kia Coapa", 22, False, WHITE)

# Summary boxes on dark bg
summary_items = [
    ("$90K/mes", "Inversión incremental total", AMBER),
    ("$376K/mes", "Revenue incremental (Mes 3)", GREEN),
    ("4.2x ROI", "Retorno sobre inversión", GREEN),
    ("+46 compras", "Incrementales por mes (Mes 3)", BRAND_BLUE),
]

for i, (val, desc, color) in enumerate(summary_items):
    x = 1.0 + i * 3.0
    add_textbox(slide, x, 4.2, 2.8, 0.5, val, 32, True, color)
    add_textbox(slide, x, 4.7, 2.8, 0.3, desc, 11, False, LIGHT_GRAY)

add_textbox(slide, 1.0, 5.5, 11, 0.4, "Seguimiento de pipe aliado migra a equipo central — sin costo adicional", 13, False, WHITE)
add_textbox(slide, 1.0, 6.0, 11, 0.3, "Next step: Aprobación de headcount — Semana 1 de Marzo 2026", 14, True, AMBER)

add_textbox(slide, 9.5, 7.0, 3.5, 0.3, "Kavak Marketplace B2B — Programa TAS", 11, False, LIGHT_GRAY, PP_ALIGN.RIGHT)

# ── Save ────────────────────────────────────────────────────────────────
output_path = os.path.expanduser("~/Downloads/SOp_Resource_Request_ISMO_Leon_Kia_Coapa_Mar2026.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
