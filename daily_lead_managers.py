from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Kavak-inspired colors
KAVAK_BLUE = RGBColor(0x00, 0x2D, 0x72)
KAVAK_LIGHT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
KAVAK_ORANGE = RGBColor(0xFF, 0x6B, 0x00)
KAVAK_GREEN = RGBColor(0x00, 0xB3, 0x4D)
KAVAK_YELLOW = RGBColor(0xFF, 0xC1, 0x07)
KAVAK_RED = RGBColor(0xE8, 0x3E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_BG = RGBColor(0x00, 0x1A, 0x40)
SECTION_BG = RGBColor(0xE8, 0xEE, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=BLACK, spacing=Pt(6), bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if bold_first and ": " in item:
            # We can't do partial bold easily, so just bold the whole line if it has a colon
            pass
    return txBox

def add_table_slide(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=KAVAK_BLUE, font_size=11):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx] if row_idx < len(data) and col_idx < len(data[row_idx]) else ""

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = BLACK

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape

# ============================================================
# SLIDE 1 - PORTADA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Accent bar top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), KAVAK_ORANGE)

# Title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
    "Daily de Capacitación", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
    "Lead Managers Rookie — Programa de Onboarding", font_size=24, color=KAVAK_LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04), KAVAK_ORANGE)

# Subtitle info
add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(0.5),
    "Equipo Central · Silvia + Soporte", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
    "Marzo 2026", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.08), KAVAK_ORANGE)

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
    "Objetivo: Que cada Lead Manager gestione compras correctamente desde el primer intento",
    font_size=14, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 - CONTEXTO Y OBJETIVO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Contexto y Objetivo", font_size=30, bold=True, color=KAVAK_BLUE)

# Left card - Contexto
card1 = add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), SECTION_BG, KAVAK_BLUE, Pt(1.5))
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(5), Inches(0.5),
    "Contexto", font_size=20, bold=True, color=KAVAK_BLUE)

context_items = [
    "• 6 Lead Managers nivel Rookie en 7 células",
    "• Todos en etapa temprana de operación",
    "• Equipo central (Silvia) brindará soporte",
    "• Células distribuidas: Soni Pachuca, Dalton,",
    "  ISMO, Plasencia, Potosina, We Cars, MISOL",
    "",
    "• Necesidad crítica: estandarizar el proceso",
    "  de compra desde el onboarding"
]
add_bullet_list(slide, Inches(1.2), Inches(2.1), Inches(5), Inches(4), context_items, font_size=14, color=BLACK)

# Right card - Objetivo
card2 = add_shape(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), SECTION_BG, KAVAK_ORANGE, Pt(1.5))
add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
    "Objetivo del Daily", font_size=20, bold=True, color=KAVAK_ORANGE)

objective_items = [
    "• Gestión de compras correcta al 1er intento",
    "",
    "  Los 3 pilares de aprendizaje:",
    "",
    "  1. Cotejar documentos recibidos de la",
    "     agencia (factura, TC, INE, tenencias)",
    "",
    "  2. Crear ticket de aprobación de compra",
    "     siguiendo procesos establecidos",
    "",
    "  3. Solicitar autorización de autofactura",
    "     al SAT correctamente"
]
add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(4), objective_items, font_size=14, color=BLACK)


# ============================================================
# SLIDE 3 - MAPA DE LEAD MANAGERS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Mapa de Lead Managers", font_size=30, bold=True, color=KAVAK_BLUE)

data = [
    ["Grupo", "Lead Manager", "Nivel", "Necesidad Principal"],
    ["Soni Pachuca", "Jorge Eduardo", "🟡 Rookie", "Data fin de turno + 1er ticket de compra + Cotejo docs"],
    ["Dalton", "Ana", "🟡 Rookie", "TID desde central + Gestión leads Plasencia + Inspector"],
    ["ISMO", "(Por asignar)", "🟡 Rookie", "Aliados inicia 10 marzo"],
    ["Plasencia", "Cristian", "🟡 Rookie", "Primera semana en piso — arranca desde cero"],
    ["Potosina", "Eduardo", "🟡 Rookie", "Aliados inicia 10 marzo"],
    ["We Cars", "Ana", "🟡 Rookie", "Por definir en primer daily"],
    ["MISOL", "Eduardo", "🟡 Rookie", "Por definir en primer daily"],
]

col_widths = [Inches(1.8), Inches(2.0), Inches(1.5), Inches(7.0)]
add_table_slide(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0), 8, 4, data, col_widths, font_size=12)


# ============================================================
# SLIDE 4 - ESTRUCTURA DEL DAILY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Estructura del Daily (30 min)", font_size=30, bold=True, color=KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.4),
    "Horario: 9:00 AM  ·  Formato: Videollamada con cámara  ·  Facilitador: Silvia",
    font_size=14, color=GRAY)

# 4 blocks
blocks = [
    ("Check-in", "5 min", "Cada LM reporta:\n• Tickets creados ayer\n• Dudas bloqueantes\n• Estatus general", KAVAK_BLUE),
    ("Micro-cápsula", "10 min", "Tema del día (rotativo):\n• Cotejo de documentos\n• Tickets de compra\n• Autofactura SAT", KAVAK_ORANGE),
    ("Caso práctico", "10 min", "Revisión en vivo de:\n• 1 ticket real (bueno o\n  con errores)\n• Feedback grupal", KAVAK_GREEN),
    ("Cierre", "5 min", "• Compromiso del día\n• Tarea práctica asignada\n• Dudas finales", KAVAK_LIGHT_BLUE),
]

for i, (title, duration, content, color) in enumerate(blocks):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.6)

    # Card background
    card = add_shape(slide, x, y, Inches(2.95), Inches(5.2), WHITE, color, Pt(2))

    # Color header bar
    add_shape(slide, x, y, Inches(2.95), Inches(0.9), color, color, Pt(0))

    # Title in header
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.6), Inches(0.5),
        title, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Duration badge
    add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.6), Inches(0.4),
        duration, font_size=14, bold=True, color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)

    # Content
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.5), Inches(3.5),
        content, font_size=13, color=BLACK)

    # Step number
    circle = add_shape(slide, x + Inches(1.1), y - Inches(0.25), Inches(0.5), Inches(0.5), color, color, Pt(0), MSO_SHAPE.OVAL)
    add_text_box(slide, x + Inches(1.1), y - Inches(0.2), Inches(0.5), Inches(0.5),
        str(i+1), font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 - PLAN DE MICRO-CÁPSULAS SEMANAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Plan de Micro-cápsulas (Ciclo Semanal)", font_size=30, bold=True, color=KAVAK_BLUE)

capsule_data = [
    ["Día", "Tema", "Objetivo", "Entregable"],
    ["Lunes", "Cotejo de documentos", "Checklist completo: factura, TC, INE,\ntenencias, verificaciones, red flags", "Checklist firmado"],
    ["Martes", "Creación de ticket de compra", "Paso a paso: campos obligatorios,\nevidencia, pricing, aprobaciones", "Ticket de práctica"],
    ["Miércoles", "Autofactura SAT", "Proceso: cuándo aplica, cómo solicitar,\ntiempos, qué hacer si se rechaza", "Solicitud simulada"],
    ["Jueves", "Gestión de pagos y aging", "Primer/segundo pago, autos estancados,\nescalamiento", "Reporte de aging"],
    ["Viernes", "Tableros + Cierre semanal", "Lectura de dashboards, KPIs de la\nsemana, reconocimiento mejor caso", "Scorecard semanal"],
]

col_widths = [Inches(1.5), Inches(3.0), Inches(4.8), Inches(2.7)]
add_table_slide(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), 6, 4, capsule_data, col_widths, font_size=12)


# ============================================================
# SLIDE 6 - RELACIÓN CON EQUIPO DE SEMINUEVOS (NEW SECTION)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
    "Relación con Equipo de Seminuevos", font_size=30, bold=True, color=KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
    "Espacio dentro del daily para que cada LM comparta contexto de su interacción con el equipo de seminuevos de cada célula",
    font_size=14, color=GRAY)

# Card: Qué se reporta
card_left = add_shape(slide, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.4), SECTION_BG, KAVAK_BLUE, Pt(1.5))
add_text_box(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.5),
    "Qué reporta cada LM (2 min por persona)", font_size=18, bold=True, color=KAVAK_BLUE)

report_items = [
    "1. Nivel de colaboración con Gte. de Seminuevos",
    "   → ¿Hay fluidez en la entrega de cotizaciones?",
    "   → ¿Entregan documentación completa o incompleta?",
    "",
    "2. Volumen y calidad de leads recibidos",
    "   → Cotizaciones con docs completos vs. incompletos",
    "   → CFDI pendientes: cuántos y desde cuándo",
    "",
    "3. Puntos de fricción o bloqueos",
    "   → Tiempos de respuesta del equipo de seminuevos",
    "   → Documentos recurrentemente faltantes",
    "   → Escalamientos necesarios",
    "",
    "4. Alianzas y avances",
    "   → Acuerdos de trabajo establecidos",
    "   → Mejoras observadas semana a semana",
]
add_bullet_list(slide, Inches(0.9), Inches(2.4), Inches(5.5), Inches(4.5), report_items, font_size=12, color=BLACK)

# Card: Template
card_right = add_shape(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.4), SECTION_BG, KAVAK_ORANGE, Pt(1.5))
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5.2), Inches(0.5),
    "Template de Reporte Semanal", font_size=18, bold=True, color=KAVAK_ORANGE)

template_items = [
    "📋 Formato que llena cada LM el viernes:",
    "",
    "Célula: ________________",
    "Gte. Seminuevos: ________________",
    "",
    "Cotizaciones recibidas esta semana: ___",
    "  - Con docs completos: ___",
    "  - Con CFDI pendiente: ___",
    "  - Sin documentación: ___",
    "",
    "Inspecciones realizadas: ___",
    "  - Docs completos: ___",
    "  - CFDI pendiente únicamente: ___",
    "",
    "Nivel de relación (1-5): ___",
    "Principal fricción: ________________",
    "Acción acordada: ________________",
]
add_bullet_list(slide, Inches(7.3), Inches(2.4), Inches(5.2), Inches(4.5), template_items, font_size=12, color=BLACK)


# ============================================================
# SLIDE 7 - DIAGNÓSTICO RELACIÓN SEMINUEVOS POR CÉLULA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
    "Diagnóstico Inicial: Relación con Seminuevos por Célula", font_size=28, bold=True, color=KAVAK_BLUE)

diag_data = [
    ["Célula", "Lead Manager", "Gte. Seminuevos", "Estatus Relación", "Acción Prioritaria"],
    ["Soni Pachuca", "Jorge Eduardo", "Lluvia & Ana", "🟡 En construcción", "Establecer reporte diario de\ncotizaciones e inspecciones"],
    ["Dalton", "Ana", "(Por confirmar)", "🟡 En construcción", "Definir flujo de entrega de\ndocumentos con equipo local"],
    ["ISMO Ags.", "(Por asignar)", "(Pre-arranque)", "⚪ No iniciada", "Mapear equipo de seminuevos\nantes del 10 de marzo"],
    ["Plasencia", "Cristian", "(Por confirmar)", "🔴 Desde cero", "Primera semana: presentación\nformal + acuerdos de trabajo"],
    ["Potosina", "Eduardo", "(Pre-arranque)", "⚪ No iniciada", "Mapear equipo de seminuevos\nantes del 10 de marzo"],
    ["We Cars", "Ana", "(Por confirmar)", "🟡 Por evaluar", "Diagnóstico en primer daily"],
    ["MISOL", "Eduardo", "(Por confirmar)", "🟡 Por evaluar", "Diagnóstico en primer daily"],
]

col_widths = [Inches(1.6), Inches(1.8), Inches(2.0), Inches(2.2), Inches(4.4)]
add_table_slide(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), 8, 5, diag_data, col_widths, font_size=11)

add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
    "Nota: El diagnóstico se actualiza semanalmente con el input del daily. Semáforo: 🔴 Crítico  🟡 En desarrollo  🟢 Estable  ⚪ No iniciada",
    font_size=11, color=GRAY)


# ============================================================
# SLIDE 8 - PLAN ESPECÍFICO POR LM (Semanas 1-2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Plan Específico por LM — Semanas 1-2", font_size=30, bold=True, color=KAVAK_BLUE)

# Cards for each LM
lm_plans = [
    ("Jorge Eduardo", "Soni Pachuca", KAVAK_BLUE, [
        "• Data fin de turno: template de reporte",
        "• Volumen cotizaciones (Lluvia & Ana)",
        "• Inspecciones con docs completos",
        "• Primer ticket supervisado por Silvia",
        "• Foco: experto en cotejo de docs",
    ]),
    ("Ana", "Dalton", KAVAK_ORANGE, [
        "• TID desde central",
        "• Gestión leads Plasencia",
        "• Priorizar asignación inspector",
        "• Meta bono: $3,500 con 2do pago",
        "• Calmar ansiedad de esquema",
    ]),
    ("Cristian", "Plasencia", KAVAK_RED, [
        "• Arranca desde cero en piso",
        "• Días 1-3: observación + cotejo",
        "• Días 4-5: primer ticket supervisado",
        "• Presentación con equipo seminuevos",
        "• Acompañamiento intensivo Silvia",
    ]),
    ("ISMO / Potosina", "Ags. + Potosina", KAVAK_GREEN, [
        "• Pre-arranque: cápsulas grabadas",
        "• Aliados inicia 10 de marzo",
        "• Semana previa: cotejo + tickets",
        "• Acompañamiento intensivo S1",
        "• Mapear equipo seminuevos local",
    ]),
    ("Ana / Eduardo", "We Cars + MISOL", KAVAK_LIGHT_BLUE, [
        "• Diagnosticar en primer daily",
        "• Integrar al ciclo estándar",
        "• Evaluar relación con seminuevos",
        "• Definir plan personalizado S2",
        "• Asignar buddy del equipo central",
    ]),
]

for i, (name, cell, color, items) in enumerate(lm_plans):
    col = i % 5
    x = Inches(0.3 + col * 2.6)
    y = Inches(1.2)

    card = add_shape(slide, x, y, Inches(2.45), Inches(5.5), WHITE, color, Pt(2))

    # Header
    add_shape(slide, x, y, Inches(2.45), Inches(1.0), color, color, Pt(0))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.2), Inches(0.45),
        name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.2), Inches(0.4),
        cell, font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.1), y + Inches(1.2 + j * 0.7), Inches(2.2), Inches(0.65),
            item, font_size=10, color=BLACK)


# ============================================================
# SLIDE 9 - MATERIALES DE APOYO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Materiales de Apoyo a Crear", font_size=30, bold=True, color=KAVAK_BLUE)

materials = [
    ("1", "Checklist de Cotejo\nde Documentos", "1 pager imprimible", [
        "• Lista de docs por tipo de operación",
        "• Red flags visuales por documento",
        "• Flujo: doc completo vs incompleto",
        "• Campos obligatorios a validar",
    ], KAVAK_BLUE),
    ("2", "Guía Paso a Paso:\nTicket de Compra", "Con capturas de pantalla", [
        "• Cada campo explicado",
        "• Errores más comunes",
        "• Evidencia a adjuntar",
        "• Checklist pre-envío",
    ], KAVAK_ORANGE),
    ("3", "Proceso de\nAutofactura SAT", "Diagrama de flujo", [
        "• Cuándo aplica autofactura",
        "• Portal/herramienta a usar",
        "• Tiempos de respuesta",
        "• Escalamiento si se rechaza",
    ], KAVAK_GREEN),
    ("4", "Template Reporte\nSeminuevos", "Google Sheet compartido", [
        "• Cotizaciones recibidas",
        "• Docs completos vs pendientes",
        "• Nivel de relación (semáforo)",
        "• Fricciones y acuerdos",
    ], KAVAK_RED),
]

for i, (num, title, subtitle, items, color) in enumerate(materials):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.2)

    card = add_shape(slide, x, y, Inches(3.0), Inches(5.8), WHITE, color, Pt(2))

    # Number circle
    circle = add_shape(slide, x + Inches(1.15), y + Inches(0.3), Inches(0.7), Inches(0.7), color, color, Pt(0), MSO_SHAPE.OVAL)
    add_text_box(slide, x + Inches(1.15), y + Inches(0.35), Inches(0.7), Inches(0.6),
        num, font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.6), Inches(0.8),
        title, font_size=16, bold=True, color=KAVAK_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(2.0), Inches(2.6), Inches(0.4),
        subtitle, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), y + Inches(2.6 + j * 0.65), Inches(2.4), Inches(0.6),
            item, font_size=12, color=BLACK)


# ============================================================
# SLIDE 10 - MÉTRICAS DE ÉXITO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Métricas de Éxito del Programa", font_size=30, bold=True, color=KAVAK_BLUE)

metrics_data = [
    ["Métrica", "Meta Semana 2", "Meta Semana 4", "Cómo se mide"],
    ["% Tickets correctos al 1er intento", "60%", "90%", "Revisión por Silvia en daily"],
    ["Tiempo promedio cotejo de docs", "< 20 min", "< 10 min", "Autoreporte + observación"],
    ["Autofacturas solicitadas correctamente", "70%", "95%", "Tracking en sistema"],
    ["Asistencia al daily", "100%", "100%", "Registro de asistencia"],
    ["Nivel relación seminuevos (1-5)", "3.0", "4.0", "Template semanal"],
    ["Docs entregados completos por seminuevos", "50%", "75%", "Reporte LM en daily"],
]

col_widths = [Inches(3.5), Inches(2.2), Inches(2.2), Inches(4.0)]
add_table_slide(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5), 7, 4, metrics_data, col_widths, font_size=12)

# Graduation criteria
add_shape(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), SECTION_BG, KAVAK_GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(6.15), Inches(11), Inches(0.7),
    "🎯 Meta de graduación Rookie → Operativo:  5 tickets consecutivos correctos al primer intento  +  Relación estable con equipo seminuevos (nivel ≥ 4)",
    font_size=15, bold=True, color=KAVAK_BLUE)


# ============================================================
# SLIDE 11 - CANALES DE COMUNICACIÓN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
    "Canales de Comunicación y Soporte", font_size=30, bold=True, color=KAVAK_BLUE)

channels = [
    ("Daily 9:00 AM", "Videollamada diaria\n30 min máximo\nTodos los LMs + Silvia\nCámara encendida obligatoria", KAVAK_BLUE, "Espacio principal"),
    ("Canal WhatsApp/Slack", "Exclusivo para LMs Rookie\nDudas de cotejo con fotos\nValidación pre-envío de tickets\nRespuesta < 30 min por Silvia", KAVAK_ORANGE, "Soporte en tiempo real"),
    ("Reporte Semanal", "Template Google Sheets\nCada viernes antes del daily\nMétricas + relación seminuevos\nBase para coaching 1:1", KAVAK_GREEN, "Seguimiento formal"),
    ("Sesión 1:1 Quincenal", "Silvia + cada LM individual\n30 min cada quincena\nRevisión de avance personal\nPlan de acción personalizado", KAVAK_LIGHT_BLUE, "Desarrollo individual"),
]

for i, (title, content, color, tag) in enumerate(channels):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.2)

    card = add_shape(slide, x, y, Inches(3.0), Inches(5.5), WHITE, color, Pt(2))

    # Header
    add_shape(slide, x, y, Inches(3.0), Inches(0.8), color, color, Pt(0))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.8), Inches(0.7),
        title, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Tag
    add_text_box(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.8), Inches(0.4),
        tag, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)

    # Content
    add_text_box(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.6), Inches(3.5),
        content, font_size=13, color=BLACK)


# ============================================================
# SLIDE 12 - PRÓXIMOS PASOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KAVAK_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "Próximos Pasos", font_size=30, bold=True, color=WHITE)

steps = [
    ("Semana 1 (3-7 Mar)", "Arranque del daily · Jorge Eduardo y Ana primer ticket · Cristian observación + cotejo · Diagnóstico relación seminuevos en cada célula"),
    ("Semana 2 (10-14 Mar)", "ISMO y Potosina arrancan con Aliados · Cristian primer ticket supervisado · Todos: ciclo completo de micro-cápsulas · Primer reporte semanal de seminuevos"),
    ("Semana 3-4", "Medición de métricas S2 · Ajuste de plan por LM · Primera sesión 1:1 quincenal · Revisión de relación seminuevos (meta: nivel 3+)"),
    ("Mes 2", "Graduación de primeros LMs · Transición a daily bisemanal para graduados · Foco en LMs rezagados · Meta: 90% tickets correctos al 1er intento"),
]

for i, (week, content) in enumerate(steps):
    y = Inches(1.3 + i * 1.45)

    # Timeline dot
    add_shape(slide, Inches(1.8), y + Inches(0.15), Inches(0.3), Inches(0.3), KAVAK_ORANGE, KAVAK_ORANGE, Pt(0), MSO_SHAPE.OVAL)

    # Line connector
    if i < len(steps) - 1:
        add_shape(slide, Inches(1.9), y + Inches(0.45), Inches(0.06), Inches(1.0), KAVAK_ORANGE)

    # Week label
    add_text_box(slide, Inches(2.3), y, Inches(3), Inches(0.4),
        week, font_size=16, bold=True, color=KAVAK_ORANGE)

    # Content
    add_text_box(slide, Inches(2.3), y + Inches(0.45), Inches(9.5), Inches(0.9),
        content, font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))


# Bottom message
add_text_box(slide, Inches(2), Inches(6.8), Inches(9), Inches(0.5),
    "El éxito de este programa se mide en la autonomía de cada Lead Manager",
    font_size=16, bold=True, color=KAVAK_ORANGE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.08), KAVAK_ORANGE)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.expanduser("~/Documents/claude-code/Daily_Capacitacion_Lead_Managers_Rookie.pptx")
prs.save(output_path)
print(f"Presentación guardada en: {output_path}")
