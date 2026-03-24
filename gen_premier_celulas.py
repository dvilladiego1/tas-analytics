#!/usr/bin/env python3
"""
Grupo Premier — Funnel por Célula (3 slides: Culiacán, Hermosillo, Mazatlán)
Cada slide muestra el breakdown por marca dentro de esa célula.
"""
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
NAVY = RGBColor(0x00, 0x4E, 0x98)
TABLE_HEADER = RGBColor(0x2B, 0x47, 0x8B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GRAY_TEXT = RGBColor(0x7A, 0x7A, 0x7A)
SEC_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
LECTURA_BLUE = RGBColor(0xE3, 0xEE, 0xFB)

# Célula accent colors
CELULA_COLORS = {
    'Culiacán':   RGBColor(0x1B, 0x2A, 0x4A),  # Navy
    'Hermosillo': RGBColor(0x3B, 0x7D, 0xDD),  # Blue
    'Mazatlán':   RGBColor(0x00, 0xB4, 0x8A),  # Teal
}

BRAND_COLORS = {
    'CHEVROLET': RGBColor(0x1B, 0x2A, 0x4A),
    'BYD': RGBColor(0x22, 0x6B, 0x22),
    'TOYOTA': RGBColor(0xCC, 0x00, 0x00),
    'KIA': RGBColor(0x00, 0x4E, 0x98),
    'MERCEDESBENZ': RGBColor(0x33, 0x33, 0x33),
    'HYUNDAI': RGBColor(0x00, 0x2C, 0x5F),
    'BUICK': RGBColor(0x8B, 0x00, 0x00),
    'AUTOCOUNTRY': RGBColor(0x00, 0x80, 0x80),
    'BMW': RGBColor(0x1C, 0x69, 0xD4),
    'GMC': RGBColor(0xC0, 0x00, 0x00),
    'JAC': RGBColor(0x66, 0x66, 0x66),
    'STELLANTIS': RGBColor(0x7B, 0x2D, 0x8E),
    'LANDROVER': RGBColor(0x00, 0x5A, 0x2B),
}

# ── Data Loading ────────────────────────────────────────────────
CSV_PATH = os.path.expanduser("~/Downloads/Summary MKP _ TAS & BULK - Snapshot_Diario (1).csv")
df = pd.read_csv(CSV_PATH, encoding='latin-1')

channel_col = [c for c in df.columns if 'BULK' in c and 'TAS' in c][0]
df.rename(columns={channel_col: 'channel'}, inplace=True)
fecha_col = [c for c in df.columns if 'Fecha de creaci' in c][0]

mask = (df['Grupo'].str.upper().str.contains('PREMIER', na=False)) & (df['channel'] != 'BULK')
prem = df[mask].copy()

prem['creation_date'] = pd.to_datetime(prem[fecha_col], format='%d/%m/%Y', errors='coerce')
prem['purchase_dt'] = pd.to_datetime(prem['purchase_date'], format='%d/%m/%Y', errors='coerce')
prem['creation_month'] = prem['creation_date'].dt.to_period('M')
prem['purchase_month'] = prem['purchase_dt'].dt.to_period('M')

def extract_marca(opp):
    s = str(opp).upper()
    if 'RECHAZO' in s:
        idx = s.find('GRUPO PREMIER-')
        if idx >= 0: s = s[idx:]
    parts = s.split('PREMIER-')
    if len(parts) < 2: return 'OTROS'
    tokens = parts[1].strip().split()
    if not tokens: return 'OTROS'
    brand = tokens[0]
    if brand in ('BULK', 'TIPIFICACI'): return 'OTROS'
    return brand

def celula_short(opp):
    s = str(opp).upper()
    if 'CULIAC' in s: return 'Culiacán'
    if 'HERMOSILLO' in s: return 'Hermosillo'
    if 'MAZATL' in s: return 'Mazatlán'
    return 'Otro'

prem['marca'] = prem['Nombre de la oportunidad'].apply(extract_marca)
prem['celula_short'] = prem['Nombre de la oportunidad'].apply(celula_short)
prem = prem[(prem['marca'] != 'OTROS') & (prem['celula_short'] != 'Otro')].copy()

JAN = pd.Period('2026-01', 'M')
FEB = pd.Period('2026-02', 'M')

def get_funnel(data, period):
    cr = data[data['creation_month'] == period]
    pu = data[data['purchase_month'] == period]
    q = int(cr['scheduled.1'].sum())
    m = int(cr['made.1'].sum())
    a = int(cr['approved'].sum())
    p = int(pu['purchased'].sum())
    return q, m, a, p

def calc_rates(q, m, a, p):
    return {
        'Q→M%': (m/q*100) if q else 0,
        'M→A%': (a/m*100) if m else 0,
        'A→P%': (p/a*100) if a else 0,
        'Q→P%': (p/q*100) if q else 0,
    }

def fmt_pct(v): return f"{v:.1f}%"
def fmt_delta(v, is_pct=False):
    sign = '+' if v >= 0 else ''
    return f"{sign}{v:.1f}pp" if is_pct else f"{sign}{int(v)}"

# ── Group-level rates (for "vs Grupo" benchmarks) ──────────────
gj_q, gj_m, gj_a, gj_p = get_funnel(prem, JAN)
gf_q, gf_m, gf_a, gf_p = get_funnel(prem, FEB)
grupo_feb_rates = calc_rates(gf_q, gf_m, gf_a, gf_p)

# ── PPTX helpers ────────────────────────────────────────────────
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, l, t, w, h, text, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb

def add_rect(slide, l, t, w, h, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s

def add_rounded_rect(slide, l, t, w, h, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    s.line.width = Pt(0.5)
    return s

# ── Build presentation ──────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

celulas_order = ['Culiacán', 'Hermosillo', 'Mazatlán']

for cel_name in celulas_order:
    cel_data = prem[prem['celula_short'] == cel_name]
    cel_color = CELULA_COLORS[cel_name]

    # Célula-level funnel
    cj_q, cj_m, cj_a, cj_p = get_funnel(cel_data, JAN)
    cf_q, cf_m, cf_a, cf_p = get_funnel(cel_data, FEB)
    cel_jan_rates = calc_rates(cj_q, cj_m, cj_a, cj_p)
    cel_feb_rates = calc_rates(cf_q, cf_m, cf_a, cf_p)

    # Brands in this célula with Feb purchases, sorted desc
    feb_cel = cel_data[cel_data['purchase_month'] == FEB]
    brand_purchases_feb = feb_cel.groupby('marca')['purchased'].sum().sort_values(ascending=False)
    # Include brands with Jan purchases too (even if 0 in Feb)
    jan_cel = cel_data[cel_data['purchase_month'] == JAN]
    brand_purchases_jan = jan_cel.groupby('marca')['purchased'].sum()
    all_brands = sorted(set(brand_purchases_feb.index) | set(brand_purchases_jan.index),
                        key=lambda x: -int(brand_purchases_feb.get(x, 0)))
    # Only brands with at least 1 purchase in Jan or Feb
    active_brands = [b for b in all_brands if int(brand_purchases_feb.get(b, 0)) + int(brand_purchases_jan.get(b, 0)) > 0]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, 0, 0, 13.333, 0.06, cel_color)

    # Title
    add_textbox(slide, 0.6, 0.15, 9, 0.4, f"{cel_name} — Funnel Comparativo por Marca", 24, True, SEC_TEXT)
    add_textbox(slide, 0.6, 0.55, 9, 0.25, "Enero vs Febrero 2026 | Compras por purchase_date · Funnel por fecha de creación", 10, False, GRAY_TEXT)

    # Célula total badge (top right)
    delta_p = cf_p - cj_p
    pct_chg = (delta_p / cj_p * 100) if cj_p > 0 else 0
    badge = add_rect(slide, 10.5, 0.12, 2.5, 0.75, cel_color)
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{cel_name.upper()}"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = f"Feb 2026"
    p2.font.size = Pt(8)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.LEFT

    add_textbox(slide, 10.5, 0.6, 2.5, 0.35, f"{cf_p} compras", 20, True, cel_color, PP_ALIGN.CENTER)
    sign = '+' if delta_p >= 0 else ''
    d_color = GREEN if delta_p > 0 else (RED if delta_p < 0 else MUTED)
    add_textbox(slide, 10.5, 0.92, 2.5, 0.25, f"{sign}{delta_p} vs Ene ({sign}{pct_chg:.0f}%)", 9, False, d_color, PP_ALIGN.CENTER)

    # ── Brand tables (up to 4 across, or stacked if >4) ──
    n_brands = len(active_brands)
    max_cols = min(n_brands, 4)
    table_w = min(3.0, (12.1 - (max_cols - 1) * 0.12) / max_cols)
    table_gap = 0.12
    table_start_x = 0.6
    table_start_y = 1.25

    for bi, marca in enumerate(active_brands):
        col = bi % max_cols
        row_group = bi // max_cols
        tx = table_start_x + col * (table_w + table_gap)
        ty = table_start_y + row_group * 3.2  # stack down if >4 brands

        brand_data = cel_data[cel_data['marca'] == marca]
        jq, jm, ja, jp = get_funnel(brand_data, JAN)
        fq, fm, fa, fp = get_funnel(brand_data, FEB)
        jan_r = calc_rates(jq, jm, ja, jp)
        feb_r = calc_rates(fq, fm, fa, fp)

        brand_color = BRAND_COLORS.get(marca, NAVY)

        # Brand header bar
        hdr = add_rect(slide, tx, ty, table_w, 0.28, brand_color)
        htf = hdr.text_frame
        hp = htf.paragraphs[0]
        hp.text = f"{marca}  ({fp} compras)"
        hp.font.size = Pt(9)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER

        # Table: 9 rows x 7 cols (Métrica, Ene, Feb, Δ MoM, vs Grupo, vs Célula, vs Marca)
        # Simplified to 5 cols: Métrica, Ene, Feb, Δ MoM, vs Grupo
        tbl_rows = 9
        tbl_cols = 5
        ts = slide.shapes.add_table(tbl_rows, tbl_cols, Inches(tx), Inches(ty + 0.28), Inches(table_w), Inches(2.6))
        tbl = ts.table

        col_widths = [0.28, 0.17, 0.17, 0.18, 0.20]
        for j in range(tbl_cols):
            tbl.columns[j].width = Inches(table_w * col_widths[j])

        # Header
        th = ['Métrica', 'Ene', 'Feb', 'Δ MoM', 'vs Grupo']
        for j, h in enumerate(th):
            c = tbl.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = TABLE_HEADER
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

        # Rows
        tbl_data = [
            ('Quotes', jq, fq, False, None),
            ('Inspecciones', jm, fm, False, None),
            ('Aprobados', ja, fa, False, None),
            ('Compras', jp, fp, False, None),
            ('Q→M%', jan_r['Q→M%'], feb_r['Q→M%'], True, grupo_feb_rates['Q→M%']),
            ('M→A%', jan_r['M→A%'], feb_r['M→A%'], True, grupo_feb_rates['M→A%']),
            ('A→P%', jan_r['A→P%'], feb_r['A→P%'], True, grupo_feb_rates['A→P%']),
            ('Q→P%', jan_r['Q→P%'], feb_r['Q→P%'], True, grupo_feb_rates['Q→P%']),
        ]

        for ri, (name, jv, fv, is_rate, grupo_ref) in enumerate(tbl_data):
            row_idx = ri + 1
            bg = WHITE if row_idx % 2 == 0 else ROW_ALT
            delta = fv - jv

            if is_rate:
                jstr = fmt_pct(jv)
                fstr = fmt_pct(fv)
                dstr = fmt_delta(delta, True)
                vs_grupo = fmt_delta(fv - grupo_ref, True) if grupo_ref is not None else '—'
            else:
                jstr = str(int(jv))
                fstr = str(int(fv))
                dstr = fmt_delta(delta)
                vs_grupo = '—'

            vals = [name, jstr, fstr, dstr, vs_grupo]
            for j, v in enumerate(vals):
                c = tbl.cell(row_idx, j)
                c.text = v
                c.fill.solid()
                c.fill.fore_color.rgb = bg
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(7)
                    p.font.color.rgb = SEC_TEXT
                    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                    if j == 3:
                        p.font.color.rgb = GREEN if delta > 0 else (RED if delta < 0 else SEC_TEXT)
                    if j == 4 and is_rate and grupo_ref is not None:
                        diff = fv - grupo_ref
                        p.font.color.rgb = GREEN if diff > 0 else (RED if diff < 0 else SEC_TEXT)
                    if j == 2 and not is_rate and ri == 3:
                        p.font.bold = True
                        p.font.color.rgb = brand_color

    # ── Consolidated KPI cards at bottom ──
    cons_y = 4.65 if n_brands <= 4 else 7.5  # only if space
    if n_brands <= 4:
        add_textbox(slide, 0.6, cons_y - 0.3, 5, 0.25, f"{cel_name} Consolidado", 13, True, cel_color)

        # 4 volume KPI cards
        card_w = 2.2
        kpis = [
            ('Quotes', cf_q, cj_q),
            ('Inspecciones', cf_m, cj_m),
            ('Aprobados', cf_a, cj_a),
            ('Compras', cf_p, cj_p),
        ]
        for ki, (kl, fv, jv) in enumerate(kpis):
            kx = 0.6 + ki * (card_w + 0.15)
            add_rounded_rect(slide, kx, cons_y, card_w, 0.8, WHITE)
            add_textbox(slide, kx + 0.1, cons_y + 0.02, card_w - 0.2, 0.15, kl, 8, False, MUTED)
            add_textbox(slide, kx + 0.1, cons_y + 0.12, card_w - 0.2, 0.35, str(fv), 20, True, SEC_TEXT, PP_ALIGN.CENTER)
            d = fv - jv
            dtxt = f"{'+' if d >= 0 else ''}{d} vs Ene"
            dc = GREEN if d > 0 else (RED if d < 0 else MUTED)
            add_textbox(slide, kx + 0.1, cons_y + 0.52, card_w - 0.2, 0.2, dtxt, 8, True, dc, PP_ALIGN.CENTER)

        # 4 rate cards
        rate_y = cons_y + 0.95
        rates = [
            ('Q→M%', cel_feb_rates['Q→M%'], cel_jan_rates['Q→M%']),
            ('M→A%', cel_feb_rates['M→A%'], cel_jan_rates['M→A%']),
            ('A→P%', cel_feb_rates['A→P%'], cel_jan_rates['A→P%']),
            ('Q→P%', cel_feb_rates['Q→P%'], cel_jan_rates['Q→P%']),
        ]
        for ki, (kl, fv, jv) in enumerate(rates):
            kx = 0.6 + ki * (card_w + 0.15)
            add_rounded_rect(slide, kx, rate_y, card_w, 0.7, WHITE)
            add_textbox(slide, kx + 0.1, rate_y + 0.02, card_w - 0.2, 0.15, kl, 8, False, MUTED)
            vc = GREEN if fv >= jv else RED
            add_textbox(slide, kx + 0.1, rate_y + 0.1, card_w - 0.2, 0.3, fmt_pct(fv), 16, True, vc, PP_ALIGN.CENTER)
            dpp = fv - jv
            add_textbox(slide, kx + 0.1, rate_y + 0.45, card_w - 0.2, 0.2, f"{fmt_delta(dpp, True)} vs Ene", 8, False, MUTED, PP_ALIGN.CENTER)

        # Lectura box
        lx = 9.8
        lw = 3.2
        lh = 1.65
        lbox = add_rounded_rect(slide, lx, cons_y, lw, lh, LECTURA_BLUE)
        add_rect(slide, lx, cons_y, 0.06, lh, cel_color)

        add_textbox(slide, lx + 0.15, cons_y + 0.05, lw - 0.3, 0.2, f"Lectura {cel_name}", 10, True, SEC_TEXT)

        # Generate insights
        insights = []
        if delta_p > 0:
            insights.append(f"• {cel_name} creció {'+' if delta_p > 0 else ''}{delta_p} compras MoM ({pct_chg:+.0f}%)")
        elif delta_p < 0:
            insights.append(f"• {cel_name} cayó {delta_p} compras MoM ({pct_chg:.0f}%)")
        else:
            insights.append(f"• {cel_name} se mantuvo estable ({cf_p} compras)")

        # Top brand
        if active_brands:
            top_b = active_brands[0]
            top_p = int(brand_purchases_feb.get(top_b, 0))
            share = (top_p / cf_p * 100) if cf_p else 0
            insights.append(f"• {top_b} lidera con {top_p} compras ({share:.0f}% share)")

        # Brand growth
        for b in active_brands[:3]:
            bp_feb = int(brand_purchases_feb.get(b, 0))
            bp_jan = int(brand_purchases_jan.get(b, 0))
            bd = bp_feb - bp_jan
            if bd >= 3:
                insights.append(f"• {b} creció {'+' if bd > 0 else ''}{bd} ({bp_jan}→{bp_feb})")
                break

        # Q→M delta
        qm_d = cel_feb_rates['Q→M%'] - cel_jan_rates['Q→M%']
        if abs(qm_d) >= 3:
            direction = "mejoró" if qm_d > 0 else "se deterioró"
            insights.append(f"• Q→M% {direction} ({fmt_delta(qm_d, True)})")

        # A→P
        ap_d = cel_feb_rates['A→P%'] - cel_jan_rates['A→P%']
        if abs(ap_d) >= 5:
            direction = "mejoró" if ap_d > 0 else "se deterioró"
            insights.append(f"• A→P% {direction} ({fmt_delta(ap_d, True)})")

        # Q→P
        qp_f = cel_feb_rates['Q→P%']
        qp_d = qp_f - cel_jan_rates['Q→P%']
        insights.append(f"• Q→P% consolidado: {fmt_pct(qp_f)} ({fmt_delta(qp_d, True)})")

        insight_text = '\n'.join(insights[:6])
        add_textbox(slide, lx + 0.15, cons_y + 0.28, lw - 0.3, lh - 0.35, insight_text, 9, False, SEC_TEXT)

    # Footer
    add_textbox(slide, 0.6, 7.15, 10, 0.25, f"{cel_name} | Grupo Premier | TAS Analytics | Febrero 2026", 9, False, GRAY_TEXT)

# ── Save ────────────────────────────────────────────────────────
output_path = os.path.expanduser("~/Downloads/Grupo_Premier_Funnel_por_Celula.pptx")
prs.save(output_path)
print(f"✅ Archivo generado: {output_path}")
print(f"   Slides: {len(prs.slides)}")
for cel in celulas_order:
    cd = prem[prem['celula_short'] == cel]
    fp = int(cd[cd['purchase_month'] == FEB]['purchased'].sum())
    brands = cd[cd['purchase_month'] == FEB].groupby('marca')['purchased'].sum().sort_values(ascending=False)
    print(f"   {cel}: {fp} compras — {dict(brands)}")
