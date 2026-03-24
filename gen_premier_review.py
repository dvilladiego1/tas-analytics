#!/usr/bin/env python3
"""
Grupo Premier — Review por Marca — Febrero 2026
Replica el formato exacto del PDF de Grupo Continental.
"""
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
NAVY = RGBColor(0x00, 0x4E, 0x98)
TABLE_HEADER = RGBColor(0x2B, 0x47, 0x8B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GRAY_TEXT = RGBColor(0x7A, 0x7A, 0x7A)
SEC_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)
INSIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
HIGHLIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LOWLIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LECTURA_BLUE = RGBColor(0xE3, 0xEE, 0xFB)
TRAJ_GREEN = RGBColor(0xE8, 0xF5, 0xE9)

# Brand colors for badges
BRAND_COLORS = {
    'CHEVROLET': RGBColor(0x1B, 0x2A, 0x4A),
    'BYD': RGBColor(0x22, 0x6B, 0x22),
    'TOYOTA': RGBColor(0xCC, 0x00, 0x00),
    'KIA': RGBColor(0x00, 0x4E, 0x98),
    'MERCEDESBENZ': RGBColor(0x33, 0x33, 0x33),
    'HYUNDAI': RGBColor(0x00, 0x2C, 0x5F),
    'BUICK': RGBColor(0x8B, 0x00, 0x00),
    'AUTOCOUNTRY': RGBColor(0x00, 0x80, 0x80),
    'BMW': RGBColor(0x1C, 0x69, 0xD4),
}

# ── Data Loading ────────────────────────────────────────────────
CSV_PATH = os.path.expanduser("~/Downloads/Summary MKP _ TAS & BULK - Snapshot_Diario (1).csv")

df = pd.read_csv(CSV_PATH, encoding='latin-1')

# Normalize column name for channel type
channel_col = [c for c in df.columns if 'BULK' in c and 'TAS' in c][0]
df.rename(columns={channel_col: 'channel'}, inplace=True)

# Find the actual column names (encoding may mangle accents)
fecha_col = [c for c in df.columns if 'Fecha de creaci' in c][0]
celula_col = [c for c in df.columns if 'lula' in c][0]

# Filter Grupo Premier, TAS only
mask = (df['Grupo'].str.upper().str.contains('PREMIER', na=False)) & (df['channel'] != 'BULK')
prem = df[mask].copy()

# Parse dates
prem['creation_date'] = pd.to_datetime(prem[fecha_col], format='%d/%m/%Y', errors='coerce')
prem['purchase_dt'] = pd.to_datetime(prem['purchase_date'], format='%d/%m/%Y', errors='coerce')
prem['creation_month'] = prem['creation_date'].dt.to_period('M')
prem['purchase_month'] = prem['purchase_dt'].dt.to_period('M')

# Extract brand from 'Nombre de la oportunidad' e.g. 'GRUPO PREMIER-BYD CULIACÁN'
def extract_marca(opp):
    s = str(opp).upper()
    if 'GRUPO PREMIER-' not in s and 'PREMIER-' not in s:
        return 'OTROS'
    # Remove "Rechazo 7 días " prefix if present
    if 'RECHAZO' in s:
        idx = s.find('GRUPO PREMIER-')
        if idx >= 0:
            s = s[idx:]
    parts = s.split('PREMIER-')
    if len(parts) < 2:
        return 'OTROS'
    after = parts[1].strip()
    # Brand is the first word before city name
    tokens = after.split()
    if len(tokens) == 0:
        return 'OTROS'
    brand = tokens[0]
    if brand in ('BULK', 'TIPIFICACI'):
        return 'OTROS'
    return brand

prem['marca'] = prem['Nombre de la oportunidad'].apply(extract_marca)
# Filter out OTROS
prem = prem[prem['marca'] != 'OTROS'].copy()

# Normalize célula: extract city from oportunidad
def celula_short_from_opp(opp):
    s = str(opp).upper()
    if 'CULIAC' in s: return 'Culiacán'
    if 'HERMOSILLO' in s: return 'Hermosillo'
    if 'MAZATL' in s: return 'Mazatlán'
    return 'Otro'

prem['celula_short'] = prem['Nombre de la oportunidad'].apply(celula_short_from_opp)
prem = prem[prem['celula_short'] != 'Otro'].copy()

# ── Helper functions ────────────────────────────────────────────
def funnel_by_group(data, group_col, period_col_creation='creation_month', periods=None):
    """Compute funnel for given grouping. Purchases by purchase_date, rest by creation_date."""
    results = {}
    for period in (periods or []):
        cr = data[data[period_col_creation] == period]
        pu = data[data['purchase_month'] == period]
        results[str(period)] = {
            'quotes': int(cr['scheduled.1'].sum()) if 'scheduled.1' in cr.columns else 0,
            'made': int(cr['made.1'].sum()) if 'made.1' in cr.columns else 0,
            'approved': int(cr['approved'].sum()),
            'purchased': int(pu['purchased'].sum()),
        }
    return results

def calc_rates(q, m, a, p):
    qm = (m/q*100) if q > 0 else 0
    ma = (a/m*100) if m > 0 else 0
    ap = (p/a*100) if a > 0 else 0
    qp = (p/q*100) if q > 0 else 0
    return {'Q→M%': qm, 'M→A%': ma, 'A→P%': ap, 'Q→P%': qp}

JAN = pd.Period('2026-01', 'M')
FEB = pd.Period('2026-02', 'M')

# ── Compute all metrics ─────────────────────────────────────────
# 1. Feb purchases by brand (sorted)
feb_purchases = prem[prem['purchase_month'] == FEB]
jan_purchases = prem[prem['purchase_month'] == JAN]
brand_feb = feb_purchases.groupby('marca')['purchased'].sum().sort_values(ascending=False)
brand_jan = jan_purchases.groupby('marca')['purchased'].sum()

brands_sorted = []
for marca in brand_feb.index:
    if brand_feb[marca] > 0:
        feb_p = int(brand_feb[marca])
        jan_p = int(brand_jan.get(marca, 0))
        # Which células
        cells = sorted(feb_purchases[feb_purchases['marca'] == marca]['celula_short'].unique())
        brands_sorted.append({'marca': marca, 'feb': feb_p, 'jan': jan_p, 'celulas': cells})

total_feb_purchases = int(brand_feb.sum())
total_brands = len(brands_sorted)

# 2. Historical purchases by célula
hist_months = [pd.Period(f'2025-{m:02d}', 'M') for m in range(8, 13)] + [pd.Period(f'2026-{m:02d}', 'M') for m in range(1, 3)]
celulas_order = ['Culiacán', 'Hermosillo', 'Mazatlán']

hist_data = {}
for cel in celulas_order:
    hist_data[cel] = {}
    cel_data = prem[prem['celula_short'] == cel]
    for mo in hist_months:
        hist_data[cel][str(mo)] = int(cel_data[cel_data['purchase_month'] == mo]['purchased'].sum())

# 3. Consolidated funnel Jan vs Feb
def get_funnel(data, period):
    cr = data[data['creation_month'] == period]
    pu = data[data['purchase_month'] == period]
    q = int(cr['scheduled.1'].sum())
    m = int(cr['made.1'].sum())
    a = int(cr['approved'].sum())
    p = int(pu['purchased'].sum())
    return q, m, a, p

jan_q, jan_m, jan_a, jan_p = get_funnel(prem, JAN)
feb_q, feb_m, feb_a, feb_p = get_funnel(prem, FEB)
jan_rates = calc_rates(jan_q, jan_m, jan_a, jan_p)
feb_rates = calc_rates(feb_q, feb_m, feb_a, feb_p)

# 4. Brand x Célula detail
brand_celula_data = {}
for b in brands_sorted:
    marca = b['marca']
    brand_data = prem[prem['marca'] == marca]
    # Find which células have this brand (in either Jan or Feb)
    all_cells = sorted(brand_data['celula_short'].unique())

    celula_detail = {}
    for cel in all_cells:
        cd = brand_data[brand_data['celula_short'] == cel]
        jq, jm, ja, jp = get_funnel(cd, JAN)
        fq, fm, fa, fp = get_funnel(cd, FEB)
        celula_detail[cel] = {
            'jan': {'q': jq, 'm': jm, 'a': ja, 'p': jp, 'rates': calc_rates(jq, jm, ja, jp)},
            'feb': {'q': fq, 'm': fm, 'a': fa, 'p': fp, 'rates': calc_rates(fq, fm, fa, fp)},
        }

    # Brand consolidated
    jq, jm, ja, jp = get_funnel(brand_data, JAN)
    fq, fm, fa, fp = get_funnel(brand_data, FEB)
    brand_celula_data[marca] = {
        'celulas': celula_detail,
        'consolidated': {
            'jan': {'q': jq, 'm': jm, 'a': ja, 'p': jp, 'rates': calc_rates(jq, jm, ja, jp)},
            'feb': {'q': fq, 'm': fm, 'a': fa, 'p': fp, 'rates': calc_rates(fq, fm, fa, fp)},
        }
    }

# ── PPTX Generation ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.5)
    return shape

def fmt_delta(val, is_pct=False):
    if is_pct:
        sign = '+' if val >= 0 else ''
        return f"{sign}{val:.1f}pp"
    else:
        sign = '+' if val >= 0 else ''
        return f"{sign}{int(val)}"

def fmt_pct(val):
    return f"{val:.1f}%"

def delta_color(val):
    return GREEN if val >= 0 else RED

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Top accent bar
add_shape_rect(slide, 0, 0, 13.333, 0.12, NAVY)

# Title content
add_textbox(slide, 0, 2.5, 13.333, 1.0, "GRUPO PREMIER", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.5, 13.333, 0.6, f"{total_feb_purchases} Compras Netas · {total_brands} Marcas · Febrero 2026", 22, False, MUTED, PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.3, 13.333, 0.4, "TAS Analytics · Review por Marca", 14, False, MUTED, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: OVERVIEW — Marcas badges
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, 0, 1.2, 13.333, 0.8, "GRUPO PREMIER", 42, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.1, 13.333, 0.5, "Vista por Marca — Funnel Comparativo por Célula", 18, False, MUTED, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.6, 13.333, 0.4, "Enero vs Febrero 2026", 14, False, MUTED, PP_ALIGN.CENTER)

# Brand badges
badge_w = 1.35
badge_h = 0.55
gap = 0.12
n_brands = len(brands_sorted)
total_w = n_brands * badge_w + (n_brands - 1) * gap
start_x = (13.333 - total_w) / 2
y = 3.8

for i, b in enumerate(brands_sorted):
    x = start_x + i * (badge_w + gap)
    color = BRAND_COLORS.get(b['marca'], NAVY)
    shape = add_shape_rect(slide, x, y, badge_w, badge_h, color)
    # Brand name
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = b['marca']
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Purchase count
    p2 = tf.add_paragraph()
    p2.text = str(b['feb'])
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

add_textbox(slide, 0, 5.0, 13.333, 0.3, "TAS Analytics · Marzo 2026", 11, False, MUTED, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: EVOLUCIÓN HISTÓRICA — Compras por Célula
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Top bar
add_shape_rect(slide, 0, 0, 13.333, 0.06, TABLE_HEADER)

# Title
add_textbox(slide, 0.6, 0.3, 10, 0.5, "Evolución Histórica — Compras por Célula", 28, True, SEC_TEXT)
add_textbox(slide, 0.6, 0.8, 10, 0.3, "Noviembre 2025 – Febrero 2026  |  Grupo Premier (TAS)", 13, False, GRAY_TEXT)

# Chart data
from pptx.chart.data import CategoryChartData
chart_data = CategoryChartData()
# Only show months with data (Nov 25 onward)
active_months = [m for m in hist_months if sum(hist_data[c].get(str(m), 0) for c in celulas_order) > 0]
month_labels = []
for m in active_months:
    mo = str(m)
    if '2025-11' in mo: month_labels.append('Nov 25')
    elif '2025-12' in mo: month_labels.append('Dic 25')
    elif '2026-01' in mo: month_labels.append('Ene 26')
    elif '2026-02' in mo: month_labels.append('Feb 26')
    else: month_labels.append(mo)

chart_data.categories = month_labels

series_colors_map = {
    'Culiacán': RGBColor(0x1B, 0x2A, 0x4A),
    'Hermosillo': RGBColor(0x3B, 0x7D, 0xDD),
    'Mazatlán': RGBColor(0x00, 0xB4, 0x8A),
}

for cel in celulas_order:
    vals = [hist_data[cel].get(str(m), 0) for m in active_months]
    chart_data.add_series(cel, vals)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, Inches(0.6), Inches(1.3), Inches(8.0), Inches(4.8), chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.include_in_layout = False

plot = chart.plots[0]
plot.gap_width = 80

for i, cel in enumerate(celulas_order):
    series = plot.series[i]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = series_colors_map[cel]
    series.has_data_labels = True
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = WHITE
    series.data_labels.number_format = '0'
    series.data_labels.show_value = True

# Right side table
tbl_left = 9.2
tbl_top = 1.5
# Trimmed months for table: Dic 25, Ene 26, Feb 26
tbl_months = ['Dic 25', 'Ene 26', 'Feb 26']
tbl_period_keys = ['2025-12', '2026-01', '2026-02']

rows = len(celulas_order) + 2  # header + cells + total
cols = 4  # Célula + 3 months
table_shape = slide.shapes.add_table(rows, cols, Inches(tbl_left), Inches(tbl_top), Inches(3.8), Inches(2.0))
table = table_shape.table

# Header
headers = ['Célula'] + tbl_months
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Data rows
for i, cel in enumerate(celulas_order):
    row_idx = i + 1
    table.cell(row_idx, 0).text = cel
    for p in table.cell(row_idx, 0).text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = SEC_TEXT
    for j, pk in enumerate(tbl_period_keys):
        val = hist_data[cel].get(pk, 0)
        c = table.cell(row_idx, j + 1)
        c.text = str(val)
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = SEC_TEXT
            p.alignment = PP_ALIGN.CENTER

# Total row
total_row = len(celulas_order) + 1
table.cell(total_row, 0).text = 'TOTAL'
table.cell(total_row, 0).fill.solid()
table.cell(total_row, 0).fill.fore_color.rgb = NAVY
for p in table.cell(total_row, 0).text_frame.paragraphs:
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE

for j, pk in enumerate(tbl_period_keys):
    total_val = sum(hist_data[c].get(pk, 0) for c in celulas_order)
    c = table.cell(total_row, j + 1)
    c.text = str(total_val)
    c.fill.solid()
    c.fill.fore_color.rgb = NAVY
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Trajectory box
total_hist = sum(sum(hist_data[c].get(str(m), 0) for c in celulas_order) for m in active_months)
traj_text = f"Trayectoria: Grupo Premier arrancó en Nov 25. Desde entonces, de 2 compras (Nov 25) a {total_feb_purchases} (Feb 26). Pico: {total_feb_purchases} en Feb 26. Acumulado: {total_hist} compras en {len(active_months)} meses."
box = add_shape_rect(slide, 0.6, 6.3, 12.1, 0.6, TRAJ_GREEN)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = traj_text
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = SEC_TEXT

add_textbox(slide, 0.6, 7.1, 10, 0.3, "Grupo Premier — Febrero 2026 | Fuente: Snapshot Diario CSV", 9, False, GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: BREAKDOWN POR MARCA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.06, TABLE_HEADER)

add_textbox(slide, 0.6, 0.3, 10, 0.5, "Breakdown por Marca — Febrero 2026", 28, True, SEC_TEXT)
add_textbox(slide, 0.6, 0.8, 10, 0.3, f"Distribución de {total_feb_purchases} compras netas por marca automotriz", 13, False, GRAY_TEXT)

# Table
rows = len(brands_sorted) + 1
cols = 4
table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.3), Inches(7.5), Inches(0.4 * rows))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(2.5)

headers = ['Marca', 'Compras', '% Total', 'Célula(s)']
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        if j > 0:
            p.alignment = PP_ALIGN.CENTER

for i, b in enumerate(brands_sorted):
    row_idx = i + 1
    pct = b['feb'] / total_feb_purchases * 100

    table.cell(row_idx, 0).text = b['marca']
    table.cell(row_idx, 1).text = str(b['feb'])
    table.cell(row_idx, 2).text = f"{pct:.1f}%"
    table.cell(row_idx, 3).text = ', '.join(b['celulas'])

    bg_color = WHITE if row_idx % 2 == 0 else ROW_ALT
    for j in range(4):
        c = table.cell(row_idx, j)
        c.fill.solid()
        c.fill.fore_color.rgb = bg_color
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = SEC_TEXT
            if j in [1, 2]:
                p.alignment = PP_ALIGN.CENTER

# Horizontal bar chart on the right
bar_x = 8.8
bar_y = 1.3
max_bar_w = 4.0
max_val = brands_sorted[0]['feb'] if brands_sorted else 1

for i, b in enumerate(brands_sorted):
    y_pos = bar_y + i * 0.5
    w = max(0.3, (b['feb'] / max_val) * max_bar_w)
    color = BRAND_COLORS.get(b['marca'], NAVY)
    bar = add_shape_rect(slide, bar_x, y_pos, w, 0.35, color)
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{b['marca']}: {b['feb']}"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

add_textbox(slide, 0.6, 6.8, 10, 0.3, "Grupo Premier — Breakdown Febrero 2026", 9, False, GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: CONVERSIONES FUNNEL CONSOLIDADO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.06, TABLE_HEADER)

add_textbox(slide, 0.6, 0.2, 10, 0.5, "Conversiones Funnel — Grupo Premier Consolidado", 24, True, SEC_TEXT)
add_textbox(slide, 0.6, 0.65, 10, 0.3, "Febrero 2026 vs Enero 2026 | TAS", 13, False, GRAY_TEXT)

# 4 KPI cards top row
kpi_labels = ['Quotes', 'Inspecciones', 'Aprobados', 'Compras']
feb_vals = [feb_q, feb_m, feb_a, feb_p]
jan_vals = [jan_q, jan_m, jan_a, jan_p]
card_w = 2.8
card_h = 1.2
gap = 0.3
start_x = 0.6
y_top = 1.05

for i, (label, fv, jv) in enumerate(zip(kpi_labels, feb_vals, jan_vals)):
    x = start_x + i * (card_w + gap)
    card = add_rounded_rect(slide, x, y_top, card_w, card_h, WHITE)
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_textbox(slide, x + 0.15, y_top + 0.05, card_w - 0.3, 0.25, label, 10, False, MUTED)
    add_textbox(slide, x + 0.15, y_top + 0.25, card_w - 0.3, 0.55, str(fv), 28, True, NAVY, PP_ALIGN.CENTER)

    delta = fv - jv
    d_text = f"▲ {fmt_delta(delta)}" if delta >= 0 else f"▼ {fmt_delta(delta)}"
    d_color = GREEN if delta >= 0 else RED
    add_textbox(slide, x + 0.15, y_top + 0.85, card_w - 0.3, 0.25, d_text, 10, True, d_color, PP_ALIGN.CENTER)

# 4 rate cards
rate_labels = ['Q→M%', 'M→A%', 'A→P%', 'Q→P%']
feb_rate_vals = [feb_rates[k] for k in rate_labels]
jan_rate_vals = [jan_rates[k] for k in rate_labels]
y_rate = 2.45

for i, (label, fr, jr) in enumerate(zip(rate_labels, feb_rate_vals, jan_rate_vals)):
    x = start_x + i * (card_w + gap)
    card = add_rounded_rect(slide, x, y_rate, card_w, 1.0, WHITE)
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_textbox(slide, x + 0.15, y_rate + 0.02, card_w - 0.3, 0.2, label, 10, False, MUTED)

    val_color = GREEN if fr >= jr else RED
    add_textbox(slide, x + 0.15, y_rate + 0.2, card_w - 0.3, 0.45, fmt_pct(fr), 24, True, val_color, PP_ALIGN.CENTER)

    delta_pp = fr - jr
    d_text = f"vs Ene: {fmt_delta(delta_pp, True)}"
    add_textbox(slide, x + 0.15, y_rate + 0.7, card_w - 0.3, 0.25, d_text, 9, False, MUTED, PP_ALIGN.CENTER)

# Comparison table
tbl_y = 3.7
rows = 9
cols = 6
table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(tbl_y), Inches(12.1), Inches(3.2))
table = table_shape.table

# Header
comp_headers = ['Métrica', 'Ene 2026', 'Feb 2026', 'Cambio', 'Var%', 'Tend.']
for j, h in enumerate(comp_headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Data
metric_rows = [
    ('Quotes', jan_q, feb_q, False),
    ('Inspecciones', jan_m, feb_m, False),
    ('Aprobados', jan_a, feb_a, False),
    ('Compras', jan_p, feb_p, False),
    ('Q→M%', jan_rates['Q→M%'], feb_rates['Q→M%'], True),
    ('M→A%', jan_rates['M→A%'], feb_rates['M→A%'], True),
    ('A→P%', jan_rates['A→P%'], feb_rates['A→P%'], True),
    ('Q→P%', jan_rates['Q→P%'], feb_rates['Q→P%'], True),
]

for i, (name, jv, fv, is_rate) in enumerate(metric_rows):
    row_idx = i + 1
    bg = WHITE if row_idx % 2 == 0 else ROW_ALT

    delta = fv - jv
    if is_rate:
        jan_str = fmt_pct(jv)
        feb_str = fmt_pct(fv)
        change_str = fmt_delta(delta, True)
        var_str = ''
    else:
        jan_str = str(int(jv))
        feb_str = str(int(fv))
        change_str = fmt_delta(delta)
        var_str = f"{delta/jv*100:+.1f}%" if jv > 0 else '—'

    tend = '↑' if delta > 0 else ('↓' if delta < 0 else '→')

    vals = [name, jan_str, feb_str, change_str, var_str, tend]
    for j, v in enumerate(vals):
        cell = table.cell(row_idx, j)
        cell.text = v
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = SEC_TEXT
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            # Color deltas
            if j == 3:
                p.font.color.rgb = GREEN if delta > 0 else (RED if delta < 0 else SEC_TEXT)
            if j == 5:
                p.font.color.rgb = GREEN if delta > 0 else (RED if delta < 0 else SEC_TEXT)
                p.font.size = Pt(14)

add_textbox(slide, 0.6, 7.1, 10, 0.3, "Grupo Premier | TAS Analytics | Febrero 2026", 9, False, GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDES 6+: PER-BRAND FUNNEL COMPARATIVO
# ═══════════════════════════════════════════════════════════════
# Only create detailed slides for brands with >= 2 purchases in Feb
major_brands = [b for b in brands_sorted if b['feb'] >= 2]
minor_brands = [b for b in brands_sorted if b['feb'] < 2]

for brand_info in major_brands:
    marca = brand_info['marca']
    bd = brand_celula_data[marca]
    cells = list(bd['celulas'].keys())
    cons = bd['consolidated']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Top bar with brand color
    brand_color = BRAND_COLORS.get(marca, NAVY)
    add_shape_rect(slide, 0, 0, 13.333, 0.06, brand_color)

    # Title
    add_textbox(slide, 0.6, 0.15, 9, 0.4, f"{marca} — Funnel Comparativo por Célula", 24, True, SEC_TEXT)
    add_textbox(slide, 0.6, 0.55, 9, 0.25, "Enero vs Febrero 2026 | Compras por purchase_date · Funnel por fecha de creación", 10, False, GRAY_TEXT)

    # Brand total badge (top right)
    badge = add_shape_rect(slide, 10.8, 0.15, 2.2, 0.65, brand_color)
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{marca} TOTAL"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = f"Feb 2026"
    p2.font.size = Pt(8)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.LEFT

    # Purchase count next to badge
    feb_brand_p = cons['feb']['p']
    jan_brand_p = cons['jan']['p']
    delta_p = feb_brand_p - jan_brand_p
    pct_change = (delta_p / jan_brand_p * 100) if jan_brand_p > 0 else 0

    add_textbox(slide, 11.0, 0.55, 2.0, 0.4, f"{feb_brand_p} compras", 18, True, WHITE, PP_ALIGN.RIGHT)
    sign = '+' if delta_p >= 0 else ''
    add_textbox(slide, 10.8, 0.85, 2.2, 0.3, f"{sign}{delta_p} vs Ene ({sign}{pct_change:.0f}%)", 9, False, MUTED if delta_p == 0 else (GREEN if delta_p > 0 else RED), PP_ALIGN.CENTER)

    # Célula tables - up to 3 side by side
    n_cells = len(cells)
    if n_cells == 0:
        continue

    table_w = min(4.0, 12.1 / n_cells - 0.15)
    table_gap = 0.15
    total_tables_w = n_cells * table_w + (n_cells - 1) * table_gap
    table_start_x = 0.6

    for ci, cel in enumerate(cells):
        cd = bd['celulas'][cel]
        jan_d = cd['jan']
        feb_d = cd['feb']

        tx = table_start_x + ci * (table_w + table_gap)
        ty = 1.2

        # Célula header bar
        cel_header = add_shape_rect(slide, tx, ty, table_w, 0.3, brand_color)
        tf = cel_header.text_frame
        p = tf.paragraphs[0]
        p.text = cel
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Table: 9 rows x 4 cols (Métrica, Ene, Feb, Δ MoM)
        tbl_rows = 9
        tbl_cols = 4
        ts = slide.shapes.add_table(tbl_rows, tbl_cols, Inches(tx), Inches(ty + 0.3), Inches(table_w), Inches(2.8))
        tbl = ts.table

        # Set column widths
        tbl.columns[0].width = Inches(table_w * 0.3)
        tbl.columns[1].width = Inches(table_w * 0.2)
        tbl.columns[2].width = Inches(table_w * 0.2)
        tbl.columns[3].width = Inches(table_w * 0.3)

        # Header
        th = ['Métrica', 'Ene', 'Feb', 'Δ MoM']
        for j, h in enumerate(th):
            c = tbl.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = TABLE_HEADER
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

        # Rows
        tbl_data = [
            ('Quotes', jan_d['q'], feb_d['q'], False),
            ('Inspecciones', jan_d['m'], feb_d['m'], False),
            ('Aprobados', jan_d['a'], feb_d['a'], False),
            ('Compras', jan_d['p'], feb_d['p'], False),
            ('Q→M%', jan_d['rates']['Q→M%'], feb_d['rates']['Q→M%'], True),
            ('M→A%', jan_d['rates']['M→A%'], feb_d['rates']['M→A%'], True),
            ('A→P%', jan_d['rates']['A→P%'], feb_d['rates']['A→P%'], True),
            ('Q→P%', jan_d['rates']['Q→P%'], feb_d['rates']['Q→P%'], True),
        ]

        for ri, (name, jv, fv, is_rate) in enumerate(tbl_data):
            row_idx = ri + 1
            bg = WHITE if row_idx % 2 == 0 else ROW_ALT
            delta = fv - jv

            if is_rate:
                jstr = fmt_pct(jv)
                fstr = fmt_pct(fv)
                dstr = fmt_delta(delta, True)
            else:
                jstr = str(int(jv))
                fstr = str(int(fv))
                dstr = fmt_delta(delta)

            vals = [name, jstr, fstr, dstr]
            for j, v in enumerate(vals):
                c = tbl.cell(row_idx, j)
                c.text = v
                c.fill.solid()
                c.fill.fore_color.rgb = bg
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.color.rgb = SEC_TEXT
                    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                    if j == 3:
                        p.font.color.rgb = GREEN if delta > 0 else (RED if delta < 0 else SEC_TEXT)
                    if j == 2 and not is_rate and ri == 3:  # Compras Feb value bold
                        p.font.bold = True
                        p.font.color.rgb = brand_color

    # Consolidated section at bottom
    cons_y = 4.5
    add_textbox(slide, 0.6, cons_y - 0.3, 5, 0.3, f"{marca} Consolidado", 14, True, brand_color)

    # 4 KPI cards
    card_w2 = 2.2
    card_h2 = 0.9
    kpi_data = [
        ('Quotes', cons['feb']['q'], cons['jan']['q']),
        ('Inspecciones', cons['feb']['m'], cons['jan']['m']),
        ('Aprobados', cons['feb']['a'], cons['jan']['a']),
        ('Compras', cons['feb']['p'], cons['jan']['p']),
    ]

    for ki, (kl, fv, jv) in enumerate(kpi_data):
        kx = 0.6 + ki * (card_w2 + 0.2)
        card = add_rounded_rect(slide, kx, cons_y, card_w2, card_h2, WHITE)
        add_textbox(slide, kx + 0.1, cons_y + 0.02, card_w2 - 0.2, 0.15, kl, 8, False, MUTED)
        add_textbox(slide, kx + 0.1, cons_y + 0.15, card_w2 - 0.2, 0.4, str(fv), 22, True, SEC_TEXT, PP_ALIGN.CENTER)
        delta = fv - jv
        d_text = f"{'+' if delta >= 0 else ''}{delta} vs Ene"
        d_color = GREEN if delta > 0 else (RED if delta < 0 else MUTED)
        add_textbox(slide, kx + 0.1, cons_y + 0.6, card_w2 - 0.2, 0.2, d_text, 8, True, d_color, PP_ALIGN.CENTER)

    # 4 Rate cards below
    rate_y = cons_y + 1.05
    rate_data = [
        ('Q→M%', cons['feb']['rates']['Q→M%'], cons['jan']['rates']['Q→M%']),
        ('M→A%', cons['feb']['rates']['M→A%'], cons['jan']['rates']['M→A%']),
        ('A→P%', cons['feb']['rates']['A→P%'], cons['jan']['rates']['A→P%']),
        ('Q→P%', cons['feb']['rates']['Q→P%'], cons['jan']['rates']['Q→P%']),
    ]

    for ki, (kl, fv, jv) in enumerate(rate_data):
        kx = 0.6 + ki * (card_w2 + 0.2)
        card = add_rounded_rect(slide, kx, rate_y, card_w2, 0.75, WHITE)
        add_textbox(slide, kx + 0.1, rate_y + 0.02, card_w2 - 0.2, 0.15, kl, 8, False, MUTED)
        val_color = GREEN if fv >= jv else RED
        add_textbox(slide, kx + 0.1, rate_y + 0.12, card_w2 - 0.2, 0.35, fmt_pct(fv), 18, True, val_color, PP_ALIGN.CENTER)
        delta_pp = fv - jv
        d_text = f"{fmt_delta(delta_pp, True)} vs Ene"
        add_textbox(slide, kx + 0.1, rate_y + 0.5, card_w2 - 0.2, 0.2, d_text, 8, False, d_color, PP_ALIGN.CENTER)

    # Lectura box (right side)
    lectura_x = 9.8
    lectura_y = cons_y
    lectura_w = 3.2
    lectura_h = 1.8

    lbox = add_rounded_rect(slide, lectura_x, lectura_y, lectura_w, lectura_h, LECTURA_BLUE)
    # Left accent bar
    add_shape_rect(slide, lectura_x, lectura_y, 0.06, lectura_h, brand_color)

    add_textbox(slide, lectura_x + 0.15, lectura_y + 0.05, lectura_w - 0.3, 0.2, f"Lectura {marca}", 10, True, SEC_TEXT)

    # Generate insights
    feb_bp = cons['feb']['p']
    jan_bp = cons['jan']['p']
    delta_bp = feb_bp - jan_bp
    pct_chg = (delta_bp / jan_bp * 100) if jan_bp > 0 else 0

    insights = []
    if delta_bp > 0:
        insights.append(f"• {marca} creció {'+' if delta_bp > 0 else ''}{delta_bp} compras MoM ({pct_chg:+.0f}%)")
    elif delta_bp < 0:
        insights.append(f"• {marca} cayó {delta_bp} compras MoM ({pct_chg:.0f}%)")
    else:
        insights.append(f"• {marca} se mantuvo estable ({feb_bp} compras)")

    # Top célula
    if len(cells) > 1:
        cell_purchases = [(cel, bd['celulas'][cel]['feb']['p']) for cel in cells]
        cell_purchases.sort(key=lambda x: -x[1])
        top_cel = cell_purchases[0]
        share = (top_cel[1] / feb_bp * 100) if feb_bp > 0 else 0
        insights.append(f"• {top_cel[0]} lidera con {top_cel[1]} compras ({share:.0f}% share)")

        # Most growth
        cell_growth = [(cel, bd['celulas'][cel]['feb']['p'] - bd['celulas'][cel]['jan']['p']) for cel in cells]
        cell_growth.sort(key=lambda x: -x[1])
        if cell_growth[0][1] > 0:
            insights.append(f"• Mayor crecimiento: {cell_growth[0][0]} (+{cell_growth[0][1]})")

    # Q→M change
    qm_delta = cons['feb']['rates']['Q→M%'] - cons['jan']['rates']['Q→M%']
    if abs(qm_delta) >= 2:
        direction = "mejoró" if qm_delta > 0 else "se deterioró"
        insights.append(f"• Q→M% {direction} ({fmt_delta(qm_delta, True)})")

    # Q→P
    qp_feb = cons['feb']['rates']['Q→P%']
    qp_delta = qp_feb - cons['jan']['rates']['Q→P%']
    insights.append(f"• Q→P% consolidado: {fmt_pct(qp_feb)} ({fmt_delta(qp_delta, True)})")

    insight_text = '\n'.join(insights[:6])
    add_textbox(slide, lectura_x + 0.15, lectura_y + 0.28, lectura_w - 0.3, lectura_h - 0.3, insight_text, 9, False, SEC_TEXT)

    # Footer
    add_textbox(slide, 0.6, 7.1, 10, 0.3, f"{marca} | Grupo Premier | TAS Analytics | Febrero 2026", 9, False, GRAY_TEXT)

# ── Save ────────────────────────────────────────────────────────
output_path = os.path.expanduser("~/Downloads/Grupo_Premier_Review Febrero 2026.pptx")
prs.save(output_path)
print(f"✅ Archivo generado: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   Marcas con slide detallado: {len(major_brands)}")
print(f"   Marcas menores (sin slide individual): {[b['marca'] for b in minor_brands]}")
