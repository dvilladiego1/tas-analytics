#!/usr/bin/env python3
"""
Plan de Mejora Abril 2026 — TAS Marketplace
Generates: ~/Downloads/Plan_Mejora_Abril_2026_TAS_Marketplace.pptx
"""

import csv, os
from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Paths ─────────────────────────────────────────────────────────────────
TAS_CSV = os.path.expanduser("~/Downloads/Summary MKP _ TAS & BULK - Snapshot_Diario (8).csv")
RETAIL_CSV = os.path.expanduser("~/Downloads/MX KPIS [Oficial] - DB Supply (1).csv")
OUTPUT = os.path.expanduser("~/Downloads/Plan_Mejora_Abril_2026_TAS_Marketplace.pptx")

# ── Colours ───────────────────────────────────────────────────────────────
KAVAK_BLUE   = RGBColor(0x00, 0x4E, 0x98)
KAVAK_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_HEADER = RGBColor(0x2B, 0x47, 0x8B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
RED          = RGBColor(0xEF, 0x44, 0x44)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE       = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
CARD_BG      = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_LIGHT   = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_LIGHT  = RGBColor(0xDC, 0xFC, 0xE7)
RED_LIGHT    = RGBColor(0xFE, 0xE2, 0xE2)
TAS_COLOR    = RGBColor(0x1B, 0x2A, 0x4A)
ALIADO_COLOR = RGBColor(0x3B, 0x7D, 0xDD)
RETAIL_COLOR = MED_GRAY

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Date constants ────────────────────────────────────────────────────────
HOLIDAYS = {date(2026, 3, 16)}  # Benito Juárez
MAR_MTD_END = date(2026, 3, 21)
FEB_END = date(2026, 2, 28)
MAR_BIZ_DAYS = 25
APR_BIZ_DAYS = 22  # Semana Santa Apr 13-18

MONTH_RANGES = {
    'Ene': (date(2026,1,1), date(2026,1,31)),
    'Feb': (date(2026,2,1), date(2026,2,28)),
    'Mar': (date(2026,3,1), date(2026,3,31)),
}

WEEKS = {
    'w-5': (date(2026,2,16), date(2026,2,21)),
    'w-4': (date(2026,2,23), date(2026,2,28)),
    'w-3': (date(2026,3,2),  date(2026,3,7)),
    'w-2': (date(2026,3,9),  date(2026,3,14)),
    'w-1': (date(2026,3,16), date(2026,3,21)),
}

GRUPO_MAP = {
    'GRUPO CONTINENTAL': 'CONTINENTAL', 'GRUPO PREMIER': 'PREMIER',
    'GRUPO ANDRADE': 'ANDRADE', 'GRUPO PLASENCIA': 'PLASENCIA',
    'GRUPO ISMO TLALNEPANTLA': 'ISMO', 'GRUPO GP AUTO': 'GP AUTO',
    'GRUPO AUTOPOLIS': 'AUTOPOLIS', 'GRUPO SONI': 'SONI',
    'GRUPO ISMO AGUASCALIENTES': 'ISMO', 'GRUPO TOLLOCAN': 'TOLLOCAN',
    'GRUPO POTOSINA': 'POTOSINA', 'GRUPO ISMO LEON': 'ISMO',
    'GRUPO TORRES CORZO': 'TORRES CORZO', 'GRUPO WECARS': 'WECARS',
    'GRUPO ISMO': 'ISMO', 'GRUPO MEGA': 'MEGA', 'GRUPO MISOL': 'MISOL',
    'ALIADO COAPA': 'COAPA', 'ALIADO LEON': 'LEON',
    'ALIADO TLAHUAC': 'TLAHUAC', 'ALIADO CONTINENTAL': 'CONTINENTAL',
    'ALIADO TOLLOCAN': 'TOLLOCAN', 'ALIADO MEGA': 'MEGA',
    'ALIADO TORRES CORZO': 'TORRES CORZO', 'ALIADO ANDRADE': 'ANDRADE',
    'ALIADO ISMO': 'ISMO', 'ALIADO SONI': 'SONI',
    'ALIADO POTOSINA': 'POTOSINA', 'ALIADO PREMIER': 'PREMIER',
    'ALIADO WECARS': 'WECARS', 'ALIADO MISOL': 'MISOL',
    'ALIADO AGUASCALIENTES': 'ISMO',
}
EXCLUDE_GROUPS = {'ALBACAR','AUTOKLIC','B2B','MENA','NO ES TAS | Oportunidad de Compra','Rechazo 7 días'}

# ── Helpers (date / math) ─────────────────────────────────────────────────
def parse_date(s):
    if not s or '/' not in s: return None
    parts = s.strip().split('/')
    if len(parts) != 3: return None
    try: return date(int(parts[2]), int(parts[1]), int(parts[0]))
    except: return None

def is_biz_day(d):
    return d.weekday() < 6 and d not in HOLIDAYS

def biz_days_in_range(start, end):
    count, d = 0, start
    while d <= end:
        if is_biz_day(d): count += 1
        d += timedelta(days=1)
    return count

def in_range(d, s, e): return s <= d <= e
def rate(n, d): return n/d*100 if d else 0
def fmt(v): return f'{v:,}'
def fmtR(v): return f'{v:.1f}%'
def delta_color(cur, prev, higher_is_better=True):
    if cur > prev: return GREEN if higher_is_better else RED
    if cur < prev: return RED if higher_is_better else GREEN
    return MED_GRAY
def delta_str(cur, prev):
    if prev == 0: return '+∞' if cur > 0 else '—'
    return f'{(cur-prev)/prev*100:+.1f}%'
def delta_pp(cur, prev): return f'{cur-prev:+.1f}pp'

# ── PPTX Helpers ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid(); shape.line.color.rgb = line_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid(); shape.line.color.rgb = line_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(text)
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font_name; p.alignment = alignment
    return txBox

def add_blue_bar(slide):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=KAVAK_BLUE)

def add_slide_title(slide, title_text, subtitle_text=None):
    add_blue_bar(slide)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.5),
                 title_text, font_size=24, bold=True, color=KAVAK_DARK)
    if subtitle_text:
        add_text_box(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
                     subtitle_text, font_size=12, color=MED_GRAY)

def add_kpi_card(slide, left, top, width, height, label, value, delta, delta_color_val, accent_color=KAVAK_BLUE):
    add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, line_color=LIGHT_GRAY)
    add_shape(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.06), fill_color=accent_color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), Inches(0.3),
                 label, font_size=10, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), Inches(0.5),
                 value, font_size=28, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.0), width - Inches(0.3), Inches(0.3),
                 delta, font_size=11, bold=True, color=delta_color_val, alignment=PP_ALIGN.CENTER)

def set_cell_text(cell, text, font_size=9, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]; p.text = str(text)
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = "Calibri"; p.alignment = alignment; cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def style_table_header(table, cols, fill_color=TABLE_HEADER):
    for i in range(cols):
        cell = table.cell(0, i); cell.fill.solid(); cell.fill.fore_color.rgb = fill_color
        set_cell_text(cell, cell.text, font_size=9, bold=True, color=WHITE)

def style_table_rows(table, rows, cols):
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

def add_insight_box(slide, left, top, width, height, text, accent_color=KAVAK_BLUE, bg_color=BLUE_LIGHT):
    add_rounded_rect(slide, left, top, width, height, fill_color=bg_color)
    add_shape(slide, left, top + Inches(0.1), Inches(0.06), height - Inches(0.2), fill_color=accent_color)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), width - Inches(0.35), height - Inches(0.16))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith(">> "):
            p.text = line[3:]; p.font.bold = True
        else:
            p.text = line
        p.font.size = Pt(9); p.font.color.rgb = KAVAK_DARK; p.font.name = "Calibri"

def add_footer(slide, text="Plan de Mejora Abril 2026 — TAS Marketplace"):
    add_text_box(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
                 text, font_size=8, color=MED_GRAY)

# ── LOAD TAS DATA ─────────────────────────────────────────────────────────
print("Loading TAS data...")

def load_tas_csv(path):
    rows = []
    with open(path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f); headers = next(reader)
        for row in reader:
            if len(row) < 55: continue
            grupo_raw = row[30].strip(); celula_raw = row[31].strip(); channel_raw = row[36].strip()
            if celula_raw == 'TEST-TEST': continue
            if channel_raw not in ('TAS', 'ALIADO'): continue
            if 'POTOSINA' in grupo_raw: channel_raw = 'TAS'
            if celula_raw == 'ANDRADE AEROPUERTO' and channel_raw == 'ALIADO': channel_raw = 'TAS'
            if celula_raw == 'SONI PACHUCA' and channel_raw == 'ALIADO': channel_raw = 'TAS'
            grupo_norm = GRUPO_MAP.get(grupo_raw, grupo_raw.replace('GRUPO ','').replace('ALIADO ',''))
            rows.append({
                'grupo_raw': grupo_raw, 'grupo': grupo_norm, 'celula': celula_raw,
                'channel': channel_raw,
                'fecha_creacion': parse_date(row[2]), 'purchase_date': parse_date(row[46]),
                'is_quote': row[48].strip() not in ('','0','FALSE','false') if len(row)>48 else False,
                'is_made': row[49].strip() not in ('','0','FALSE','false') if len(row)>49 else False,
                'is_approved': row[50].strip() not in ('','0','FALSE','false') if len(row)>50 else False,
                'is_purchased': row[54].strip() not in ('','0','FALSE','false') if len(row)>54 else False,
            })
    return rows

tas_records = load_tas_csv(TAS_CSV)
print(f"  TAS records: {len(tas_records)}")

# ── AGGREGATE ─────────────────────────────────────────────────────────────
def aggregate(records, channel_filter=None, grupo_filter=None, celula_filter=None):
    """Aggregate funnel metrics. Q/M/A by fecha_creacion, P by purchase_date."""
    filt = [r for r in records if r['grupo'] not in EXCLUDE_GROUPS]
    if channel_filter: filt = [r for r in filt if r['channel'] == channel_filter]
    if grupo_filter: filt = [r for r in filt if r['grupo'] == grupo_filter]
    if celula_filter: filt = [r for r in filt if r['celula'] == celula_filter]

    result = {}
    # Monthly
    for mo, (ms, me) in MONTH_RANGES.items():
        q = sum(1 for r in filt if r['is_quote'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, me))
        m = sum(1 for r in filt if r['is_made'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, me))
        a = sum(1 for r in filt if r['is_approved'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, me))
        p = sum(1 for r in filt if r['is_purchased'] and r['purchase_date'] and in_range(r['purchase_date'], ms, me))
        result[mo] = [q, m, a, p]

    # Weekly
    for wk, (ws, we) in WEEKS.items():
        q = sum(1 for r in filt if r['is_quote'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ws, we))
        m = sum(1 for r in filt if r['is_made'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ws, we))
        a = sum(1 for r in filt if r['is_approved'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ws, we))
        p = sum(1 for r in filt if r['is_purchased'] and r['purchase_date'] and in_range(r['purchase_date'], ws, we))
        result[wk] = [q, m, a, p]

    # Mar MTD (through Mar 21)
    ms = date(2026,3,1)
    q = sum(1 for r in filt if r['is_quote'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, MAR_MTD_END))
    m = sum(1 for r in filt if r['is_made'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, MAR_MTD_END))
    a = sum(1 for r in filt if r['is_approved'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], ms, MAR_MTD_END))
    p = sum(1 for r in filt if r['is_purchased'] and r['purchase_date'] and in_range(r['purchase_date'], ms, MAR_MTD_END))
    result['Mar MTD'] = [q, m, a, p]

    # Feb MTD same-period (Feb 1-20, ~17 biz days to match Mar 1-21)
    feb_mtd_end = date(2026, 2, 20)
    feb_mtd_s = date(2026, 2, 1)
    q = sum(1 for r in filt if r['is_quote'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], feb_mtd_s, feb_mtd_end))
    m = sum(1 for r in filt if r['is_made'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], feb_mtd_s, feb_mtd_end))
    a = sum(1 for r in filt if r['is_approved'] and r['fecha_creacion'] and in_range(r['fecha_creacion'], feb_mtd_s, feb_mtd_end))
    p = sum(1 for r in filt if r['is_purchased'] and r['purchase_date'] and in_range(r['purchase_date'], feb_mtd_s, feb_mtd_end))
    result['Feb MTD'] = [q, m, a, p]

    return result

# Compute all aggregates
print("Computing aggregates...")
agg_total = aggregate(tas_records)
agg_tas = aggregate(tas_records, channel_filter='TAS')
agg_ali = aggregate(tas_records, channel_filter='ALIADO')

# Per-celula aggregates
celulas = sorted(set(r['celula'] for r in tas_records if r['grupo'] not in EXCLUDE_GROUPS))
agg_by_celula = {}
for c in celulas:
    agg_by_celula[c] = aggregate(tas_records, celula_filter=c)

# Per-grupo_raw aggregates
grupos_raw = sorted(set(r['grupo_raw'] for r in tas_records if r['grupo'] not in EXCLUDE_GROUPS))
agg_by_grupo = {}
for g in grupos_raw:
    recs = [r for r in tas_records if r['grupo_raw'] == g]
    ch = recs[0]['channel'] if recs else 'TAS'
    agg_by_grupo[g] = {'agg': aggregate(tas_records, grupo_filter=GRUPO_MAP.get(g, g.replace('GRUPO ','').replace('ALIADO ',''))),
                        'channel': ch, 'grupo_raw': g}

# ── LOAD RETAIL DATA ──────────────────────────────────────────────────────
print("Loading Retail data...")

def load_retail_csv(path):
    """Parse multi-block Retail CSV. Returns dict of monthly funnel metrics."""
    import csv as csv_mod
    with open(path, 'r', encoding='utf-8') as f:
        reader = csv_mod.reader(f)
        headers = next(reader)
        data = list(reader)

    def extract_block(date_col, region_col, val_col):
        """Sum values by month for region='All'."""
        monthly = {}
        for row in data:
            if len(row) <= max(date_col, region_col, val_col): continue
            region = row[region_col].strip()
            if region != 'All': continue
            try:
                dt = date.fromisoformat(row[date_col].strip())
            except: continue
            try:
                val = float(row[val_col].strip())
            except: continue
            key = f"{dt.year}-{dt.month:02d}"
            monthly[key] = monthly.get(key, 0) + val
        return monthly

    # Block indices from header analysis
    made = extract_block(25, 26, 27)       # fecha_inspection_made, region, inspection_made
    approved = extract_block(36, 37, 38)   # fecha_inspection_approved, region, inspection_approved
    purchases = extract_block(45, 46, 47)  # fecha_compra, region, net_purchases_retail
    schedules = extract_block(14, 15, 18)  # fecha_inspection_scheduled, region, schedule_confirmed

    result = {}
    for mo_key, mo_label in [('2026-01','Ene'), ('2026-02','Feb'), ('2026-03','Mar')]:
        result[mo_label] = {
            'schedules': int(schedules.get(mo_key, 0)),
            'made': int(made.get(mo_key, 0)),
            'approved': int(approved.get(mo_key, 0)),
            'purchases': int(purchases.get(mo_key, 0)),
        }
    return result

retail = load_retail_csv(RETAIL_CSV)
print(f"  Retail Mar: Made={retail['Mar']['made']}, App={retail['Mar']['approved']}, Pur={retail['Mar']['purchases']}")

# ── RUN RATE ──────────────────────────────────────────────────────────────
w1_data = agg_total['w-1']  # Mar 16-21
w1_biz = biz_days_in_range(date(2026,3,16), date(2026,3,21))
daily_rate_p = w1_data[3] / w1_biz if w1_biz else 0
mar_projected = round(daily_rate_p * MAR_BIZ_DAYS)
apr_baseline = round(daily_rate_p * APR_BIZ_DAYS)
apr_target = 470  # User input: March closes 470-500

print(f"  RR 7d: {daily_rate_p:.1f}/day, Mar projected: {mar_projected}, Apr baseline (22d): {apr_baseline}")

# ── CELL-LEVEL DATA FOR RANKINGS & IMPACT ─────────────────────────────────
print("Computing cell-level metrics...")

cell_data = []
for c in celulas:
    a = agg_by_celula[c]
    mar = a.get('Mar MTD', [0,0,0,0])
    feb = a.get('Feb', [0,0,0,0])
    q, m, ap, p = mar
    # Find channel for this celula
    ch = 'TAS'
    for r in tas_records:
        if r['celula'] == c and r['grupo'] not in EXCLUDE_GROUPS:
            ch = r['channel']; break
    cell_data.append({
        'celula': c, 'channel': ch,
        'q': q, 'm': m, 'a': ap, 'p': p,
        'q_feb': feb[0], 'm_feb': feb[1], 'a_feb': feb[2], 'p_feb': feb[3],
        'qm': rate(m, q), 'ma': rate(ap, m), 'ap_pct': rate(p, ap), 'e2e': rate(p, q),
        'qm_feb': rate(feb[1], feb[0]), 'ap_feb': rate(feb[3], feb[2]), 'e2e_feb': rate(feb[3], feb[0]),
    })

# ══════════════════════════════════════════════════════════════════════════
# SLIDE GENERATION
# ══════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────
print("Generating Slide 1: Title...")
s1 = prs.slides.add_slide(blank_layout)
add_shape(s1, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=KAVAK_DARK)
add_shape(s1, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), fill_color=KAVAK_BLUE)
add_text_box(s1, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Plan de Mejora Abril 2026", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_text_box(s1, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "TAS Marketplace — Mantener volumen + mejorar E2E +15% con 22 días hábiles",
             font_size=18, color=RGBColor(0xCB,0xD5,0xE1), alignment=PP_ALIGN.LEFT)
add_text_box(s1, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Abril: 22 días hábiles (Semana Santa 13-18 Abr) | Target: 470+ compras | E2E: 8.8% → 10.1%",
             font_size=13, color=MED_GRAY, alignment=PP_ALIGN.LEFT)
add_text_box(s1, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "24 de Marzo 2026", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

# ── SLIDE 2: SUMMARY Mar MTD vs Feb MTD (same period) ────────────────────
print("Generating Slide 2: Summary...")
s2 = prs.slides.add_slide(blank_layout)
add_slide_title(s2, "Resumen Marzo MTD vs Febrero MTD", "Mar 1-21 vs Feb 1-20 (mismo período ~17 días hábiles) | Run Rate últimos 7 días")

mar = agg_total['Mar MTD']
feb_mtd = agg_total['Feb MTD']

# Volume KPI cards (row 1) — using Feb MTD same-period
cards_vol = [
    ("Quotes", fmt(mar[0]), delta_str(mar[0], feb_mtd[0]), delta_color(mar[0], feb_mtd[0])),
    ("Made", fmt(mar[1]), delta_str(mar[1], feb_mtd[1]), delta_color(mar[1], feb_mtd[1])),
    ("Approved", fmt(mar[2]), delta_str(mar[2], feb_mtd[2]), delta_color(mar[2], feb_mtd[2])),
    ("Purchases", fmt(mar[3]), delta_str(mar[3], feb_mtd[3]), delta_color(mar[3], feb_mtd[3])),
]
for i, (lbl, val, dlt, dc) in enumerate(cards_vol):
    add_kpi_card(s2, Inches(0.5 + i*3.1), Inches(1.1), Inches(2.9), Inches(1.3), lbl, val, dlt, dc)

# Conversion KPI cards (row 2) — using Feb MTD same-period
mar_qm, mar_ma, mar_ap, mar_e2e = rate(mar[1],mar[0]), rate(mar[2],mar[1]), rate(mar[3],mar[2]), rate(mar[3],mar[0])
feb_mtd_qm = rate(feb_mtd[1],feb_mtd[0])
feb_mtd_ma = rate(feb_mtd[2],feb_mtd[1])
feb_mtd_ap = rate(feb_mtd[3],feb_mtd[2])
feb_mtd_e2e = rate(feb_mtd[3],feb_mtd[0])

cards_conv = [
    ("Q→M%", fmtR(mar_qm), delta_pp(mar_qm, feb_mtd_qm), delta_color(mar_qm, feb_mtd_qm)),
    ("M→A%", fmtR(mar_ma), delta_pp(mar_ma, feb_mtd_ma), delta_color(mar_ma, feb_mtd_ma)),
    ("A→P%", fmtR(mar_ap), delta_pp(mar_ap, feb_mtd_ap), delta_color(mar_ap, feb_mtd_ap)),
    ("E2E Q→P%", fmtR(mar_e2e), delta_pp(mar_e2e, feb_mtd_e2e), delta_color(mar_e2e, feb_mtd_e2e)),
]
for i, (lbl, val, dlt, dc) in enumerate(cards_conv):
    add_kpi_card(s2, Inches(0.5 + i*3.1), Inches(2.55), Inches(2.9), Inches(1.3), lbl, val, dlt, dc)

# Row 3 left: Mini-tabla TAS Only vs Aliado
t_mar = agg_tas['Mar MTD']
a_mar = agg_ali['Mar MTD']
ta_headers = ["Métrica", "TAS Only", "Aliado", "Δ"]
ta_data = [
    ["Quotes", fmt(t_mar[0]), fmt(a_mar[0]), f"{t_mar[0]/(t_mar[0]+a_mar[0])*100:.0f}% / {a_mar[0]/(t_mar[0]+a_mar[0])*100:.0f}%"],
    ["Made", fmt(t_mar[1]), fmt(a_mar[1]), ""],
    ["Approved", fmt(t_mar[2]), fmt(a_mar[2]), ""],
    ["Purchases", fmt(t_mar[3]), fmt(a_mar[3]), f"{t_mar[3]/(t_mar[3]+a_mar[3])*100:.0f}% / {a_mar[3]/(t_mar[3]+a_mar[3])*100:.0f}%"],
    ["Q→M%", fmtR(rate(t_mar[1],t_mar[0])), fmtR(rate(a_mar[1],a_mar[0])), delta_pp(rate(t_mar[1],t_mar[0]), rate(a_mar[1],a_mar[0]))],
    ["M→A%", fmtR(rate(t_mar[2],t_mar[1])), fmtR(rate(a_mar[2],a_mar[1])), delta_pp(rate(t_mar[2],t_mar[1]), rate(a_mar[2],a_mar[1]))],
    ["A→P%", fmtR(rate(t_mar[3],t_mar[2])), fmtR(rate(a_mar[3],a_mar[2])), delta_pp(rate(t_mar[3],t_mar[2]), rate(a_mar[3],a_mar[2]))],
    ["E2E Q→P%", fmtR(rate(t_mar[3],t_mar[0])), fmtR(rate(a_mar[3],a_mar[0])), delta_pp(rate(t_mar[3],t_mar[0]), rate(a_mar[3],a_mar[0]))],
]
add_text_box(s2, Inches(0.5), Inches(3.95), Inches(6), Inches(0.25),
             "TAS Only vs Aliado (Mar MTD)", font_size=11, bold=True, color=KAVAK_DARK)
ta_tbl = s2.shapes.add_table(len(ta_data)+1, 4, Inches(0.5), Inches(4.2), Inches(6), Inches(0.26)*(len(ta_data)+1)).table
for ci, h in enumerate(ta_headers):
    set_cell_text(ta_tbl.cell(0, ci), h, font_size=8, bold=True, color=WHITE)
style_table_header(ta_tbl, 4)
for ri, row in enumerate(ta_data):
    for ci, val in enumerate(row):
        bold = (ci == 0 or ri == 3 or ri == 7)
        color = DARK_GRAY
        if ci == 3 and 'pp' in val:
            color = GREEN if val.startswith('+') else RED
        set_cell_text(ta_tbl.cell(ri+1, ci), val, font_size=8, bold=bold, color=color)
style_table_rows(ta_tbl, len(ta_data)+1, 4)

# Row 3 right: Run Rate insight box (compact)
rr_text = (f">> Run Rate últimos 7 días (Mar 16-21)\n"
           f"Compras w-1: {w1_data[3]} en {w1_biz} días = {daily_rate_p:.1f}/día\n"
           f"Proyección Mar (25d): ~{mar_projected} | Abr (22d): ~{apr_baseline}\n"
           f">> Target Abril: 470+ = {470/APR_BIZ_DAYS:.1f}/día\n"
           f"Mejora E2E +15% compensa Semana Santa.")
add_insight_box(s2, Inches(6.8), Inches(4.2), Inches(6), Inches(2.2), rr_text)
add_footer(s2)

# ── SLIDE 3: THREE FUNNELS ────────────────────────────────────────────────
print("Generating Slide 3: Three Funnels...")
s3 = prs.slides.add_slide(blank_layout)
add_slide_title(s3, "Comparativo de Funnels: Retail vs TAS vs Aliado", "Marzo MTD (1-21)")

# Retail funnel data (including Schedules = equivalent to TAS Quotes)
r_sch = retail['Mar']['schedules']
r_made = retail['Mar']['made']
r_app = retail['Mar']['approved']
r_pur = retail['Mar']['purchases']
r_sm = rate(r_made, r_sch)  # S→M% = equivalent to TAS Q→M%
r_ma = rate(r_app, r_made)
r_ap = rate(r_pur, r_app)
r_e2e = rate(r_pur, r_sch)

# TAS data
t = agg_tas['Mar MTD']
t_qm, t_ma, t_ap, t_e2e = rate(t[1],t[0]), rate(t[2],t[1]), rate(t[3],t[2]), rate(t[3],t[0])

# Aliado data
al = agg_ali['Mar MTD']
a_qm, a_ma, a_ap, a_e2e = rate(al[1],al[0]), rate(al[2],al[1]), rate(al[3],al[2]), rate(al[3],al[0])

# Draw funnel columns
def draw_funnel_col(slide, left, label, color, stages, conv_rates):
    """Draw a vertical funnel column."""
    top_start = Inches(1.3)
    col_w = Inches(3.5)
    add_rounded_rect(slide, left, top_start - Inches(0.1), col_w, Inches(0.45), fill_color=color)
    add_text_box(slide, left, top_start - Inches(0.05), col_w, Inches(0.35),
                 label, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    for i, (stage_name, vol) in enumerate(stages):
        y = top_start + Inches(0.5) + i * Inches(0.85)
        # Stage box
        add_rounded_rect(slide, left + Inches(0.1), y, col_w - Inches(0.2), Inches(0.6),
                        fill_color=WHITE, line_color=LIGHT_GRAY)
        add_text_box(slide, left + Inches(0.2), y + Inches(0.05), Inches(1.5), Inches(0.25),
                     stage_name, font_size=9, bold=True, color=MED_GRAY)
        add_text_box(slide, left + Inches(1.8), y + Inches(0.05), Inches(1.4), Inches(0.25),
                     fmt(vol), font_size=14, bold=True, color=KAVAK_DARK, alignment=PP_ALIGN.RIGHT)
        # Conversion rate arrow
        if i < len(conv_rates):
            add_text_box(slide, left + Inches(0.2), y + Inches(0.32), col_w - Inches(0.4), Inches(0.22),
                         f"↓ {fmtR(conv_rates[i])}", font_size=9, bold=True,
                         color=GREEN if conv_rates[i] >= 50 else (AMBER if conv_rates[i] >= 30 else RED),
                         alignment=PP_ALIGN.CENTER)

# Retail column (now includes Schedules as equivalent to Quotes)
draw_funnel_col(s3, Inches(0.5), "RETAIL", RETAIL_COLOR,
    [("Schedules", r_sch), ("Made", r_made), ("Approved", r_app), ("Purchases", r_pur)],
    [r_sm, r_ma, r_ap])

# TAS Only column
draw_funnel_col(s3, Inches(4.5), "TAS ONLY", TAS_COLOR,
    [("Quotes", t[0]), ("Made", t[1]), ("Approved", t[2]), ("Purchases", t[3])],
    [t_qm, t_ma, t_ap])

# Aliado column
draw_funnel_col(s3, Inches(8.8), "ALIADO", ALIADO_COLOR,
    [("Quotes", al[0]), ("Made", al[1]), ("Approved", al[2]), ("Purchases", al[3])],
    [a_qm, a_ma, a_ap])

# Gap comparison table
gap_top = Inches(5.3)
gap_data = [
    ["Métrica", "Retail", "TAS Only", "Aliado", "Gap TAS vs Retail"],
    ["S→M% / Q→M%", fmtR(r_sm), fmtR(t_qm), fmtR(a_qm), delta_pp(t_qm, r_sm)],
    ["M→A%", fmtR(r_ma), fmtR(t_ma), fmtR(a_ma), delta_pp(t_ma, r_ma)],
    ["A→P%", fmtR(r_ap), fmtR(t_ap), fmtR(a_ap), delta_pp(t_ap, r_ap)],
    ["E2E", fmtR(r_e2e), fmtR(t_e2e), fmtR(a_e2e), delta_pp(t_e2e, r_e2e)],
]
tbl = s3.shapes.add_table(len(gap_data), 5, Inches(0.5), gap_top, Inches(12.3), Inches(1.2)).table
for ri, row_data in enumerate(gap_data):
    for ci, val in enumerate(row_data):
        set_cell_text(tbl.cell(ri, ci), val, font_size=10, bold=(ri==0),
                     color=WHITE if ri==0 else DARK_GRAY)
style_table_header(tbl, 5)
style_table_rows(tbl, len(gap_data), 5)
# Color the gap column
for ri in range(1, len(gap_data)):
    cell = tbl.cell(ri, 4)
    val_text = gap_data[ri][4]
    if '+' in val_text:
        set_cell_text(cell, val_text, font_size=10, bold=True, color=GREEN)
    elif '-' in val_text or 'target' in val_text:
        set_cell_text(cell, val_text, font_size=10, bold=True, color=RED)
add_footer(s3)

# ── SLIDE 4: TOP CELLS RANKING ───────────────────────────────────────────
print("Generating Slide 4: Top Cells Ranking...")
s4 = prs.slides.add_slide(blank_layout)
add_slide_title(s4, "Qué salió bien en Marzo", "Rankings por Volumen y por Conversión")

# Table A: Top by volume (purchases)
vol_sorted = sorted([c for c in cell_data if c['p'] > 0], key=lambda x: -x['p'])[:10]
headers_a = ["#", "Célula", "Canal", "P(Mar)", "P(Feb)", "Δ", "A→P%", "Q→P%"]
tbl_a = s4.shapes.add_table(len(vol_sorted)+1, len(headers_a), Inches(0.3), Inches(1.1), Inches(6.2), Inches(0.3)*(len(vol_sorted)+1)).table
for ci, h in enumerate(headers_a):
    set_cell_text(tbl_a.cell(0, ci), h, font_size=9, bold=True, color=WHITE)
style_table_header(tbl_a, len(headers_a))
for ri, c in enumerate(vol_sorted):
    row_i = ri + 1
    p_delta = delta_str(c['p'], c['p_feb'])
    set_cell_text(tbl_a.cell(row_i, 0), str(ri+1), font_size=9)
    set_cell_text(tbl_a.cell(row_i, 1), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT)
    set_cell_text(tbl_a.cell(row_i, 2), c['channel'], font_size=8)
    set_cell_text(tbl_a.cell(row_i, 3), str(c['p']), font_size=9, bold=True)
    set_cell_text(tbl_a.cell(row_i, 4), str(c['p_feb']), font_size=9)
    set_cell_text(tbl_a.cell(row_i, 5), p_delta, font_size=9, color=GREEN if c['p']>=c['p_feb'] else RED)
    set_cell_text(tbl_a.cell(row_i, 6), fmtR(c['ap_pct']), font_size=9)
    set_cell_text(tbl_a.cell(row_i, 7), fmtR(c['e2e']), font_size=9)
style_table_rows(tbl_a, len(vol_sorted)+1, len(headers_a))

add_text_box(s4, Inches(0.3), Inches(1.0) - Inches(0.15), Inches(6), Inches(0.2),
             "Top 10 por Volumen (Compras)", font_size=11, bold=True, color=KAVAK_DARK)

# Table B: Top by conversion (E2E) — filter Made >= 20
conv_filtered = [c for c in cell_data if c['m'] >= 20]
conv_sorted = sorted(conv_filtered, key=lambda x: -x['e2e'])[:10]
headers_b = ["#", "Célula", "Canal", "Made", "Q→M%", "M→A%", "A→P%", "Q→P%"]
tbl_b = s4.shapes.add_table(len(conv_sorted)+1, len(headers_b), Inches(6.8), Inches(1.1), Inches(6.2), Inches(0.3)*(len(conv_sorted)+1)).table
for ci, h in enumerate(headers_b):
    set_cell_text(tbl_b.cell(0, ci), h, font_size=9, bold=True, color=WHITE)
style_table_header(tbl_b, len(headers_b))
for ri, c in enumerate(conv_sorted):
    row_i = ri + 1
    set_cell_text(tbl_b.cell(row_i, 0), str(ri+1), font_size=9)
    set_cell_text(tbl_b.cell(row_i, 1), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT)
    set_cell_text(tbl_b.cell(row_i, 2), c['channel'], font_size=8)
    set_cell_text(tbl_b.cell(row_i, 3), str(c['m']), font_size=9)
    set_cell_text(tbl_b.cell(row_i, 4), fmtR(c['qm']), font_size=9, color=GREEN if c['qm']>30 else DARK_GRAY)
    set_cell_text(tbl_b.cell(row_i, 5), fmtR(c['ma']), font_size=9, color=GREEN if c['ma']>85 else DARK_GRAY)
    set_cell_text(tbl_b.cell(row_i, 6), fmtR(c['ap_pct']), font_size=9, color=GREEN if c['ap_pct']>50 else DARK_GRAY)
    set_cell_text(tbl_b.cell(row_i, 7), fmtR(c['e2e']), font_size=9, bold=True, color=GREEN if c['e2e']>10 else DARK_GRAY)
style_table_rows(tbl_b, len(conv_sorted)+1, len(headers_b))

add_text_box(s4, Inches(6.8), Inches(1.0) - Inches(0.15), Inches(6), Inches(0.2),
             "Top 10 por Conversión (E2E%, ≥20 Made)", font_size=11, bold=True, color=KAVAK_DARK)

# Insight box
top3_vol = ', '.join(c['celula'] for c in vol_sorted[:3])
top3_conv = ', '.join(c['celula'] for c in conv_sorted[:3])
insight_4 = (f">> Top 3 volumen: {top3_vol}\n"
             f">> Top 3 conversión: {top3_conv}\n"
             f"Células con alta conversión Y volumen son las mejores prácticas a replicar.")
add_insight_box(s4, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.2), insight_4, accent_color=GREEN, bg_color=GREEN_LIGHT)
add_footer(s4)

# ── SLIDE 5: Q→M% IMPACT MATRIX ──────────────────────────────────────────
print("Generating Slide 5: Q→M% Impact Matrix...")
s5 = prs.slides.add_slide(blank_layout)
total_mar = agg_total['Mar MTD']
total_qm = rate(total_mar[1], total_mar[0])
add_slide_title(s5, f"Plan Q→M%: Matriz de Impacto por Célula",
                f"Actual: {fmtR(total_qm)} | Target Abril: {fmtR(total_qm*1.15)} | Gap vs Retail (39.1%): {delta_pp(total_qm, 39.1)}")

# Impact calculation: if cell improves Q→M by 5pp, how many incremental purchases?
qm_impact = []
for c in cell_data:
    if c['q'] < 10: continue  # skip tiny cells
    incr_made = c['q'] * 0.05  # 5pp improvement
    if c['m'] > 0 and c['a'] > 0:
        incr_p = incr_made * (c['a']/c['m']) * (c['p']/c['a']) if c['a'] > 0 else 0
    else:
        incr_p = 0
    qm_impact.append({**c, 'incr_made': round(incr_made, 1), 'incr_p': round(incr_p, 1),
                       'gap_retail': 39.1 - c['qm'],
                       'priority': 'CRITICA' if c['qm'] < 18 else ('ALTA' if c['qm'] < 25 else ('MEDIA' if c['qm'] < 32 else 'OK'))})
qm_impact.sort(key=lambda x: -x['incr_p'])

# Table
headers_5 = ["Célula", "Canal", "Quotes", "Q→M%", "Gap vs Retail", "+Made (+5pp)", "+Compras Est.", "Prioridad"]
rows_5 = qm_impact[:12]
tbl5 = s5.shapes.add_table(len(rows_5)+1, len(headers_5), Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.3)*(len(rows_5)+1)).table
for ci, h in enumerate(headers_5):
    set_cell_text(tbl5.cell(0, ci), h, font_size=9, bold=True, color=WHITE)
style_table_header(tbl5, len(headers_5))
for ri, c in enumerate(rows_5):
    r = ri + 1
    prio_color = RED if c['priority']=='CRITICA' else (AMBER if c['priority']=='ALTA' else (RGBColor(0xD9,0x77,0x06) if c['priority']=='MEDIA' else GREEN))
    set_cell_text(tbl5.cell(r, 0), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT)
    set_cell_text(tbl5.cell(r, 1), c['channel'], font_size=8)
    set_cell_text(tbl5.cell(r, 2), str(c['q']), font_size=9)
    set_cell_text(tbl5.cell(r, 3), fmtR(c['qm']), font_size=9, color=RED if c['qm']<25 else DARK_GRAY)
    set_cell_text(tbl5.cell(r, 4), f"{c['gap_retail']:+.1f}pp", font_size=9, color=RED)
    set_cell_text(tbl5.cell(r, 5), f"+{c['incr_made']:.0f}", font_size=9, bold=True, color=GREEN)
    set_cell_text(tbl5.cell(r, 6), f"+{c['incr_p']:.1f}", font_size=9, bold=True, color=GREEN)
    set_cell_text(tbl5.cell(r, 7), c['priority'], font_size=9, bold=True, color=prio_color)
style_table_rows(tbl5, len(rows_5)+1, len(headers_5))

total_incr_p_qm = sum(c['incr_p'] for c in qm_impact)
insight_5 = (f">> Impacto total si todas las células mejoran Q→M% en +5pp: +{total_incr_p_qm:.0f} compras incrementales\n"
             f"Palancas: Confirmación 24h, bots VAPI, reducción no-shows, pre-screening de leads\n"
             f"Foco en células CRITICA/ALTA que concentran el mayor volumen de quotes")
add_insight_box(s5, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.3), insight_5)
add_footer(s5)

# ── SLIDE 6: A→P% IMPACT MATRIX ──────────────────────────────────────────
print("Generating Slide 6: A→P% Impact Matrix...")
s6 = prs.slides.add_slide(blank_layout)
total_ap = rate(total_mar[3], total_mar[2])
add_slide_title(s6, f"Plan A→P%: Matriz de Impacto por Célula",
                f"Actual: {fmtR(total_ap)} | Target Abril: {fmtR(total_ap*1.15)} | Gap vs Retail (59.3%): {delta_pp(total_ap, 59.3)}")

# Impact: if cell improves A→P by 5pp, incremental purchases = Approved * 0.05
ap_impact = []
for c in cell_data:
    if c['a'] < 5: continue
    incr_p = c['a'] * 0.05
    ap_impact.append({**c, 'incr_p': round(incr_p, 1),
                       'gap_retail': 59.3 - c['ap_pct'],
                       'priority': 'CRITICA' if c['ap_pct'] < 25 else ('ALTA' if c['ap_pct'] < 40 else ('MEDIA' if c['ap_pct'] < 50 else 'OK'))})
ap_impact.sort(key=lambda x: -x['incr_p'])

headers_6 = ["Célula", "Canal", "Approved", "A→P%", "Gap vs Retail", "+Compras (+5pp)", "Perdidos", "Prioridad"]
rows_6 = ap_impact[:12]
tbl6 = s6.shapes.add_table(len(rows_6)+1, len(headers_6), Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.3)*(len(rows_6)+1)).table
for ci, h in enumerate(headers_6):
    set_cell_text(tbl6.cell(0, ci), h, font_size=9, bold=True, color=WHITE)
style_table_header(tbl6, len(headers_6))
for ri, c in enumerate(rows_6):
    r = ri + 1
    prio_color = RED if c['priority']=='CRITICA' else (AMBER if c['priority']=='ALTA' else (RGBColor(0xD9,0x77,0x06) if c['priority']=='MEDIA' else GREEN))
    lost = c['a'] - c['p']
    set_cell_text(tbl6.cell(r, 0), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT)
    set_cell_text(tbl6.cell(r, 1), c['channel'], font_size=8)
    set_cell_text(tbl6.cell(r, 2), str(c['a']), font_size=9)
    set_cell_text(tbl6.cell(r, 3), fmtR(c['ap_pct']), font_size=9, color=RED if c['ap_pct']<40 else DARK_GRAY)
    set_cell_text(tbl6.cell(r, 4), f"{c['gap_retail']:+.1f}pp", font_size=9, color=RED)
    set_cell_text(tbl6.cell(r, 5), f"+{c['incr_p']:.1f}", font_size=9, bold=True, color=GREEN)
    set_cell_text(tbl6.cell(r, 6), str(lost), font_size=9, color=RED if lost > 20 else DARK_GRAY)
    set_cell_text(tbl6.cell(r, 7), c['priority'], font_size=9, bold=True, color=prio_color)
style_table_rows(tbl6, len(rows_6)+1, len(headers_6))

total_incr_p_ap = sum(c['incr_p'] for c in ap_impact)
criticos = [c for c in ap_impact if c['priority'] == 'CRITICA']
crit_names = ', '.join(c['celula'] for c in criticos[:5])
insight_6 = (f">> Impacto total si +5pp A→P% en todas las células: +{total_incr_p_ap:.0f} compras incrementales\n"
             f"Células CRITICAS (A→P <25%): {crit_names}\n"
             f"Palancas: Velocidad de cierre (<48h), pricing competitivo, capacitación supply agents")
add_insight_box(s6, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.3), insight_6, accent_color=RED, bg_color=RED_LIGHT)
add_footer(s6)

# ── SLIDE 7: COMBINED TARGETS ─────────────────────────────────────────────
print("Generating Slide 7: Combined Targets...")
s7 = prs.slides.add_slide(blank_layout)
add_slide_title(s7, "Targets Consolidados Abril 2026",
                f"Objetivo: mantener 470+ compras con 22 días hábiles | Mejora E2E +15%")

# Summary table
summary_data = [
    ["Métrica", "Mar Actual", "Target Abr", "Δ Necesario", "Palanca Principal"],
    ["Q→M%", fmtR(total_qm), fmtR(total_qm*1.15), f"+{total_qm*0.15:.1f}pp", "No-shows, contactabilidad, VAPI"],
    ["M→A%", fmtR(rate(total_mar[2],total_mar[1])), "Mantener", "—", "Ya es fortaleza (+15pp vs Retail)"],
    ["A→P%", fmtR(total_ap), fmtR(total_ap*1.15), f"+{total_ap*0.15:.1f}pp", "Cierre comercial, pricing, agents"],
    ["E2E Q→P%", fmtR(rate(total_mar[3],total_mar[0])), fmtR(rate(total_mar[3],total_mar[0])*1.15), "+15%", "Resultado de Q→M + A→P"],
    ["Compras/mes", fmt(total_mar[3])+" MTD", "470+", f"+{470-total_mar[3]}", "Eficiencia + volumen quotes"],
    ["Compras/día", f"{daily_rate_p:.1f}", f"{470/APR_BIZ_DAYS:.1f}", f"+{470/APR_BIZ_DAYS-daily_rate_p:.1f}", "22 días hábiles (Semana Santa)"],
]
tbl7 = s7.shapes.add_table(len(summary_data), 5, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.3)*len(summary_data)).table
for ri, row in enumerate(summary_data):
    for ci, val in enumerate(row):
        set_cell_text(tbl7.cell(ri, ci), val, font_size=10, bold=(ri==0 or ci==0),
                     color=WHITE if ri==0 else DARK_GRAY)
style_table_header(tbl7, 5)
style_table_rows(tbl7, len(summary_data), 5)

# Per-cell targets (top 10 by opportunity)
combined = {}
for c in qm_impact:
    combined[c['celula']] = {'celula': c['celula'], 'channel': c['channel'],
                              'qm': c['qm'], 'ap': c['ap_pct'], 'q': c['q'], 'a': c['a'],
                              'qm_target': min(c['qm']+5, 39.1),
                              'ap_target': min(c['ap_pct']+5, 59.3),
                              'incr_qm': c['incr_p']}
for c in ap_impact:
    if c['celula'] in combined:
        combined[c['celula']]['incr_ap'] = c['incr_p']
        combined[c['celula']]['total_incr'] = combined[c['celula']].get('incr_qm',0) + c['incr_p']
    else:
        combined[c['celula']] = {'celula': c['celula'], 'channel': c['channel'],
                                  'qm': c.get('qm',0), 'ap': c['ap_pct'], 'q': c.get('q',0), 'a': c['a'],
                                  'incr_ap': c['incr_p'], 'incr_qm': 0,
                                  'total_incr': c['incr_p'],
                                  'qm_target': c.get('qm',0)+5, 'ap_target': min(c['ap_pct']+5, 59.3)}

top_combined = sorted(combined.values(), key=lambda x: -x.get('total_incr',0))[:10]
headers_7b = ["Célula", "Canal", "Q→M% Actual", "Q→M% Target", "A→P% Actual", "A→P% Target", "+Compras Est."]
tbl7b = s7.shapes.add_table(len(top_combined)+1, len(headers_7b), Inches(0.3), Inches(3.6), Inches(12.7), Inches(0.28)*(len(top_combined)+1)).table
for ci, h in enumerate(headers_7b):
    set_cell_text(tbl7b.cell(0, ci), h, font_size=9, bold=True, color=WHITE)
style_table_header(tbl7b, len(headers_7b))
for ri, c in enumerate(top_combined):
    r = ri + 1
    set_cell_text(tbl7b.cell(r, 0), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT)
    set_cell_text(tbl7b.cell(r, 1), c['channel'], font_size=8)
    set_cell_text(tbl7b.cell(r, 2), fmtR(c['qm']), font_size=9)
    set_cell_text(tbl7b.cell(r, 3), fmtR(c.get('qm_target', c['qm']+5)), font_size=9, bold=True, color=KAVAK_BLUE)
    set_cell_text(tbl7b.cell(r, 4), fmtR(c['ap']), font_size=9)
    set_cell_text(tbl7b.cell(r, 5), fmtR(c.get('ap_target', c['ap']+5)), font_size=9, bold=True, color=KAVAK_BLUE)
    set_cell_text(tbl7b.cell(r, 6), f"+{c.get('total_incr',0):.1f}", font_size=9, bold=True, color=GREEN)
style_table_rows(tbl7b, len(top_combined)+1, len(headers_7b))

total_combined = sum(c.get('total_incr',0) for c in combined.values())
insight_7 = (f">> Impacto combinado Q→M% + A→P% (+5pp cada uno): +{total_combined:.0f} compras incrementales/mes\n"
             f"Con {fmt(total_mar[3])} compras MTD (21 días) + mejora de eficiencia → 470+ compras alcanzable en Abril")
add_insight_box(s7, Inches(0.3), Inches(6.2), Inches(12.7), Inches(1.0), insight_7)
add_footer(s7)

# ── SLIDE 8: SCORECARD + NEXT STEPS ──────────────────────────────────────
print("Generating Slide 8: Scorecard + Next Steps...")
s8 = prs.slides.add_slide(blank_layout)
add_shape(s8, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=KAVAK_DARK)
add_shape(s8, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=KAVAK_BLUE)
add_text_box(s8, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5),
             "Scorecard Semanal + Acciones Prioritarias Abril", font_size=24, bold=True, color=WHITE)

# Scorecard table
score_cells = sorted([c for c in cell_data if c['m'] >= 10], key=lambda x: -x['p'])[:10]
sc_headers = ["Grupo/Célula", "Q→M%", "M→A%", "A→P%", "E2E%", "Semáforo"]
sc_tbl = s8.shapes.add_table(len(score_cells)+1, len(sc_headers), Inches(0.4), Inches(0.85), Inches(6), Inches(0.28)*(len(score_cells)+1)).table
for ci, h in enumerate(sc_headers):
    set_cell_text(sc_tbl.cell(0, ci), h, font_size=8, bold=True, color=WHITE)
style_table_header(sc_tbl, len(sc_headers), fill_color=RGBColor(0x2B,0x47,0x8B))

target_qm = total_qm * 1.15
target_ap_val = total_ap * 1.15
target_e2e = rate(total_mar[3], total_mar[0]) * 1.15

for ri, c in enumerate(score_cells):
    r = ri + 1
    # Semaphore: green if all above target, red if any below target-2pp, yellow otherwise
    if c['qm'] >= target_qm and c['ap_pct'] >= target_ap_val:
        sem = "🟢"
    elif c['qm'] < target_qm - 5 or c['ap_pct'] < target_ap_val - 5:
        sem = "🔴"
    else:
        sem = "🟡"
    set_cell_text(sc_tbl.cell(r, 0), c['celula'], font_size=8, alignment=PP_ALIGN.LEFT, color=WHITE)
    set_cell_text(sc_tbl.cell(r, 1), fmtR(c['qm']), font_size=8, color=GREEN if c['qm']>=target_qm else (RED if c['qm']<target_qm-5 else AMBER))
    set_cell_text(sc_tbl.cell(r, 2), fmtR(c['ma']), font_size=8, color=GREEN if c['ma']>=80 else RED)
    set_cell_text(sc_tbl.cell(r, 3), fmtR(c['ap_pct']), font_size=8, color=GREEN if c['ap_pct']>=target_ap_val else (RED if c['ap_pct']<target_ap_val-5 else AMBER))
    set_cell_text(sc_tbl.cell(r, 4), fmtR(c['e2e']), font_size=8, color=GREEN if c['e2e']>=target_e2e else DARK_GRAY)
    set_cell_text(sc_tbl.cell(r, 5), sem, font_size=10)
    # Dark background for rows
    for ci2 in range(len(sc_headers)):
        sc_tbl.cell(r, ci2).fill.solid()
        sc_tbl.cell(r, ci2).fill.fore_color.rgb = RGBColor(0x22,0x22,0x3A) if r%2==0 else RGBColor(0x1E,0x1E,0x34)

# 5 Priority Actions
actions = [
    ("1", "Confirmación 24h + bot VAPI en Continental, Ismo, GP Auto", "+2pp Q→M", "Ops + AI"),
    ("2", "Intervención directa: Wecars, Potosina, Soni, Andrade (A→P <25%)", "+54 uds/mes", "Supply Lead"),
    ("3", "Subir MKD inicial en Andrade (12%) y Continental CDMX (13.1%)", "+3pp A→P", "Pricing"),
    ("4", "Auditoría Aliados con 0% A→P (Hermosillo, Mazatlán, Cuautitlán)", "+15 uds/mes", "Aliado Lead"),
    ("5", "Scorecard semanal con semáforos y escalamiento automático", "Visibilidad", "Analytics"),
]
act_headers = ["#", "Acción", "Impacto Est.", "Responsable"]
act_tbl = s8.shapes.add_table(len(actions)+1, 4, Inches(6.8), Inches(0.85), Inches(6.2), Inches(0.32)*(len(actions)+1)).table
for ci, h in enumerate(act_headers):
    set_cell_text(act_tbl.cell(0, ci), h, font_size=8, bold=True, color=WHITE)
style_table_header(act_tbl, 4, fill_color=KAVAK_BLUE)
for ri, (num, action, impact, resp) in enumerate(actions):
    r = ri + 1
    set_cell_text(act_tbl.cell(r, 0), num, font_size=9, bold=True, color=WHITE)
    set_cell_text(act_tbl.cell(r, 1), action, font_size=8, alignment=PP_ALIGN.LEFT, color=WHITE)
    set_cell_text(act_tbl.cell(r, 2), impact, font_size=9, bold=True, color=GREEN)
    set_cell_text(act_tbl.cell(r, 3), resp, font_size=8, color=RGBColor(0xCB,0xD5,0xE1))
    for ci2 in range(4):
        act_tbl.cell(r, ci2).fill.solid()
        act_tbl.cell(r, ci2).fill.fore_color.rgb = RGBColor(0x22,0x22,0x3A) if r%2==0 else RGBColor(0x1E,0x1E,0x34)

# Checkpoint
add_text_box(s8, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             "Target Abril: 470+ compras | E2E: 10.1% | Revisión S1: 6-11 Abril 2026",
             font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Semaphore legend
add_text_box(s8, Inches(0.6), Inches(6.5), Inches(12), Inches(0.3),
             "🟢 Above target  |  🟡 Within 5pp  |  🔴 Below target >5pp  |  Grupo en 🔴 2 semanas → plan correctivo",
             font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────
print(f"\nSaving to {OUTPUT}...")
prs.save(OUTPUT)
print(f"✅ Done! Output: {OUTPUT}")

# ── VALIDATION ────────────────────────────────────────────────────────────
print("\n── Validation ──")
print(f"TAS Total Mar MTD: Q={total_mar[0]}, M={total_mar[1]}, A={total_mar[2]}, P={total_mar[3]}")
print(f"  Q→M={fmtR(total_qm)}, M→A={fmtR(rate(total_mar[2],total_mar[1]))}, A→P={fmtR(total_ap)}, E2E={fmtR(rate(total_mar[3],total_mar[0]))}")
print(f"TAS Only Mar MTD: Q={t[0]}, M={t[1]}, A={t[2]}, P={t[3]}")
print(f"Aliado Mar MTD: Q={al[0]}, M={al[1]}, A={al[2]}, P={al[3]}")
print(f"Retail Mar: Made={retail['Mar']['made']}, App={retail['Mar']['approved']}, Pur={retail['Mar']['purchases']}")
print(f"RR: {daily_rate_p:.1f}/day, Mar proj={mar_projected}, Apr baseline={apr_baseline}")
